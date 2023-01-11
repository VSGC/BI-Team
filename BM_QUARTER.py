# -*- coding: utf-8 -*-
"""
Created on Tue Sep 13 16:01:58 2022

@author: VAIBHAV.SRIVASTAV01
"""






from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import pandas as pd
import os
from datetime import date
import calendar
import numpy as np
from pptx import Presentation

import pyodbc
pyodbc.drivers()
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
config = {}
with open(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\config.txt","r") as file:
    for line in file:
        key, value = line.strip().split(",")
        config[key] = value

# del key,line,value,file
# mydata=pd.read_csv(r"C:\Users\VAIBHAV.SRIVASTAV01\.spyder-py3\DTCRON_MASTER_03072022.csv")
# mydata=pd.read_excel(r"C:\Users\vaibhav.srivastav01\Downloads\DTCRON (18 AUG).xlsx")
# #fd=pd.read_csv(r"C:\Users\VAIBHAV.SRIVASTAV01\.spyder-py3\FINDISBURSEMENTDETAILS (30th June).csv")
# fd=pd.read_csv(r"C:\Users\vaibhav.srivastav01\Desktop\findisb1808.csv")


del key,line,value,file

conx = pyodbc.connect('DRIVER={SnowflakeDSIIDriver};SERVER='+str(config['server'])+';Warehouse = COMPUTE_WH;DATABASE='+str(config['db1'])+';UID='+ str(config['user1']) +';PWD='+ str(config['password']),autocommit=True)
cur = conx.cursor()


cur.execute("USE DATABASE SAMPLE_DB")

query = "select * from dtcron_master_05012023;"
q2="select * from V_FINDISBURSEMENTDETAILS;"
mydata = pd.read_sql(query, conx)
fd=pd.read_sql(q2,conx)
df=mydata.copy()
df1= pd.read_sql("select * from dtcron_master_04122022;", conx)
df2= pd.read_sql("select * from dtcron_master_05112022;", conx)
# m1target=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Ops Review_Q2 targets.xlsx", sheet_name="Jul_Internal")
# m2target=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Ops Review_Q2 targets.xlsx", sheet_name="Aug_Internal")
# m3target=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Ops Review_Q2 targets.xlsx", sheet_name="Sep_Internal")




# def change(df):
#     if ('AMD' in df.index)|('NCR' in df.index)|('Total ' in df.index)|('Mumbai ' in df.index)|('Bangalore ' in df.index):
#         a=df.index.get_indexer(['AMD','NCR','Total ','Mumbai ','Bangalore '])
#         b=df.index.to_list()
#         b[a[0]]='Ahmedabad'
#         b[a[1]]='Delhi'
#         b[a[2]]='Total'
#         b[a[3]]='Mumbai'
#         b[a[4]]='Bangalore'
#         df.index=b
#     return df
# def month_target(target):
#     target_df=target.copy()
#     target_df=target_df[['Unnamed: 17','Unnamed: 18','Unnamed: 19','Unnamed: 20']]
#     length=int(len(target_df)/2)
#     mtd_target_df=target_df.iloc[2:length-1]
#     mtd_target_df.reset_index(drop=True,inplace=True)
#     ytd_target_df=target_df.iloc[length+3:]
#     ytd_target_df.reset_index(drop=True,inplace=True)
#     ytd_target_df.loc[0,'Unnamed: 20']='AUM'
#     mtd_target_df.loc[0,'Unnamed: 20']='AUM'
#     ytd_target_df.loc[0,'Unnamed: 17']='BRANCH'
#     mtd_target_df.loc[0,'Unnamed: 17']='BRANCH'
#     ytd_target_df.columns = ytd_target_df.iloc[0]
#     mtd_target_df.columns = mtd_target_df.iloc[0]
#     ytd_target_df.drop(0,inplace=True)
#     mtd_target_df.drop(0,inplace=True)
    
#     mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
#     mtd_target_df.set_index('BRANCH',inplace=True)
    
#     target_df=target.copy()
#     target_df=target_df[['Unnamed: 2','Unnamed: 3','Unnamed: 4','Unnamed: 5']]
#     length=int(len(target_df)/2)
#     GPL_mtd_target_df=target_df.iloc[2:length-1]
#     GPL_mtd_target_df.reset_index(drop=True,inplace=True)
#     GPL_ytd_target_df=target_df.iloc[length+3:]
#     GPL_ytd_target_df.reset_index(drop=True,inplace=True)
#     GPL_ytd_target_df.loc[0,'Unnamed: 5']='AUM'
#     GPL_mtd_target_df.loc[0,'Unnamed: 5']='AUM'
#     GPL_ytd_target_df.loc[0,'Unnamed: 2']='BRANCH'
#     GPL_mtd_target_df.loc[0,'Unnamed: 2']='BRANCH'
#     GPL_ytd_target_df.columns = GPL_ytd_target_df.iloc[0]
#     GPL_mtd_target_df.columns = GPL_mtd_target_df.iloc[0]
#     GPL_ytd_target_df.drop(0,inplace=True)
#     GPL_mtd_target_df.drop(0,inplace=True)
#     GPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
#     GPL_mtd_target_df.loc[6,'BRANCH']=0
#     GPL_mtd_target_df.set_index('BRANCH',inplace=True)
    
#     target_df=target.copy()
#     target_df=target_df[['Unnamed: 7','Unnamed: 8','Unnamed: 9','Unnamed: 10']]
#     length=int(len(target_df)/2)
#     NGPL_mtd_target_df=target_df.iloc[2:length-1]
#     NGPL_mtd_target_df.reset_index(drop=True,inplace=True)
#     NGPL_ytd_target_df=target_df.iloc[length+3:]
#     NGPL_ytd_target_df.reset_index(drop=True,inplace=True)
#     NGPL_ytd_target_df.loc[0,'Unnamed: 10']='AUM'
#     NGPL_mtd_target_df.loc[0,'Unnamed: 10']='AUM'
#     NGPL_ytd_target_df.loc[0,'Unnamed: 7']='BRANCH'
#     NGPL_mtd_target_df.loc[0,'Unnamed: 7']='BRANCH'
#     NGPL_ytd_target_df.columns = NGPL_ytd_target_df.iloc[0]
#     NGPL_mtd_target_df.columns = NGPL_mtd_target_df.iloc[0]
#     NGPL_ytd_target_df.drop(0,inplace=True)
#     NGPL_mtd_target_df.drop(0,inplace=True)
#     NGPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
#     NGPL_mtd_target_df.set_index('BRANCH',inplace=True)
    
#     target_df=target.copy()
#     target_df=target_df[['Unnamed: 22','Unnamed: 23','Unnamed: 24','Unnamed: 25']]
#     length=int(len(target_df)/2)
#     HL_mtd_target_df=target_df.iloc[2:length-1]
#     HL_mtd_target_df.reset_index(drop=True,inplace=True)
#     HL_ytd_target_df=target_df.iloc[length+3:]
#     HL_ytd_target_df.reset_index(drop=True,inplace=True)
#     HL_ytd_target_df.loc[0,'Unnamed: 25']='AUM'
#     HL_mtd_target_df.loc[0,'Unnamed: 25']='AUM'
#     HL_ytd_target_df.loc[0,'Unnamed: 22']='BRANCH'
#     HL_mtd_target_df.loc[0,'Unnamed: 22']='BRANCH'
#     HL_ytd_target_df.columns = HL_ytd_target_df.iloc[0]
#     HL_mtd_target_df.columns = HL_mtd_target_df.iloc[0]
#     HL_ytd_target_df.drop(0,inplace=True)
#     HL_mtd_target_df.drop(0,inplace=True)
#     HL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
#     HL_mtd_target_df.set_index('BRANCH',inplace=True)
    
#     target_df=target.copy()
#     target_df=target_df[['Unnamed: 12','Unnamed: 13','Unnamed: 14','Unnamed: 15']]
#     length=int(len(target_df)/2)
#     LAP_mtd_target_df=target_df.iloc[2:length-1]
#     LAP_mtd_target_df.reset_index(drop=True,inplace=True)
#     LAP_ytd_target_df=target_df.iloc[length+3:]
#     LAP_ytd_target_df.reset_index(drop=True,inplace=True)
#     LAP_ytd_target_df.loc[0,'Unnamed: 15']='AUM'
#     LAP_mtd_target_df.loc[0,'Unnamed: 15']='AUM'
#     LAP_ytd_target_df.loc[0,'Unnamed: 12']='BRANCH'
#     LAP_mtd_target_df.loc[0,'Unnamed: 12']='BRANCH'
#     LAP_ytd_target_df.columns = LAP_ytd_target_df.iloc[0]
#     LAP_mtd_target_df.columns = LAP_mtd_target_df.iloc[0]
#     LAP_ytd_target_df.drop(0,inplace=True)
#     LAP_mtd_target_df.drop(0,inplace=True)
#     LAP_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
#     LAP_mtd_target_df.set_index('BRANCH',inplace=True)
    
#     month_target={}
#     LAP_mtd_target_df=change(LAP_mtd_target_df)
#     NGPL_mtd_target_df=change(NGPL_mtd_target_df)
#     GPL_mtd_target_df=change(GPL_mtd_target_df)
#     mtd_target_df=change(mtd_target_df)
#     HL_mtd_target_df=change(HL_mtd_target_df)
#     LAP_mtd_target_df=LAP_mtd_target_df.round(decimals=0)
#     mtd_target_df=mtd_target_df.round(decimals=0)
#     GPL_mtd_target_df=GPL_mtd_target_df.round(decimals=0)
#     NGPL_mtd_target_df=NGPL_mtd_target_df.round(decimals=0)
#     HL_mtd_target_df=HL_mtd_target_df.round(decimals=0)
#     month_target['HL_mtd_target_df']=HL_mtd_target_df
#     month_target['mtd_target_df']=mtd_target_df
#     month_target['GPL_mtd_target_df']=GPL_mtd_target_df
#     month_target['NGPL_mtd_target_df']=NGPL_mtd_target_df
#     month_target['LAP_mtd_target_df']=LAP_mtd_target_df
#     return (month_target)

# m1trgt=month_target(m1target)
# m2trgt=month_target(m2target)
# m3trgt=month_target(m3target)

target3=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\December Targets.xlsx")
target2=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\For vaibhav.xlsx",sheet_name="November targets")
target1=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\For vaibhav.xlsx",sheet_name="October targets")
def change(df):
    if ('AMD' in df.index)|('NCR' in df.index)|('Total ' in df.index)|('Mumbai ' in df.index)|('Bangalore ' in df.index):
        a=df.index.get_indexer(['AMD','NCR','Total','Mumbai','Bangalore'])
        b=df.index.to_list()
        b[a[0]]='Ahmedabad'
        b[a[1]]='Delhi'
        b[a[2]]='Total'
        b[a[3]]='Mumbai'
        b[a[4]]='Bangalore'
        df.index=b
    return df

def get_target(target):
    target_df=target.copy()
    target_df=target_df[['Unnamed: 20','Unnamed: 21','Unnamed: 22','Unnamed: 23']]
    length=int(len(target_df)/2)
    mtd_target_df=target_df.iloc[4:length]
    mtd_target_df.reset_index(drop=True,inplace=True)
    ytd_target_df=target_df.iloc[length+3:]
    ytd_target_df.reset_index(drop=True,inplace=True)
    ytd_target_df.loc[0,'Unnamed: 23']='AUM'
    mtd_target_df.loc[0,'Unnamed: 23']='AUM'
    ytd_target_df.loc[0,'Unnamed: 20']='BRANCH'
    mtd_target_df.loc[0,'Unnamed: 20']='BRANCH'
    ytd_target_df.columns = ytd_target_df.iloc[0]
    mtd_target_df.columns = mtd_target_df.iloc[0]
    ytd_target_df.drop(0,inplace=True)
    mtd_target_df.drop(0,inplace=True)
    
    mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    mtd_target_df.set_index('BRANCH',inplace=True)
    
    target_df=target.copy()
    target_df=target_df[['Unnamed: 3','Unnamed: 4','Unnamed: 5','Unnamed: 6']]
    length=int(len(target_df)/2)
    GPL_mtd_target_df=target_df.iloc[4:length]
    GPL_mtd_target_df.reset_index(drop=True,inplace=True)
    GPL_ytd_target_df=target_df.iloc[length+3:]
    GPL_ytd_target_df.reset_index(drop=True,inplace=True)
    GPL_ytd_target_df.loc[0,'Unnamed: 6']='AUM'
    GPL_mtd_target_df.loc[0,'Unnamed: 6']='AUM'
    GPL_ytd_target_df.loc[0,'Unnamed: 3']='BRANCH'
    GPL_mtd_target_df.loc[0,'Unnamed: 3']='BRANCH'
    GPL_ytd_target_df.columns = GPL_ytd_target_df.iloc[0]
    GPL_mtd_target_df.columns = GPL_mtd_target_df.iloc[0]
    GPL_ytd_target_df.drop(0,inplace=True)
    GPL_mtd_target_df.drop(0,inplace=True)
    GPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    GPL_mtd_target_df.loc[6,'BRANCH']=0
    GPL_mtd_target_df.set_index('BRANCH',inplace=True)
    
    target_df=target.copy()
    target_df=target_df[['Unnamed: 9','Unnamed: 10','Unnamed: 11','Unnamed: 12']]
    length=int(len(target_df)/2)
    NGPL_mtd_target_df=target_df.iloc[4:length]
    NGPL_mtd_target_df.reset_index(drop=True,inplace=True)
    NGPL_ytd_target_df=target_df.iloc[length+3:]
    NGPL_ytd_target_df.reset_index(drop=True,inplace=True)
    NGPL_ytd_target_df.loc[0,'Unnamed: 12']='AUM'
    NGPL_mtd_target_df.loc[0,'Unnamed: 12']='AUM'
    NGPL_ytd_target_df.loc[0,'Unnamed: 9']='BRANCH'
    NGPL_mtd_target_df.loc[0,'Unnamed: 9']='BRANCH'
    NGPL_ytd_target_df.columns = NGPL_ytd_target_df.iloc[0]
    NGPL_mtd_target_df.columns = NGPL_mtd_target_df.iloc[0]
    NGPL_ytd_target_df.drop(0,inplace=True)
    NGPL_mtd_target_df.drop(0,inplace=True)
    NGPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    NGPL_mtd_target_df.set_index('BRANCH',inplace=True)
    
    target_df=target.copy()
    target_df=target_df[['Unnamed: 15','Unnamed: 16','Unnamed: 17','Unnamed: 18']]
    length=int(len(target_df)/2)
    LAP_mtd_target_df=target_df.iloc[4:length]
    LAP_mtd_target_df.reset_index(drop=True,inplace=True)
    LAP_ytd_target_df=target_df.iloc[length+3:]
    LAP_ytd_target_df.reset_index(drop=True,inplace=True)
    LAP_ytd_target_df.loc[0,'Unnamed: 18']='AUM'
    LAP_mtd_target_df.loc[0,'Unnamed: 18']='AUM'
    LAP_ytd_target_df.loc[0,'Unnamed: 15']='BRANCH'
    LAP_mtd_target_df.loc[0,'Unnamed: 15']='BRANCH'
    LAP_ytd_target_df.columns = LAP_ytd_target_df.iloc[0]
    LAP_mtd_target_df.columns = LAP_mtd_target_df.iloc[0]
    LAP_ytd_target_df.drop(0,inplace=True)
    LAP_mtd_target_df.drop(0,inplace=True)
    LAP_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    LAP_mtd_target_df.set_index('BRANCH',inplace=True)
    
    
    target_df=target.copy()
    target_df=target_df[['Unnamed: 25','Unnamed: 26','Unnamed: 27','Unnamed: 28']]
    length=int(len(target_df)/2)
    HL_mtd_target_df=target_df.iloc[4:length]
    HL_mtd_target_df.reset_index(drop=True,inplace=True)
    HL_ytd_target_df=target_df.iloc[length+3:]
    HL_ytd_target_df.reset_index(drop=True,inplace=True)
    HL_ytd_target_df.loc[0,'Unnamed: 28']='AUM'
    HL_mtd_target_df.loc[0,'Unnamed: 28']='AUM'
    HL_ytd_target_df.loc[0,'Unnamed: 25']='BRANCH'
    HL_mtd_target_df.loc[0,'Unnamed: 25']='BRANCH'
    HL_ytd_target_df.columns = HL_ytd_target_df.iloc[0]
    HL_mtd_target_df.columns = HL_mtd_target_df.iloc[0]
    HL_ytd_target_df.drop(0,inplace=True)
    HL_mtd_target_df.drop(0,inplace=True)
    HL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    HL_mtd_target_df.set_index('BRANCH',inplace=True)
    
    HL_mtd_target_df=change(HL_mtd_target_df)
    LAP_mtd_target_df=change(LAP_mtd_target_df)
    NGPL_mtd_target_df=change(NGPL_mtd_target_df)
    GPL_mtd_target_df=change(GPL_mtd_target_df)
    mtd_target_df=change(mtd_target_df)
    HL_mtd_target_df=HL_mtd_target_df.round(decimals=0)
    LAP_mtd_target_df=LAP_mtd_target_df.round(decimals=0)
    mtd_target_df=mtd_target_df.round(decimals=0)
    GPL_mtd_target_df=GPL_mtd_target_df.round(decimals=0)
    NGPL_mtd_target_df=NGPL_mtd_target_df.round(decimals=0)
    LAP_mtd_target_df=LAP_mtd_target_df.transpose()
    LAP_mtd_target_df=LAP_mtd_target_df[['Mumbai', 'Ahmedabad', 'Surat','Pune', 'Indore',  'Delhi', 'Chandigarh', 'Jaipur','Bangalore', 'Chennai','Hyderabad', 'South', 'Bangalore','Total']]
    LAP_mtd_target_df=LAP_mtd_target_df.transpose()
    HL_mtd_target_df=HL_mtd_target_df.transpose()
    HL_mtd_target_df=HL_mtd_target_df[['Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune','Total']]
    HL_mtd_target_df=HL_mtd_target_df.transpose()
    a={}
    a['HL_mtd_target_df']=HL_mtd_target_df
    a['LAP_mtd_target_df']=LAP_mtd_target_df
    a['mtd_target_df']=mtd_target_df
    a['GPL_mtd_target_df']=GPL_mtd_target_df
    a['NGPL_mtd_target_df']=NGPL_mtd_target_df
    return a
m1trgt=get_target(target1)
m2trgt=get_target(target2)
m3trgt=get_target(target3)
'''
#--------------LAP_BUSINESS METRICS-------------------

#------------BOOK-----------

LP_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')

#------------ROI------------
LP_ROI<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
               & BM$BOOK_YEAR_MONTH == MTD)
#------------PF-------------
LP_PF<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
                     & BM$BOOK_YEAR_MONTH == MTD)

#------------LTD (BOOKING TO LOGIN)----------
LTD_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')
LTD_LOGIN<-subset (BM,BM$FINTYPE %in% c('LP','NP') & BM$LOGINSTATUS == 'A) Login')

LTD_BOOK1<-ddply(LTD_BOOK,c('FINBRANCH'),summarise,
              Book_value=sum(BOOKED_AMOUNT/10000000), 
              Book_vol=length(FINREFERENCE))
LTD_LOGIN1<-ddply(LTD_LOGIN,c('FINBRANCH'),summarise,
                  Log_value=sum(REQLOANAMT/1000000000),
                  Log_Vol=length(FINREFERENCE))

LTD_BOOK_LOG_mrg<-merge(LTD_LOGIN1,LTD_BOOK1,by.x = c('FINBRANCH'),
                        by.y = c('FINBRANCH'),all.x = T)
'''







def buss_metrics(LOAN,df,fd,todays_date):
    
    MTD=str(todays_date.year)+"-"+str(todays_date.month)
    
    YTD=[]
    cur_month=todays_date.month
    while cur_month>=4:
        YTD.append(cur_month)
        cur_month=cur_month-1
    ytd=[]
    for i in YTD:
        ytd.append(str(todays_date.year)+"-"+str(i))
    
    
    BM_df = df[[ 'FINREFERENCE','NET_PREMIUM','BOOK_YEAR_MONTH','LOGIN_YEAR_MONTH','FINBRANCH','FINTYPE','STATUS', 'WROI','BOOKED_AMOUNT' ,'PROCESSING_FEE','REQLOANAMT','LOGINSTATUS','GPLFLAG_SANCTIONS','PRINCIPAL_OUTSTANDING','DISB_AMOUNT']]
    findisb=fd[['FINREFERENCE','DISBSEQ','DISBAMOUNT']]
    findisb=findisb[(findisb['DISBSEQ']==1)]
    findisb=findisb.sort_values(by='FINREFERENCE')
    BM_df=BM_df.sort_values(by='FINREFERENCE')
    fd_bm= pd.merge(BM_df, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')

    if LOAN== 'LAP':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'LP') | (fd_bm['FINTYPE'] =='NP'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'LP') |  (bm_df['FINTYPE'] =='NP')) ]
    elif LOAN== 'HL':
        fd_bm= fd_bm[ (fd_bm['FINBRANCH'].isin(['Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']))&(fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'HL') | (fd_bm['FINTYPE'] =='HT') | (fd_bm['FINTYPE'] =='FL') | (fd_bm['FINTYPE'] =='FT') | (fd_bm['FINTYPE'] =='LT')| (fd_bm['FINTYPE'] =='AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[(bm_df['FINBRANCH'].isin(['Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']))& ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'HL') | (bm_df['FINTYPE'] =='HT') | (bm_df['FINTYPE'] =='FL') | (bm_df['FINTYPE'] =='FT') | (bm_df['FINTYPE'] =='LT')| (bm_df['FINTYPE'] =='AHL')) ]
    elif LOAN== 'GPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'] == 'GPL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'] == 'GPL') ]
    elif LOAN== 'NGPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL'])))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL'])) ]
    elif LOAN== 'AHL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['FINTYPE'] == 'AHL') ]
    
    # bm_df_hl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &((bm_df['FINTYPE'] == 'HL') |  (bm_df['FINTYPE'] =='HT'))]
    # bm_df_gpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='GPL')]
    # bm_df_ngpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='NON GPL')]
    bm_df['one_count']=np.ones(len(bm_df), dtype = int)
    
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','FINREFERENCE','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    # def buss_metrics(bm_df):
        #bm_df_group = bm_df.groupby(['FINBRANCH'],as_index=False).sum(['WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING' ])
        #bm_df_group['BOOK_TARGET'] =
        
    fd_bm_group=fd_bm.groupby(['FINBRANCH'],as_index=False).sum(['WROI', 'BOOKED_AMOUNT', 'PROCESSING_FEE','REQLOANAMT','PRINCIPAL_OUTSTANDING', 'DISB_AMOUNT', 'DISBSEQ', 'DISBAMOUNT'])
    fd_bm_group=fd_bm_group[['FINBRANCH','DISBAMOUNT']]
    bm_df_book=bm_df[(bm_df.STATUS== 'Booked')&(bm_df.BOOK_YEAR_MONTH==MTD)]
    # bm_df_book=pd.merge(bm_df_book, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')
    bm_df_book['one_count']=np.ones(len(bm_df_book), dtype = int)
    bm_df_BOOK_group=bm_df_book.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    
    ltd=[]
    z = bm_df.BOOK_YEAR_MONTH.unique().tolist()
    z=[x for x in z if x is not None]
    z.sort()
    for p in z:
        if p<=MTD:
            ltd.append(p)
            
    b=bm_df[(bm_df.STATUS== 'Booked')& (bm_df.BOOK_YEAR_MONTH.isin(ltd))]
    l=bm_df[(bm_df.LOGINSTATUS== 'A) Login')& (bm_df.LOGIN_YEAR_MONTH.isin(ltd))]
    c=bm_df[((bm_df.STATUS== 'Booked') & (bm_df.BOOK_YEAR_MONTH.isin(ytd)))]
    c_group=c.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','WROI'])
    c_group=c_group[['FINBRANCH','BOOKED_AMOUNT','WROI']]
    c_group.rename(columns = {'BOOKED_AMOUNT':'BOOKED_AMOUNT_M1'}, inplace = True)
    c_group.rename(columns = {'WROI':'WROI_M1'}, inplace = True)
    b_group=b.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','one_count' ])
    b_group=b_group[['FINBRANCH','PRINCIPAL_OUTSTANDING','one_count','BOOKED_AMOUNT' ]]
    b_group.rename(columns = {'one_count':'one_count_tot_book'}, inplace = True)
    b_group.rename(columns = {'BOOKED_AMOUNT':'TOTAL_BOOKED_AMOUNTytd'}, inplace = True)
    
    
    l_group=l.groupby(['FINBRANCH'],as_index=False).sum(['REQLOANAMT','one_count' ])
    l_group=l_group[['FINBRANCH','one_count','REQLOANAMT' ]]
    l_group.rename(columns = {'one_count':'one_count_tot_login'}, inplace = True)
    
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','BOOKED_AMOUNT','WROI','NET_PREMIUM','one_count','PROCESSING_FEE']]
    bm_df_BOOK_group.rename(columns = {'one_count':'one_count_MTD'}, inplace = True)
    bm_df_group=bm_df_group[['FINBRANCH','one_count']]
    
    merged=pd.merge(bm_df_group, bm_df_BOOK_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged2=pd.merge(merged, c_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged3=pd.merge(merged2, b_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged4=pd.merge(merged3, fd_bm_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged4=pd.merge(merged4, l_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    if LOAN=='LAP' and (int(MTD[5])>6) :
        print('yes')
        newb=merged4.copy()
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.set_index('FINBRANCH',inplace=True)
        newb=newb.transpose()
        newb=newb[['Hyderabad','Indore', 'Jaipur', 'Chandigarh', 'Chennai', 'Surat']]
        newb=newb.transpose()
        newb.reset_index(inplace=True)
        newb.loc[11]=newb.sum()
        newb.loc[11,'FINBRANCH']='New Branches'
        newb['BOOK_VALUE'] = newb.BOOKED_AMOUNT/10000000
        #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]
    
        newb['BOOK_VOL'] = newb.one_count_MTD
        #newb['DISB AMNT Tranch 1']=newb.DISBAMOUNT/1000000000
        newb['DISB_Act (%)']=(newb.DISBAMOUNT)/newb.BOOKED_AMOUNT
        newb['AUM']=newb.PRINCIPAL_OUTSTANDING/10000000
        newb['ROI MTD'] = newb.WROI / newb.BOOKED_AMOUNT
        newb['ROI M-1'] = newb.WROI_M1 /newb.BOOKED_AMOUNT_M1
        newb['GROSS%']=100*newb.NET_PREMIUM/newb.BOOKED_AMOUNT
        newb['PF%']=100*newb.PROCESSING_FEE / newb.BOOKED_AMOUNT
        newb['LTD_volume']=100*newb.one_count_tot_book/newb.one_count_tot_login
        newb['LTD_Value']=100*(newb.TOTAL_BOOKED_AMOUNTytd*100)/newb.REQLOANAMT
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    
        newb.loc['Total']= newb.sum()
        if sum(newb.BOOKED_AMOUNT)!=0:
            newb.loc['Total','DISB_Act (%)']=sum(newb.DISBAMOUNT)/sum(newb.BOOKED_AMOUNT)
            newb.loc['Total','ROI MTD']=sum(newb.WROI ) /sum( newb.BOOKED_AMOUNT)
            newb.loc['Total','ROI M-1']=sum(newb.WROI_M1)/sum(newb.BOOKED_AMOUNT_M1)
            newb.loc['Total','GROSS%']=100*sum(newb.NET_PREMIUM)/sum(newb.BOOKED_AMOUNT)
            newb.loc['Total','PF%']=100*sum(newb.PROCESSING_FEE ) / sum(newb.BOOKED_AMOUNT)
        else:
            newb.loc['Total','DISB_Act (%)']=0
            newb.loc['Total','ROI MTD']=0
            newb.loc['Total','ROI M-1']=0
            newb.loc['Total','GROSS%']=0
            newb.loc['Total','PF%']=0
        newb.loc['Total','LTD_volume']=100*sum(newb.one_count_tot_book)/sum(newb.one_count_tot_login)
        newb.loc['Total','LTD_Value']=100*sum(newb.TOTAL_BOOKED_AMOUNTytd)*(100)/sum(newb.REQLOANAMT)
        newb.loc['Total','FINBRANCH']='Total'
        newb.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
        newb["new"] = range(1,len(newb)+1)
        newb.loc['Total','new'] = 0
        newb=newb.sort_values(by="new").drop('new', axis=1)
        newb=newb.round(decimals=2)
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
        newb.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
        newb.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
        newb.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
        newb.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
        newb.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
        newb.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
        newb.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
        newb['Booking Tar(Cr)']=' '
        newb['Value %']=' '
        newb['Booking Tar(#)']=' '
        newb['Volume %']=' '
        newb['AUM Tar(Cr)']=' '
        newb['AUM Actual %']=' '
        newb['Net']=' '
        newb['Empaneled']=' '
        newb['Active']=' '
        newb['Channel']=' '
        newb['Employee']=' '
        newb['Cumulative']=' '
        newb.columns
        newb=newb.round(decimals=2)
        newb['Value%']=round(newb['Value%'])
        newb['Volume%']=round(newb['Volume%'])
        newb=newb[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Value %','DISB_Act (%)','Booking Tar(#)', 'Booking Act(#)','Volume %', 'AUM Tar(Cr)','AUM Act(Cr)', 'AUM Actual %','ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
        #newb.set_index('FINBRANCH',inplace=True)
        
    
        finaln=newb.transpose()
        
        finaln.insert(0,'Target',0)
        finaln.loc['FINBRANCH','Target']='Target'
    
    
    merged4['BOOK_VALUE'] = merged4.BOOKED_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]

    merged4['BOOK_VOL'] = merged4.one_count_MTD
    #merged4['DISB AMNT Tranch 1']=merged4.DISBAMOUNT/1000000000
    merged4['DISB_Act (%)']=(merged4.DISBAMOUNT)/merged4.BOOKED_AMOUNT
    merged4['AUM']=merged4.PRINCIPAL_OUTSTANDING/10000000
    merged4['ROI MTD'] = merged4.WROI / merged4.BOOKED_AMOUNT
    merged4['ROI M-1'] = merged4.WROI_M1 /merged4.BOOKED_AMOUNT_M1
    merged4['GROSS%']=100*merged4.NET_PREMIUM/merged4.BOOKED_AMOUNT
    merged4['PF%']=100*merged4.PROCESSING_FEE / merged4.BOOKED_AMOUNT
    merged4['LTD_volume']=100*merged4.one_count_tot_book/merged4.one_count_tot_login
    merged4['LTD_Value']=100*(merged4.TOTAL_BOOKED_AMOUNTytd*100)/merged4.REQLOANAMT
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    merged4.loc['Total']= merged4.sum()
    merged4.loc['Total','DISB_Act (%)']=sum(merged4.DISBAMOUNT)/sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','ROI MTD']=sum(merged4.WROI ) /sum( merged4.BOOKED_AMOUNT)
    merged4.loc['Total','ROI M-1']=sum(merged4.WROI_M1)/sum(merged4.BOOKED_AMOUNT_M1)
    merged4.loc['Total','GROSS%']=100*sum(merged4.NET_PREMIUM)/sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','PF%']=100*sum(merged4.PROCESSING_FEE ) / sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','LTD_volume']=100*sum(merged4.one_count_tot_book)/sum(merged4.one_count_tot_login)
    merged4.loc['Total','LTD_Value']=100*sum(merged4.TOTAL_BOOKED_AMOUNTytd)*(100)/sum(merged4.REQLOANAMT)
    merged4.loc['Total','FINBRANCH']='Total'
    merged4.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    merged4["new"] = range(1,len(merged4)+1)
    merged4.loc['Total','new'] = 0
    merged4=merged4.sort_values(by="new").drop('new', axis=1)
    merged4=merged4.round(decimals=2)
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    merged4.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    merged4.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    merged4.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    merged4.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    merged4.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    merged4.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    merged4['Booking Tar(Cr)']=' '
    merged4['Value %']=' '
    merged4['Booking Tar(#)']=' '
    merged4['Volume %']=' '
    merged4['AUM Tar(Cr)']=' '
    merged4['AUM Actual %']=' '
    merged4['Net']=' '
    merged4['Empaneled']=' '
    merged4['Active']=' '
    merged4['Channel']=' '
    merged4['Employee']=' '
    merged4['Cumulative']=' '
    merged4.columns
    merged4=merged4.round(decimals=2)
    merged4['Value%']=round(merged4['Value%'])
    merged4['Volume%']=round(merged4['Volume%'])
    merged4=merged4[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Value %','DISB_Act (%)','Booking Tar(#)', 'Booking Act(#)','Volume %', 'AUM Tar(Cr)','AUM Act(Cr)', 'AUM Actual %','ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #merged4.set_index('FINBRANCH',inplace=True)
    

    final=merged4.transpose()
    
    final.insert(0,'Target',0)
    final.loc['FINBRANCH','Target']='Target'
    if LOAN=='LAP' and (int(MTD[5])>6) :
        final.reset_index(inplace=True)
        finaln.reset_index(inplace=True)
        finaln=finaln[['index',11]]
        final=pd.merge(final,finaln,left_on=('index'),right_on=('index'),how='outer')
        final.set_index('index',inplace=True)
    return(final)
# return(final)
lap_metrics=buss_metrics('LAP',df,fd,date(2022,8,5))
# lap_metrics=lap_metrics[['Target', 'Total', 1, 4, 8, 9, 0, 2, 3, 5, 6, 7, 10,11]]
# lap_metrics.columns=lap_metrics.loc['FINBRANCH']
def set_target(lap_metrics,LAP_mtd_target_df,month):
    
    for i in LAP_mtd_target_df.index:
        if  (i != 0): 
            lap_metrics.loc['Booking Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'Value (Cr)']
            print((lap_metrics))
            lap_metrics.loc['Booking Act(Cr)',i]=round(float(lap_metrics.loc['Booking Act(Cr)',i]))
            lap_metrics.loc['Value %',i]=str(round(100*lap_metrics.loc['Booking Act(Cr)',i]/LAP_mtd_target_df.loc[i,'Value (Cr)']))+'% '
            lap_metrics.loc['Booking Tar(#)',i]=LAP_mtd_target_df.loc[i,'#']
            lap_metrics.loc['Booking Act(#)',i]=round(lap_metrics.loc['Booking Act(#)',i])
            lap_metrics.loc['Volume %',i]=str(round(100*lap_metrics.loc['Booking Act(#)',i]/LAP_mtd_target_df.loc[i,'#']))+'% '
            lap_metrics.loc['AUM Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'AUM']
            lap_metrics.loc['AUM Act(Cr)',i]=round(float(lap_metrics.loc['AUM Act(Cr)',i]))
            lap_metrics.loc['AUM Actual %',i]=str(round(100*lap_metrics.loc['AUM Act(Cr)',i]/LAP_mtd_target_df.loc[i,'AUM']))
    return lap_metrics

# hl_metrics=buss_metrics('HL',df,fd,date(2022, 10, 10))

# hl_metrics=hl_metrics[['Target', 'Total', 1, 2, 3, 4,0]]
# hl_metrics.columns=hl_metrics.loc['FINBRANCH']
# gpl_metrics=buss_metrics('GPL',df,fd,todays_date )
# gpl_metrics=gpl_metrics[[ 'Target', 'Total', 1, 2, 3, 4,0]]
# gpl_metrics.columns=gpl_metrics.loc['FINBRANCH']
# ngpl_metrics=buss_metrics('NGPL',df,fd,todays_date )
# ngpl_metrics=ngpl_metrics[['Target', 'Total', 1, 2, 3, 4,0]]
# ngpl_metrics.columns=ngpl_metrics.loc['FINBRANCH']
# ahl_metrics=buss_metrics('AHL',df,fd,todays_date )
# ahl_metrics.columns=ahl_metrics.loc['FINBRANCH']



# lap_metrics=set_target(lap_metrics,m2trgt['LAP_mtd_target_df'])
# gpl_metrics=set_target(gpl_metrics,m2trgt['GPL_mtd_target_df'])
# ngpl_metrics=set_target(ngpl_metrics,m2trgt['NGPL_mtd_target_df'])


# lap_metrics.loc['DISB_Act (%)','Target']='90%'
# lap_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
# lap_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
# lap_metrics.loc['GROSS%','Target']='1.5%'
# lap_metrics.loc['PF%','Target']='0.9%'
# lap_metrics.loc['Cumulative','Target']='1.35%'
# lap_metrics.loc['Value%','Target']='33%'
# lap_metrics.loc['Volume%','Target']='33%'

# # hl_metrics.loc['DISB_Act (%)','Target']='90%'
# hl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
# hl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
# hl_metrics.loc['GROSS%','Target']='2%'
# hl_metrics.loc['PF%','Target']='0.02%'
# hl_metrics.loc['Cumulative','Target']='0.45%'
# hl_metrics.loc['Value%','Target']='50%'
# hl_metrics.loc['Volume%','Target']='50%'

# gpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
# gpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
# gpl_metrics.loc['GROSS%','Target']='2%'
# gpl_metrics.loc['PF%','Target']='0.02%'
# gpl_metrics.loc['Cumulative','Target']='0.45%'
# gpl_metrics.loc['Value%','Target']='50%'
# gpl_metrics.loc['Volume%','Target']='50%'


# ngpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
# ngpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
# ngpl_metrics.loc['GROSS%','Target']='2%'
# ngpl_metrics.loc['PF%','Target']='0.02%'
# ngpl_metrics.loc['Cumulative','Target']='0.45%'
# ngpl_metrics.loc['Value%','Target']='50%'
# ngpl_metrics.loc['Volume%','Target']='50%'


# MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
# with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\AUG\1"+MTD+".xlsx") as writer:
#     lap_metrics.to_excel(writer, sheet_name="LAP BM")
#     hl_metrics.to_excel(writer, sheet_name="HL BM")
#     gpl_metrics.to_excel(writer, sheet_name="GPL BM")
#     ngpl_metrics.to_excel(writer, sheet_name="NGPL BM")
#     ahl_metrics.to_excel(writer, sheet_name="AHL BM")



# ---create presentation with 1 slide---
# prs = Presentation(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\OPSREVIEWDECK.pptx")
# # prs.slide_width = 11887200
# # prs.slide_height = 6686550
def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
# def create_slide(prs,lap_metrics,title,c):
#     prs.slides[c].shapes[0].text=title
#     lap_metrics.reset_index(level=0, inplace=True)
#     lap_metrics.insert(0,' ',' ')
#     # slide = prs.slides.add_slide(prs.slide_layouts[5])
#     # slide.shapes.title.text = title
#     # slide.shapes.title.top = Inches(0.5)
#     # slide.shapes.title.left = Inches(0.5)
#     # slide.shapes.title.width = Inches(10)

    
#     # # ---add table to slide---
#     # x, y, cx, cy = Inches(0.5), Inches(1), Inches(11), Inches(3)
#     # shape = slide.shapes.add_table(len(lap_metrics),len(lap_metrics.columns),x, y, cx,cy)
    
#     # tbl =  shape._element.graphic.graphicData.tbl
#     # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
#     # tbl[0][-1].text = style_id    
#     # table = shape.table
    
#     cols=prs.slides[c].shapes[2].table.columns
#     cols[1].width=Inches(1.4)
#     prs.slides[c].shapes[2].table.width=Inches(12)
#     prs.slides[c].shapes[2].table.height=Inches(5)
#     # cols[0].width=Inches(1.4)
#     # rows=table.rows
#     for i in range(len(lap_metrics)):
#         for j in range(len(lap_metrics.columns)):
#             if ((i in (5,11,13,19,20))& (j>=2)):
#                 prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])+'%'
#             else:
#                 prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])
                    
                
#                 prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_bottom = Inches(0.08)
#                 prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_left = 0
#                 prs.slides[c].shapes[2].table.cell(i,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
#                 prs.slides[c].shapes[2].table.cell(i,j).text_frame.word_wrap = False
#                 prs.slides[c].shapes[2].table.cell(i,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#                 #prs.slides[c].shapes[2].table.cell(i,j).text_frame.fit_text()
#                 # if i==0:
#                 #     cell=table.cell(i, j)
#                 #     fill = cell.fill
#                 #     fill.solid()
#                 #     fill.fore_color.rgb = RGBColor(173, 216, 230)
    
                
    
    
#     prs.slides[c].shapes[2].table.cell(1,0).text="Business Number"
#     prs.slides[c].shapes[2].table.cell(9,0).text="ROI"
#     prs.slides[c].shapes[2].table.cell(11,0).text="%Insurance"
#     prs.slides[c].shapes[2].table.cell(13,0).text="PF%"
#     prs.slides[c].shapes[2].table.cell(14,0).text="Number of DSAs"
#     prs.slides[c].shapes[2].table.cell(16,0).text="Incentive"
#     prs.slides[c].shapes[2].table.cell(19,0).text="LTD Booking to Login"
    
#     for cell in iter_cells(prs.slides[c].shapes[2].table):
#         for paragraph in cell.text_frame.paragraphs:
#             for run in paragraph.runs:
#                 if title=="LAP Business Metrics":
#                     run.font.size = Pt(11)
#                 else:
#                     run.font.size = Pt(14)
#                 run.font.style='Calibri'
#                 run.font.color.rgb = RGBColor(0,0,0)
#                 run.font.bold = False
    
#     return(prs)
todays_date = date.today()
todays_date=date(2022,12,31)
MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
MTD='2022-12'
# create_slide(prs,lap_metrics,"LAP Business Metrics",13)            
# create_slide(prs,hl_metrics,"HL Business Metrics",14)       
# create_slide(prs,gpl_metrics,"GPL Business Metrics",15)
# create_slide(prs,ngpl_metrics,"N-GPL Business Metrics",16)
# create_slide(prs,ahl_metrics,"AHL Business Metrics",17)
# prs.save(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\AUG\OPERATING REVIEW"+MTD +".pptx")
def get_quarter(M):
    if M in [4,5,6]:
        return [4,4]
    elif M in [7,8,9]:
        return [7,1]
    elif M in [10,11,12]:
        return [10,2]
    elif M in [1,2,3]:
        return [1,3]



def create_qslide(prs,lap_metrics_dict,title,c):
    
    for s in range(c,c+int(len(lap_metrics_dict.keys()))):
        lap_metrics=lap_metrics_dict[list(lap_metrics_dict.keys())[s-c]]
        qprs.slides[s].shapes[0].text=title+":"+str(list(lap_metrics_dict.keys())[s-c])
        #lap_metrics.reset_index(level=0, inplace=True)
        lap_metrics.insert(0,' ',' ')
        # slide = prs.slides.add_slide(prs.slide_layouts[5])
        # slide.shapes.title.text = title
        # slide.shapes.title.top = Inches(0.5)
        # slide.shapes.title.left = Inches(0.5)
        # slide.shapes.title.width = Inches(10)
    
        
        # # ---add table to slide---
        # x, y, cx, cy = Inches(0.5), Inches(1), Inches(11), Inches(3)
        # shape = slide.shapes.add_table(len(lap_metrics),len(lap_metrics.columns),x, y, cx,cy)
        
        # tbl =  shape._element.graphic.graphicData.tbl
        # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
        # tbl[0][-1].text = style_id    
        # table = shape.table
        quarters={4:[4,5,6],7:[7,8,9],10:[10,11,12],1:[1,2,3]}
        Cq=10
        if Cq==1:
            Pq=10
        else:
            Pq=Cq-3
        qlist=quarters[Pq]
        cols=prs.slides[s].shapes[3].table.columns
        cols[1].width=Inches(1.4)
        prs.slides[s].shapes[3].table.width=Inches(12)
        prs.slides[s].shapes[3].table.height=Inches(5)
        # cols[0].width=Inches(1.4)
        # rows=table.rows
        for i in range(len(lap_metrics)):
            for j in range(len(lap_metrics.columns)):
 
                if ((i in (4,5,8,11,12,13,14,16,22,23))& (j>2)):
                    prs.slides[s].shapes[3].table.cell(i+1,j).text=str(lap_metrics.iloc[i,j])+'%'
                else:
                    prs.slides[s].shapes[3].table.cell(i+1,j).text=str(lap_metrics.iloc[i,j])
                        
                    
                    prs.slides[s].shapes[3].table.cell(i+1,j).text_frame.margin_bottom = Inches(0.08)
                    prs.slides[s].shapes[3].table.cell(i+1,j).text_frame.margin_left = 0
                    prs.slides[s].shapes[3].table.cell(i+1,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
                    prs.slides[s].shapes[3].table.cell(i+1,j).text_frame.word_wrap = False
                    prs.slides[s].shapes[3].table.cell(i+1,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    #prs.slides[c].shapes[3].table.cell(i,j).text_frame.fit_text()
                    # if i==0:
                    #     cell=table.cell(i, j)
                    #     fill = cell.fill
                    #     fill.solid()
                    #     fill.fore_color.rgb = RGBColor(173, 216, 230)
        
        prs.slides[s].shapes[3].table.cell(1,3).text=str(calendar.month_abbr[qlist[0]])+'-22'
        prs.slides[s].shapes[3].table.cell(1,7).text=str(calendar.month_abbr[qlist[0]])+'-22'
        prs.slides[s].shapes[3].table.cell(1,4).text=str(calendar.month_abbr[qlist[1]])+'-22'
        prs.slides[s].shapes[3].table.cell(1,8).text=str(calendar.month_abbr[qlist[1]])+'-22'  
        prs.slides[s].shapes[3].table.cell(1,5).text=str(calendar.month_abbr[qlist[2]])+'-22' 
        prs.slides[s].shapes[3].table.cell(1,9).text=str(calendar.month_abbr[qlist[2]])+'-22' 
        prs.slides[s].shapes[3].table.cell(1,6).text='Q2 FY23' 
        prs.slides[s].shapes[3].table.cell(1,10).text='Q2 FY23' 
        prs.slides[s].shapes[3].table.cell(0,7).text=str(list(lap_metrics_dict.keys())[s-c])
                  
        
        
        prs.slides[s].shapes[3].table.cell(2,0).text="Business Number"
        prs.slides[s].shapes[3].table.cell(12,0).text="ROI"
        prs.slides[s].shapes[3].table.cell(14,0).text="%Insurance"
        prs.slides[s].shapes[3].table.cell(16,0).text="PF%"
        prs.slides[s].shapes[3].table.cell(17,0).text="Number of DSAs"
        prs.slides[s].shapes[3].table.cell(19,0).text="Incentive"
        prs.slides[s].shapes[3].table.cell(22,0).text="LTD Booking to Login"
        
        for cell in iter_cells(prs.slides[s].shapes[3].table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if title=="LAP Business Metrics":
                        run.font.size = Pt(11)
                    else:
                        run.font.size = Pt(14)
                    run.font.style='Calibri'
                    run.font.color.rgb = RGBColor(0,0,0)
                    run.font.bold = False
    
    return(prs)
def quarter_lap(LOAN):
    LOAN='HL'
    YTD=[]
    cur_month=todays_date.month
    while cur_month>=4:
        YTD.append(cur_month)
        cur_month=cur_month-1
    ytd=[]
    for i in YTD:
        ytd.append(str(todays_date.year)+"-"+str(i))
    
    td=date.today()  
    Cq=get_quarter(td.month)[0]
    quarters={4:[4,5,6],7:[7,8,9],10:[10,11,12],1:[1,2,3]}
    Cq=1
    if Cq==1:
        Pq=10
    else:
        Pq=Cq-3
    qlist=quarters[Pq]
    laplist=[]
    fd_df=fd.copy()
    df=mydata.copy()
    old_branches=['Target', 'Total','Ahmedabad', 'Bangalore','Delhi','Mumbai','Pune']
    if LOAN=='LAP':
        for i in qlist:
            
            if i%3 ==1:
                mt=m1trgt
            elif i%3 ==2:
                mt=m2trgt
            else:
                mt=m3trgt
            run_date=date(2022, i, 10)
            lap_metrics=buss_metrics(LOAN,df,fd_df,run_date)
            lap_metrics.columns=lap_metrics.loc['FINBRANCH']
            if int(run_date.month)>7:
                new_branches=list(set(lap_metrics.columns)-set(old_branches))
                lap_metrics=lap_metrics[(old_branches+new_branches)]
            else:
                new_branches=list(set(lap_metrics.columns)-set(old_branches))
                lap_metrics=lap_metrics[(old_branches+new_branches)]
            
            lap_metrics=set_target(lap_metrics,mt['LAP_mtd_target_df'],i) 
            lap_metrics.loc['DISB_Act (%)','Target']='90%'
            lap_metrics.loc['ROI ('+calendar.month_name[run_date.month-1][:3]+')','Target']='8.75%'
            lap_metrics.loc['ROI (Till'+calendar.month_name[run_date.month-2][:3]+')','Target']='8.75%'
            lap_metrics.loc['GROSS%','Target']='1.5%'
            lap_metrics.loc['PF%','Target']='0.9%'
            lap_metrics.loc['Cumulative','Target']='1.35%'
            lap_metrics.loc['Value%','Target']='33%'
            lap_metrics.loc['Volume%','Target']='33%'
            laplist.append(lap_metrics)
    elif LOAN in ['HL','GPL','NGPL']:
        for i in qlist:
            if i%3 ==1:
                mt=m1trgt
                lap_metrics=buss_metrics(LOAN,df2,fd_df,date(2022, i, 10))
            elif i%3 ==2:
                mt=m2trgt
                lap_metrics=buss_metrics(LOAN,df1,fd_df,date(2022, i, 10))
            else:
                mt=m3trgt
                lap_metrics=buss_metrics(LOAN,df,fd_df,date(2022, i, 10))
            run_date=date(2022, i, 10)
            
            lap_metrics.columns=lap_metrics.loc['FINBRANCH']
            new_branches=list(set(lap_metrics.columns)-set(old_branches))
            lap_metrics=lap_metrics[(old_branches+new_branches)]
            
            if LOAN=='HL':
                lap_metrics=set_target(lap_metrics,mt['HL_mtd_target_df'],i) 
            elif LOAN=='GPL':
                lap_metrics=set_target(lap_metrics,mt['GPL_mtd_target_df'],i) 
            elif LOAN=='NGPL':
                lap_metrics=set_target(lap_metrics,mt['NGPL_mtd_target_df'],i) 
            
            lap_metrics.loc['ROI ('+calendar.month_name[run_date.month-1][:3]+')','Target']='7.75%'
            lap_metrics.loc['ROI (Till'+calendar.month_name[run_date.month-2][:3]+')','Target']='7.75%'
            lap_metrics.loc['GROSS%','Target']='2%'
            lap_metrics.loc['PF%','Target']='0.02%'
            lap_metrics.loc['Cumulative','Target']='0.45%'
            lap_metrics.loc['Value%','Target']='50%'
            lap_metrics.loc['Volume%','Target']='50%'
            laplist.append(lap_metrics)
    ####################  3 Months METRICS READY ################################################
          
    '''
    COMPUTING Quarter TOTAL METRICS
    '''            
    BM_df = df[[ 'FINREFERENCE','NET_PREMIUM','BOOK_YEAR_MONTH','FINBRANCH','FINTYPE','STATUS', 'WROI','BOOKED_AMOUNT' ,'PROCESSING_FEE','REQLOANAMT','LOGINSTATUS','GPLFLAG_SANCTIONS','PRINCIPAL_OUTSTANDING','DISB_AMOUNT']]
    findisb=fd[['FINREFERENCE','DISBSEQ','DISBAMOUNT']]
    findisb=findisb[(findisb['DISBSEQ']==1)]
    findisb=findisb.sort_values(by='FINREFERENCE')
    BM_df=BM_df.sort_values(by='FINREFERENCE')
    fd_bm= pd.merge(BM_df, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')
    QTD=[]
    for i in qlist: 
        if Pq==10:
            year=run_date.year
        else:
            year=run_date.year
        QTD.append(str(year)+'-'+str(i))
    MTD=QTD
    
    if LOAN== 'LAP':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'LP') | (fd_bm['FINTYPE'] =='NP'))&(fd_bm.BOOK_YEAR_MONTH.isin(MTD))]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'LP') |  (bm_df['FINTYPE'] =='NP'))  ]
    elif LOAN== 'HL':
        fd_bm= fd_bm[ (fd_bm['FINBRANCH'].isin(['Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']))&(fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'HL') | (fd_bm['FINTYPE'] =='HT') | (fd_bm['FINTYPE'] =='FL') | (fd_bm['FINTYPE'] =='FT') | (fd_bm['FINTYPE'] =='LT')| (fd_bm['FINTYPE'] =='AHL'))&(fd_bm.BOOK_YEAR_MONTH.isin(MTD))]
        bm_df= BM_df
        bm_df= bm_df[(bm_df['FINBRANCH'].isin(['Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']))& ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'HL') | (bm_df['FINTYPE'] =='HT') | (bm_df['FINTYPE'] =='FL') | (bm_df['FINTYPE'] =='FT') | (bm_df['FINTYPE'] =='LT')| (bm_df['FINTYPE'] =='AHL'))  ]
    elif LOAN== 'GPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'] == 'GPL'))&(fd_bm.BOOK_YEAR_MONTH.isin(MTD))]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'] == 'GPL')  ]
    elif LOAN== 'NGPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL'])))&(fd_bm.BOOK_YEAR_MONTH.isin(MTD))]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL']))  ]
    elif LOAN== 'AHL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'AHL'))&(fd_bm.BOOK_YEAR_MONTH.isin(MTD))]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['FINTYPE'] == 'AHL')  ]
    
    # bm_df_hl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &((bm_df['FINTYPE'] == 'HL') |  (bm_df['FINTYPE'] =='HT'))]
    # bm_df_gpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='GPL')]
    # bm_df_ngpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='NON GPL')]
    bm_df['one_count']=np.ones(len(bm_df), dtype = int)
    
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','FINREFERENCE','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    # def buss_metrics(bm_df):
        #bm_df_group = bm_df.groupby(['FINBRANCH'],as_index=False).sum(['WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING' ])
        #bm_df_group['BOOK_TARGET'] =
        
    fd_bm_group=fd_bm.groupby(['FINBRANCH'],as_index=False).sum(['WROI', 'BOOKED_AMOUNT', 'PROCESSING_FEE','REQLOANAMT','PRINCIPAL_OUTSTANDING', 'DISB_AMOUNT', 'DISBSEQ', 'DISBAMOUNT'])
    fd_bm_group=fd_bm_group[['FINBRANCH','DISBAMOUNT']]
    bm_df_book=bm_df[(bm_df.STATUS== 'Booked')&(bm_df.BOOK_YEAR_MONTH.isin(MTD))]
    # bm_df_book=pd.merge(bm_df_book, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')
    bm_df_book['one_count']=np.ones(len(bm_df_book), dtype = int)
    bm_df_BOOK_group=bm_df_book.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    
    
    b=bm_df[(bm_df.STATUS== 'Booked')]
    c=bm_df[((bm_df.STATUS== 'Booked') & (bm_df.BOOK_YEAR_MONTH.isin(ytd)))]
    c_group=c.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','WROI'])
    c_group=c_group[['FINBRANCH','BOOKED_AMOUNT','WROI']]
    c_group.rename(columns = {'BOOKED_AMOUNT':'BOOKED_AMOUNT_M1'}, inplace = True)
    c_group.rename(columns = {'WROI':'WROI_M1'}, inplace = True)
    b_group=b.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','one_count' ])
    b_group=b_group[['FINBRANCH','PRINCIPAL_OUTSTANDING','one_count','BOOKED_AMOUNT' ]]
    b_group.rename(columns = {'one_count':'one_count_tot_book'}, inplace = True)
    b_group.rename(columns = {'BOOKED_AMOUNT':'TOTAL_BOOKED_AMOUNTytd'}, inplace = True)
    
    
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','BOOKED_AMOUNT','WROI','NET_PREMIUM','one_count','PROCESSING_FEE']]
    bm_df_BOOK_group.rename(columns = {'one_count':'one_count_MTD'}, inplace = True)
    bm_df_group=bm_df_group[['FINBRANCH','REQLOANAMT','one_count']]
    
    merged=pd.merge(bm_df_group, bm_df_BOOK_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged2=pd.merge(merged, c_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged3=pd.merge(merged2, b_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged4=pd.merge(merged3, fd_bm_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    if LOAN=='LAP' and (int(MTD[2][5])>6) :
        print('yes')
        newb=merged4.copy()
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.set_index('FINBRANCH',inplace=True)
        newb=newb.transpose()
        newb=newb[['Hyderabad','Indore', 'Jaipur', 'Chandigarh', 'Chennai', 'Surat']]
        newb=newb.transpose()
        newb.reset_index(inplace=True)
        newb.loc[11]=newb.sum()
        newb.loc[11,'FINBRANCH']='New Branches'
        newb['BOOK_VALUE'] = newb.BOOKED_AMOUNT/10000000
        #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]
    
        newb['BOOK_VOL'] = newb.one_count_MTD
        # newb['DISB AMNT Tranch 1']=newb.DISBAMOUNT/1000000000
        newb['DISB_Act (%)']=(newb.DISBAMOUNT)/newb.BOOKED_AMOUNT
        newb['AUM']=newb.PRINCIPAL_OUTSTANDING/10000000
        newb['ROI MTD'] = newb.WROI / newb.BOOKED_AMOUNT
        newb['ROI M-1'] = newb.WROI_M1 /newb.BOOKED_AMOUNT_M1
        newb['GROSS%']=100*newb.NET_PREMIUM/newb.BOOKED_AMOUNT
        newb['PF%']=100*newb.PROCESSING_FEE / newb.BOOKED_AMOUNT
        newb['LTD_volume']=100*newb.one_count_tot_book/newb.one_count
        newb['LTD_Value']=100*(newb.TOTAL_BOOKED_AMOUNTytd*100)/newb.REQLOANAMT
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    
        newb.loc['Total']= newb.sum()
        if sum(newb.BOOKED_AMOUNT)!=0:
            newb.loc['Total','DISB_Act (%)']=sum(newb.DISBAMOUNT)/sum(newb.BOOKED_AMOUNT)
            newb.loc['Total','ROI MTD']=sum(newb.WROI ) /sum( newb.BOOKED_AMOUNT)
            newb.loc['Total','ROI M-1']=sum(newb.WROI_M1)/sum(newb.BOOKED_AMOUNT_M1)
            newb.loc['Total','GROSS%']=100*sum(newb.NET_PREMIUM)/sum(newb.BOOKED_AMOUNT)
            newb.loc['Total','PF%']=100*sum(newb.PROCESSING_FEE ) / sum(newb.BOOKED_AMOUNT)
        else:
            newb.loc['Total','DISB_Act (%)']=0
            newb.loc['Total','ROI MTD']=0
            newb.loc['Total','ROI M-1']=0
            newb.loc['Total','GROSS%']=0
            newb.loc['Total','PF%']=0
        newb.loc['Total','LTD_volume']=100*sum(newb.one_count_tot_book)/sum(newb.one_count)
        newb.loc['Total','LTD_Value']=100*sum(newb.TOTAL_BOOKED_AMOUNTytd)*(100)/sum(newb.REQLOANAMT)
        newb.loc['Total','FINBRANCH']='Total'
        newb.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
        newb["new"] = range(1,len(newb)+1)
        newb.loc['Total','new'] = 0
        newb=newb.sort_values(by="new").drop('new', axis=1)
        newb=newb.round(decimals=2)
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
        newb.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
        newb.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
        newb.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
        newb.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
        newb.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
        newb.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
        newb.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
        newb['Booking Tar(Cr)']=' '
        newb['Value %']=' '
        newb['Booking Tar(#)']=' '
        newb['Volume %']=' '
        newb['AUM Tar(Cr)']=' '
        newb['AUM Actual %']=' '
        newb['Net']=' '
        newb['Empaneled']=' '
        newb['Active']=' '
        newb['Channel']=' '
        newb['Employee']=' '
        newb['Cumulative']=' '
        newb.columns
        newb=newb.round(decimals=2)
        newb['Value%']=round(newb['Value%'])
        newb['Volume%']=round(newb['Volume%'])
        newb=newb[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Value %','DISB_Act (%)','Booking Tar(#)', 'Booking Act(#)','Volume %','AUM Tar(Cr)','AUM Act(Cr)','AUM Actual %', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]        #newb.set_index('FINBRANCH',inplace=True)
        
    
        finaln=newb.transpose()
        
        finaln.insert(0,'Target',0)
        finaln.loc['FINBRANCH','Target']='Target'
    
    
    merged4['BOOK_VALUE'] = merged4.BOOKED_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]

    merged4['BOOK_VOL'] = merged4.one_count_MTD
    # merged4['DISB AMNT Tranch 1']=merged4.DISBAMOUNT/1000000000
    merged4['DISB_Act (%)']=(merged4.DISBAMOUNT)/merged4.BOOKED_AMOUNT
    merged4['AUM']=merged4.PRINCIPAL_OUTSTANDING/10000000
    merged4['ROI MTD'] = merged4.WROI / merged4.BOOKED_AMOUNT
    merged4['ROI M-1'] = merged4.WROI_M1 /merged4.BOOKED_AMOUNT_M1
    merged4['GROSS%']=100*merged4.NET_PREMIUM/merged4.BOOKED_AMOUNT
    merged4['PF%']=100*merged4.PROCESSING_FEE / merged4.BOOKED_AMOUNT
    merged4['LTD_volume']=100*merged4.one_count_tot_book/merged4.one_count
    merged4['LTD_Value']=100*(merged4.TOTAL_BOOKED_AMOUNTytd*100)/merged4.REQLOANAMT
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    merged4.loc['Total']= merged4.sum()
    merged4.loc['Total','DISB_Act (%)']=sum(merged4.DISBAMOUNT)/sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','ROI MTD']=sum(merged4.WROI ) /sum( merged4.BOOKED_AMOUNT)
    merged4.loc['Total','ROI M-1']=sum(merged4.WROI_M1)/sum(merged4.BOOKED_AMOUNT_M1)
    merged4.loc['Total','GROSS%']=100*sum(merged4.NET_PREMIUM)/sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','PF%']=100*sum(merged4.PROCESSING_FEE ) / sum(merged4.BOOKED_AMOUNT)
    merged4.loc['Total','LTD_volume']=100*sum(merged4.one_count_tot_book)/sum(merged4.one_count)
    merged4.loc['Total','LTD_Value']=100*sum(merged4.TOTAL_BOOKED_AMOUNTytd)*(100)/sum(merged4.REQLOANAMT)
    merged4.loc['Total','FINBRANCH']='Total'
    merged4.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    merged4["new"] = range(1,len(merged4)+1)
    merged4.loc['Total','new'] = 0
    merged4=merged4.sort_values(by="new").drop('new', axis=1)
    merged4=merged4.round(decimals=2)
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    merged4.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    merged4.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    merged4.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    merged4.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    merged4.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    merged4.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    merged4['Booking Tar(Cr)']=' '
    merged4['Value %']=' '
    merged4['Booking Tar(#)']=' '
    merged4['Volume %']=' '
    merged4['AUM Tar(Cr)']=' '
    merged4['AUM Actual %']=' '
    merged4['Net']=' '
    merged4['Empaneled']=' '
    merged4['Active']=' '
    merged4['Channel']=' '
    merged4['Employee']=' '
    merged4['Cumulative']=' '
    merged4.columns
    merged4=merged4.round(decimals=2)
    merged4['Value%']=round(merged4['Value%'])
    merged4['Volume%']=round(merged4['Volume%'])
    merged4=merged4[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Value %','DISB_Act (%)','Booking Tar(#)', 'Booking Act(#)','Volume %','AUM Tar(Cr)','AUM Act(Cr)','AUM Actual %', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #merged4.set_index('FINBRANCH',inplace=True)
    

    final=merged4.transpose()
    
    final.insert(0,'Target',0)
    final.loc['FINBRANCH','Target']='Target'
    if LOAN=='LAP' and (int(MTD[2][5])>6) :
        final.reset_index(inplace=True)
        finaln.reset_index(inplace=True)
        finaln=finaln[['index',11]]
        final=pd.merge(final,finaln,left_on=('index'),right_on=('index'),how='outer')
        final.set_index('index',inplace=True)
    final.columns=final.loc['FINBRANCH']
    if int(run_date.month)>7 and LOAN=='LAP':
        new_branches=list(set(final.columns)-set(old_branches))
        final=final[(old_branches+new_branches)]
        
    elif LOAN=='LAP':
        new_branches=list(set(final.columns)-set(old_branches))
        final=final[(old_branches+new_branches)]
        final=set_target(final,mt['LAP_mtd_target_df'],10) 
    if LOAN in ['HL','GPL','NGPL']:
        new_branches=list(set(final.columns)-set(old_branches))
        final=final[(old_branches+new_branches)]
        final=set_target(final,mt['HL_mtd_target_df'],10) 
    
    final.loc['DISB_Act (%)','Target']='90%'
    final.loc['ROI ('+calendar.month_name[run_date.month-1][:3]+')','Target']='8.75%'
    final.loc['ROI (Till'+calendar.month_name[run_date.month-2][:3]+')','Target']='8.75%'
    final.loc['GROSS%','Target']='1.5%'
    final.loc['PF%','Target']='0.9%'
    final.loc['Cumulative','Target']='1.35%'
    final.loc['Value%','Target']='33%'
    final.loc['Volume%','Target']='33%'
    m1=laplist[0]
    m2=laplist[1]
    m3=laplist[2]
    brnc_q={}
    for i in m1.columns[2:]:
        bi1=m1.copy()
        bi1=bi1[['Target','Total',i]]
        if i in m2.columns[2:]:
            bi2=m2.copy()
            bi2=bi2[['Total',i]]
        else:
            bi2=m2.copy()
            bi2=bi2[['Total']]
            bi2[i]=0
        if i in m3.columns[2:]:
            bi3=m3.copy()
            bi3=bi3[['Total',i]]
        else:
            bi3=m3.copy()
            bi3=bi3[['Total']]
            bi3[i]=0
        bi1.rename(columns={'Total':'Total1'},inplace=True)
        bi1.rename(columns={i:i+"1"},inplace=True)
        bi2.rename(columns={'Total':'Total2'},inplace=True)
        bi2.rename(columns={i:i+"2"},inplace=True)
        bi3.rename(columns={'Total':'Total3'},inplace=True)
        bi3.rename(columns={i:i+"3"},inplace=True)
        bi1.reset_index(inplace=True)
        bi2.reset_index(inplace=True)
        bi3.reset_index(inplace=True)
        bi1.iloc[11,0]='ROI (Month)'
        bi1.iloc[12,0]='ROI (Till M-1)'
        bi2.iloc[11,0]='ROI (Month)'
        bi2.iloc[12,0]='ROI (Till M-1)'
        bi3.iloc[11,0]='ROI (Month)'
        bi3.iloc[12,0]='ROI (Till M-1)'
        bi1=bi1.merge(bi2,on='index')
        bi1=bi1.merge(bi3,on='index')
        qtotal=final.copy()
        qtotal=qtotal[['Total',i]]
        qtotal.rename(columns={'Total':'TotalQT'},inplace=True)
        qtotal.rename(columns={i:i+"QB"},inplace=True)
        qtotal.reset_index(inplace=True)
        qtotal.iloc[11,0]='ROI (Month)'
        qtotal.iloc[12,0]='ROI (Till M-1)'
        bi1=bi1.merge(qtotal,on='index')
        bi1.iloc[3,8]=0
        bi1.iloc[7,8]=0
        bi1.iloc[3,9]=0
        bi1.iloc[7,9]=0
        bi1['TotalQT'][1:8]=bi1['TotalQT'][1:8].astype('float')
        if i in  ['Ahmedabad','Mumbai','Pune','Bangalore','Delhi']:
            bi1[i+'QB'][1:8]=bi1[i+'QB'][1:8].astype('float')
        bi1=bi1[['index','Target' ,'Total1', 'Total2','Total3', 'TotalQT',i+'1', i+'2', i+'3',i+'QB']]
        bi1.loc[1,'TotalQT']=bi1.loc[1,'Total1']+bi1.loc[1,'Total2']+bi1.loc[1,'Total3']
        bi1.loc[1,i+'QB']=bi1.loc[1,i+'1']+bi1.loc[1,i+'2']+bi1.loc[1,i+'3']
        bi1.loc[2,'TotalQT']=bi1.loc[2,'Total1']+bi1.loc[2,'Total2']+bi1.loc[2,'Total3']
        bi1.loc[2,i+'QB']=bi1.loc[2,i+'1']+bi1.loc[2,i+'2']+bi1.loc[2,i+'3']
        bi1.loc[3,'TotalQT']=str(round(bi1.loc[2,'TotalQT']/bi1.loc[1,'TotalQT']*100))+'%'
        
        bi1.loc[5,i+'QB']=bi1.loc[5,i+'1']+bi1.loc[5,i+'2']+bi1.loc[5,i+'3']
        bi1.loc[5,'TotalQT']=bi1.loc[5,'Total1']+bi1.loc[5,'Total2']+bi1.loc[5,'Total3']
        if i in  ['Ahmedabad','Mumbai','Pune','Bangalore','Delhi']:
            bi1.loc[7,i+'QB']=str(round(bi1.loc[6,i+'QB']/bi1.loc[5,i+'QB']*100))+'%'
            bi1.loc[3,i+'QB']=str(round(bi1.loc[2,i+'QB']/bi1.loc[1,i+'QB']*100))+'%'
        bi1.loc[7,'TotalQT']=str(round(bi1.loc[6,'TotalQT']/bi1.loc[5,'TotalQT']*100))+'%'
        bi1.loc[8,i+'QB']=bi1.loc[8,i+'3']
        bi1.loc[8,'TotalQT']=bi1.loc[8,'Total3']
        bi1.loc[9,i+'QB']=bi1.loc[9,i+'3']
        bi1.loc[9,'TotalQT']=bi1.loc[9,'Total3']
        bi1.loc[10,i+'QB']=str(round(bi1.loc[9,'TotalQT']/bi1.loc[8,'TotalQT']*100))+'%'
        bi1.loc[10,i+'QB']=str(round(bi1.loc[9,i+'QB']/bi1.loc[8,i+'QB']*100))+'%'
        brnc_q[i]=bi1
    return(brnc_q)
quarter_lap_metrics=quarter_lap('LAP')
quarter_hl_metrics=quarter_lap('HL')
quarter_gpl_metrics=quarter_lap('GPL')
quarter_ngpl_metrics=quarter_lap('NGPL')
# quarter_ahl_metrics=
# if int(todays_date.month)==get_quarter(int(MTD[5]))[0]:
qprs= Presentation(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\OPSREVIEWDECKQUARTER.pptx")
create_qslide(qprs,quarter_lap_metrics,"LAP Business Metrics",13)            
create_qslide(qprs,quarter_hl_metrics,"HL Business Metrics",13)       
create_qslide(qprs,quarter_gpl_metrics,"GPL Business Metrics",13+len(quarter_lap_metrics.keys())+len(quarter_hl_metrics.keys()))
create_qslide(qprs,quarter_ngpl_metrics,"N-GPL Business Metrics",13+len(quarter_lap_metrics.keys())+len(quarter_hl_metrics.keys())+len(quarter_gpl_metrics.keys()))
# create_qslide(qprs,quarter_ahl_metrics,"AHL Business Metrics",43,6,target)
qprs.save(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\OPERATING REVIEW QUARTER_"+str(get_quarter(int(MTD[5]))[1])+".pptx")
with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\quarter_report_hl.xlsx") as writer:
    for key,val in quarter_hl_metrics.items():
        quarter_hl_metrics[key].to_excel(writer, sheet_name=key)

############################# Unit Economics #############################


import pandas as pd
import calendar
import os
import datetime  
from datetime import timedelta
import numpy as np
from datetime import date
import pyodbc
pyodbc.drivers()
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx import Presentation
config = {}
with open(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\config.txt","r") as file:
    for line in file:
        key, value = line.strip().split(",")
        config[key] = value

del key,line,value,file

conx = pyodbc.connect('DRIVER={SnowflakeDSIIDriver};SERVER='+str(config['server'])+';Warehouse = COMPUTE_WH;DATABASE='+str(config['db1'])+';UID='+ str(config['user1']) +';PWD='+ str(config['password']),autocommit=True)

cur = conx.cursor()


cur.execute("USE DATABASE SAMPLE_DB")

query = "select * from dtcron_master_02102022;"
mydata = pd.read_sql(query, conx)
# insurance=pd.read_sql("SELECT * FROM MIS_INSURANCE_V",conx)
# insurance=insurance[['LANID','INSURANCE_TYPE']]
# mydata=pd.merge(dtcron,insurance,left_on=['FINREFERENCE'],right_on=['LANID'], how='outer')
# mydata.drop(columns='LANID',inplace=True)
# mydata['INSURANCE_TYPE'].replace(to_replace=[np.nan,'',' '],value="0",inplace=True)

# #mydata=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Downloads\DTCRON_01082022_NetInsuranceIncome_avi.xlsx")
# # prs = Presentation()
# # prs.slide_width = 11887200
# # prs.slide_height = 6686550
# # def create_slide(prs,lap_metrics,title):
todays_date = date.today()  
# iname=[0,
#  'ABH Group Active Secure (PA)',
#  'ABH Group Active Secure - co borrower 1 (PA)',
#  'ABH Group Protect Secure (Cancer & Heart)',
#  'ABH Group Protect Secure-CoBorrower1 Cancer,Heart',
#  'ABH Heart secure (Group Active Secure)',
#  'ABH Personal Accident Active Secure',
#  'ABSLI - Group Asset Assure Plan / ABSLI - GAAP',
#  'ABSLI GSS Level Borrower 1',
#  'ABSLI GSS Reducing Borrower 1',
#  'ABSLI GSS Reducing Borrower 2',
#  'Aditya Birla health Group Heart Secure',
#  'BAGIC Bharat Grah Raksha Borrower 2',
#  'BAGIC Bharat Grah Raksha Borrower1',
#  'BAGIC Bharat Grah Raksha Borrower2',
#  'BAGIC Credit Linked Health Plan Borrower 1',
#  'BAGIC Credit Linked Health Plan Borrower 2',
#  'New Plan of GCSPlus without Critical Illness',
#  'TAGIC - Group Credit Secure Plus / TAGIC - GCS+',
#  'TAGIC - Group MediCare',
#  'TAGIC - Property Insurance']
# ipercent=[0.0,
#  0.65,
#  0.0,
#  0.0,
#  0.0,
#  0.0,
#  0.0,
#  0.6000000000000001,
#  0.6000000000000001,
#  0.6000000000000001,
#  0.0,
#  0.0,
#  0.0,
#  0.4,
#  0.0,
#  0.5,
#  0.0,
#  0.6000000000000001,
#  0.6000000000000001,
#  0.35,
#  0.39999999999999997]
# insur_dc=dict(zip(iname,ipercent))
# mydata['InsPercent']=0
# for i in insur_dc.keys():
#     mydata['InsPercent']=np.where(mydata['INSURANCE_TYPE']==i,insur_dc[i],mydata['InsPercent'])
# # mydata['Net insurance income']=mydata['NET_PREMIUM']*mydata['InsPercent']*0.9/1.18
mydata['Net insurance income']=0
def UE_Branch_GPL(df_roi,MTD,prod_name):
    if type(MTD)==str:
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH']==MTD)]
    elif type(MTD)==list :
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH'].isin(MTD))]        
    prod2= df_roi[(df_roi[prod_name]==1)]
    prod2.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total','BOOKED_AMOUNT','Net insurance income'],inplace=True)
    prod_pos=prod2.groupby('FINBRANCH',as_index=False).sum(['PRINCIPAL_OUTSTANDING'])
    prod_branch1=prod.groupby('FINBRANCH',as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income'])
    prod_branch1.drop(columns = ['PRINCIPAL_OUTSTANDING'],inplace=True)
    prod_branch=pd.merge(prod_branch1,prod_pos ,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    prod.loc['Total']= prod.sum()
    prod_branch.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    return(prod_branch,prod)


def cal_c(df):
    df['ROI']=df.WROI/df.BOOKED_AMOUNT
    df['PF']=100*df.PROCESSING_FEE/df.BOOKED_AMOUNT
    df['GROSS']=100*df.NET_PREMIUM/df.BOOKED_AMOUNT
    df['NET INSURANCE INCOME']=100*(df['Net insurance income']/df.BOOKED_AMOUNT)
    
    return(df)
def drop_col_gpl(df):
    try:
        df.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total','Net insurance income'],inplace=True)
    except:
        pass
    df.set_index('FINBRANCH',inplace=True)
    return(df)


def UE_Branch_NGPL(df_roi,MTD,prod_name):
    if type(MTD)==str:
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH']==MTD)]
    elif type(MTD)==list :
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH'].isin(MTD))]     
    prod2= df_roi[(df_roi[prod_name]==1)]
    prod2.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'NON_GPL_regular','NON_GPL_demi','NON_GPL_cre','NON_GPL_plot','NON_GPL_ahl','NON_GPL_total','BOOKED_AMOUNT','Net insurance income'],inplace=True)
    prod_pos=prod2.groupby('FINBRANCH',as_index=False).sum(['PRINCIPAL_OUTSTANDING'])
    prod_branch1=prod.groupby('FINBRANCH',as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income'])
    prod_branch1.drop(columns = ['PRINCIPAL_OUTSTANDING'],inplace=True)
    prod_branch=pd.merge(prod_branch1,prod_pos ,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    prod.loc['Total']= prod.sum()
    prod_branch.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    return(prod_branch,prod)

def UE_Branch_HL(df_roi,MTD):
    if type(MTD)==str:
        prod=df_roi[df_roi['BOOK_YEAR_MONTH']==MTD]
    elif type(MTD)==list :
        prod=df_roi[df_roi['BOOK_YEAR_MONTH'].isin(MTD)]     
    prod2= df_roi.copy()
    prod2.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM','BOOKED_AMOUNT','Net insurance income'],inplace=True)
    prod_pos=prod2.groupby('FINBRANCH',as_index=False).sum(['PRINCIPAL_OUTSTANDING'])
    prod_branch1=prod.groupby('FINBRANCH',as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income'])
    prod_branch1.drop(columns = ['PRINCIPAL_OUTSTANDING'],inplace=True)
    prod_branch=pd.merge(prod_branch1,prod_pos ,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    prod.loc['Total']= prod.sum()
    prod_branch.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    return(prod_branch,prod)

def UE_Branch_LAP(df_roi,MTD,prod_name):
    if type(MTD)==str:
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH']==MTD)]
    elif type(MTD)==list :
        prod=df_roi[(df_roi[prod_name]==1)&(df_roi['BOOK_YEAR_MONTH'].isin(MTD))]     
    prod2= df_roi[(df_roi[prod_name]==1)]
    prod2.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'Lap_regular','Lap_industrial','Lap_booster','Lap_topup','Lap_lrd','Lap_total','BOOKED_AMOUNT','Net insurance income'],inplace=True)
    prod_pos=prod2.groupby('FINBRANCH',as_index=False).sum(['PRINCIPAL_OUTSTANDING'])
    prod_branch1=prod.groupby('FINBRANCH',as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income'])
    prod_branch1.drop(columns = ['PRINCIPAL_OUTSTANDING'],inplace=True)
    prod_branch=pd.merge(prod_branch1,prod_pos ,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    prod.loc['Total']= prod.sum()
    prod_branch.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    return(prod_branch,prod)

def drop_col_lap(df):
    try:
        df.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'Lap_regular','Lap_industrial','Lap_booster','Lap_topup','Lap_lrd','Lap_total','Net insurance income'],inplace=True)
    except:
        pass
    df.set_index('FINBRANCH',inplace=True)
    return(df)
def drop_col_hl(df):
    try:
        df.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM','Net insurance income'],inplace=True)
    except:
        pass
    df.set_index('FINBRANCH',inplace=True)
    return(df)

def drop_col_ngpl(df):
    try:
        df.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'NON_GPL_regular','NON_GPL_demi','NON_GPL_cre','NON_GPL_plot','NON_GPL_ahl','NON_GPL_total','Net insurance income'],inplace=True)
    except:
        pass
    df.set_index('FINBRANCH',inplace=True)
    return(df)

def transpose(df):
    df=df[['ROI','NET INSURANCE INCOME','PF','BOOKED_AMOUNT']]
    df=df.transpose()
    return df
       

def branch_func(df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15,df16,df17):
    df1=transpose(df1)
    df2=transpose(df2)
    df3=transpose(df3)
    df4=transpose(df4)
    df5=transpose(df5)
    df6=transpose(df6)
    df7=transpose(df7)
    df8=transpose(df8)
    df9=transpose(df9)
    df10=transpose(df10)
    df11=transpose(df11)
    df12=transpose(df12)
    df13=transpose(df13)
    df14=transpose(df14)
    df15=transpose(df15)
    df16=transpose(df16)
    df17=transpose(df17)
    return df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15,df16,df17
    

def Unit_Economics(mydata):
    
    
    df=mydata.copy()
    
    '''
     ------ LAP ROI ############################
    LAP_IND_ROI <- subset(mydata ,mydata$FINTYPE %in% c('LP','NP') & 
                          mydata$STATUS == 'Booked' & mydata$BOOK_YEAR_MONTH == MTD)
    
     Regular = LAP_Regular_ROI <- subset(mydata, mydata$FINTYPE %in% c('LP','NP') &
                                         !mydata$LOAN_PURPOSE %in% c('LAP Balance Transfer plus Top-up ','LAP Top Up') 
                              & mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked'
                              & !mydata$SUBPRODUCT == 'BOOSTER')
    	Industrial = subset(mydata, mydata$FINTYPE %in% c('LP','NP') & mydata$STATUS == 'Booked' & mydata$BOOK_YEAR_MONTH == MTD)
      
    	Booster =  LAP_Booster_ROI <- subset(mydata, mydata$FINTYPE %in% c('LP','NP') & 
                                          mydata$SUBPRODUCT == 'BOOSTER' 
                              & mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked' )
    
    	Topup = subset(mydata, mydata$FINTYPE %in% c('LP','NP') & mydata$STATUS == 'Booked' 
                            & mydata$BOOK_YEAR_MONTH == MTD
                            & mydata$INDUSTRIAL_PROPERTY_FLAG == 0
                            & mydata$LOAN_PURPOSE %in% c('LAP Balance Transfer plus Top-up ','LAP Top Up'))
    
    	Total = LAP_Total_ROI <- subset(mydata, mydata$FINTYPE %in% c('LP','NP') &
                                     mydata$STATUS == 'Booked'
                            & mydata$BOOK_YEAR_MONTH == MTD )
    '''
    
    
    
    lap_df = df[[ 'BOOK_YEAR_MONTH','FINBRANCH','FINTYPE','STATUS', 'WROI', 'CRE','STEPFINANCE', 
                 'GPLFLAG_SANCTIONS','BOOKED_AMOUNT' ,'SUBPRODUCT','LOAN_PURPOSE','INDUSTRIAL_PROPERTY_FLAG','PROCESSING_FEE','NET_PREMIUM','PRINCIPAL_OUTSTANDING','Net insurance income']]
    lap_df_roi = lap_df.copy()
    lap_df_roi = lap_df_roi[lap_df_roi.GPLFLAG_SANCTIONS == 'NIL']
    lap_df_roi = lap_df_roi[lap_df_roi.STATUS == 'Booked']
    lap_df_roi = lap_df_roi[(lap_df_roi['FINTYPE'] == 'LP') | (lap_df_roi['FINTYPE'] == 'NP') ]
    lap_df_roi = lap_df_roi[lap_df_roi['BOOK_YEAR_MONTH'].notna()]
    lap_df_roi['Lap_total'] = np.where( (lap_df_roi['FINTYPE'] == 'LP') |  (lap_df_roi['FINTYPE'] =='NP') &  (lap_df_roi['STATUS'] == 'Booked') , 1,0 )
    lap_df_roi['Lap_industrial'] = np.where((lap_df_roi['INDUSTRIAL_PROPERTY_FLAG'] == 1), 1,0 )
    lap_df_roi['Lap_industrial'].value_counts()
    lap_df_roi['Lap_topup'] = np.where((lap_df_roi['INDUSTRIAL_PROPERTY_FLAG'] == 0)
                                       & ((lap_df_roi['LOAN_PURPOSE']=='LAP Balance Transfer plus Top-up ')|(lap_df_roi['LOAN_PURPOSE']=='LAP Top Up')) , 1,0 )
    lap_df_roi['Lap_topup'].value_counts()
    lap_df_roi['Lap_booster'] = np.where( lap_df_roi['SUBPRODUCT'] == 'BOOSTER' , 1,0)
    lap_df_roi['Lap_lrd'] = np.where( lap_df_roi['SUBPRODUCT'] == 'LRD' , 1,0)
    lap_df_roi['Lap_regular'] = np.where( (lap_df_roi['Lap_industrial'] ==0 ) &  (lap_df_roi['Lap_booster']==0 )& 
                                      (lap_df_roi['Lap_topup'] ==0 )&(lap_df_roi['Lap_lrd'] ==0 ) , 1,0 )
    lap_df_roi['Lap_regular'].value_counts()
    lap_df_roi = lap_df_roi[['BOOK_YEAR_MONTH', 'FINBRANCH', 'WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT', 'Lap_regular','Lap_industrial','Lap_booster','Lap_topup','Lap_lrd','Lap_total','PRINCIPAL_OUTSTANDING','Net insurance income']]
    lap_roi_group = lap_df_roi.groupby(['BOOK_YEAR_MONTH', 'FINBRANCH','Lap_regular','Lap_industrial','Lap_booster','Lap_topup','Lap_lrd','Lap_total' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    lap_roi_group['LAP_ROI'] = lap_roi_group.WROI / lap_roi_group.BOOKED_AMOUNT
    lap_roi_group['LAP_PF'] = lap_roi_group.PROCESSING_FEE / lap_roi_group.BOOKED_AMOUNT
    lap_roi_group['LAP_NetP'] = lap_roi_group.NET_PREMIUM / lap_roi_group.BOOKED_AMOUNT
    
    #####################
    
    lap_roi_group_total = lap_df_roi.groupby([ 'BOOK_YEAR_MONTH','Lap_regular','Lap_industrial','Lap_booster','Lap_topup','Lap_lrd','Lap_total'],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    lap_roi_group_total['LAP_ROI'] = lap_roi_group_total.WROI / lap_roi_group_total.BOOKED_AMOUNT
    lap_roi_group_total['LAP_PF'] = lap_roi_group_total.PROCESSING_FEE / lap_roi_group_total.BOOKED_AMOUNT
    lap_roi_group_total['LAP_NetP'] = lap_roi_group_total.NET_PREMIUM / lap_roi_group_total.BOOKED_AMOUNT
    lap_roi_group.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    lap_roi_group_total.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    #lap_roi_group_total.to_csv("lap_roi_group_total.csv")
    
    ######################## LAP Completed #################################
    
    
    
    
    
    
    
    
    
    '''
    GPL ROI
    
    
    DEMI =GPL_DEMI_ROI <- subset(mydata,mydata$STEPFINANCE == 1 & mydata$GPLFLAG_SANCTIONS == 'GPL' & 
                             mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    
    
    
    CRE=GPL_CRE_ROI <- subset(mydata,mydata$CRE %in% c('CRE','CRE-H') & mydata$GPLFLAG_SANCTIONS == 'GPL' & 
                            mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    
    PLOT=GPL_PLOT_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'GPL' & 
                             mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    REGULAR= GPL_REGULAR_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'GPL' & 
                                mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked' 
                              & !mydata$CRE %in% c('CRE','CRE-H') & mydata$STEPFINANCE == 0)
    TOTAL=GPL_TOTAL_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'GPL' & 
                              mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked' )
    
    '''
    
    gpl_df=df[['STEPFINANCE','GPLFLAG_SANCTIONS','BOOK_YEAR_MONTH','STATUS','CRE','LOAN_PURPOSE','FINBRANCH','WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income']]
    gpl_df_roi=gpl_df.copy()
    gpl_df_roi = gpl_df_roi[gpl_df_roi.GPLFLAG_SANCTIONS == 'GPL']
    gpl_df_roi=gpl_df_roi[gpl_df_roi.STATUS == 'Booked']
    gpl_df_roi['LOAN_PURPOSE'].value_counts()
    gpl_df_roi['GPL_total'] = np.where( (gpl_df_roi['GPLFLAG_SANCTIONS'] == 'GPL')
                                       &  (gpl_df_roi['STATUS'] == 'Booked') , 1,0 )
    gpl_df_roi['GPL_total'].value_counts()
    gpl_df_roi['GPL_plot'] = np.where( ( (gpl_df_roi['LOAN_PURPOSE'] == 'Plot purchase only' ) 
                                       | (gpl_df_roi['LOAN_PURPOSE'] == 'Plot loan BT' )) , 1,0 )
    gpl_df_roi['GPL_plot'].value_counts()
    gpl_df_roi['GPL_demi'] = np.where((gpl_df_roi['STEPFINANCE'] == 1) , 1,0 )
    gpl_df_roi['GPL_demi'].value_counts()
    gpl_df_roi['GPL_cre'] = np.where( ((gpl_df_roi['CRE']=='CRE')|(gpl_df_roi['CRE']=='CRE-H')) , 1,0)
    gpl_df_roi['GPL_cre'].value_counts()
    gpl_df_roi['GPL_regular'] = np.where( (gpl_df_roi['GPL_plot'] == 0 ) &  (gpl_df_roi['GPL_demi'] ==0 )& 
                                      (gpl_df_roi['GPL_cre'] ==0 ), 1,0 )
    # gpl_df_roi['GPL_type'].value_counts()
    gpl_df_roi = gpl_df_roi[['PRINCIPAL_OUTSTANDING','BOOK_YEAR_MONTH', 'FINBRANCH', 'WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT', 'GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total','Net insurance income']]
    gpl_roi_group = gpl_df_roi.groupby(['BOOK_YEAR_MONTH', 'FINBRANCH','GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    gpl_roi_group['GPL_ROI'] = gpl_roi_group.WROI / gpl_roi_group.BOOKED_AMOUNT
    gpl_roi_group['GPL_PF'] = gpl_roi_group.PROCESSING_FEE / gpl_roi_group.BOOKED_AMOUNT
    gpl_roi_group['GPL_NetP'] = gpl_roi_group.NET_PREMIUM / gpl_roi_group.BOOKED_AMOUNT
    
    #####################
    
    gpl_roi_group_total = gpl_df_roi.groupby(['BOOK_YEAR_MONTH','GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    gpl_roi_group_total['GPL_ROI'] = gpl_roi_group_total.WROI / gpl_roi_group_total.BOOKED_AMOUNT
    gpl_roi_group_total['GPL_PF'] = gpl_roi_group_total.PROCESSING_FEE / gpl_roi_group_total.BOOKED_AMOUNT
    gpl_roi_group_total['GPL_NetP'] = gpl_roi_group_total.NET_PREMIUM / gpl_roi_group_total.BOOKED_AMOUNT
    gpl_roi_group.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    gpl_roi_group_total.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    #gpl_roi_group_total.to_csv("gpl_roi_group_total.csv")
    
    ######################## GPL Completed #################################
    
    
    
    
    '''
    NON GPL ROI
    
    
    DEMI =NONGPL_DEMI_ROI <- subset(mydata,mydata$STEPFINANCE == 1 & mydata$GPLFLAG_SANCTIONS == 'NON GPL' & 
                             mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    
    
    
    CRE=NONGPL_CRE_ROI <- subset(mydata,mydata$CRE %in% c('CRE','CRE-H') & mydata$GPLFLAG_SANCTIONS == 'NON GPL' & 
                            mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    
    PLOT=NONGPL_PLOT_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'NON GPL' & 
                             mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked')
    REGULAR= NONGPL_REGULAR_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'NON GPL' & 
                                mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked' 
                              & !mydata$CRE %in% c('CRE','CRE-H') & mydata$STEPFINANCE == 0)
    TOTAL=NONGPL_TOTAL_ROI <- subset(mydata, mydata$GPLFLAG_SANCTIONS == 'NON GPL' & 
                              mydata$BOOK_YEAR_MONTH == MTD & mydata$STATUS == 'Booked' )
    
    '''
    
    non_gpl_df=df[['STEPFINANCE','GPLFLAG_SANCTIONS','AHL_FLAG','FINTYPE','BOOK_YEAR_MONTH','STATUS','CRE','LOAN_PURPOSE','FINBRANCH','WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income']]
    non_gpl_df_roi=non_gpl_df.copy()
    non_gpl_df_roi = non_gpl_df_roi[non_gpl_df_roi.GPLFLAG_SANCTIONS.isin(['NON GPL','AHL'])]
    non_gpl_df_roi=non_gpl_df_roi[non_gpl_df_roi.STATUS == 'Booked']
    non_gpl_df_roi['LOAN_PURPOSE'].value_counts()
    non_gpl_df_roi['NON_GPL_total'] = np.where( (non_gpl_df_roi['GPLFLAG_SANCTIONS'] == 'NON GPL')
                                       &  (non_gpl_df_roi['STATUS'] == 'Booked') , 1,0 )
    non_gpl_df_roi['NON_GPL_total'].value_counts()
    non_gpl_df_roi['NON_GPL_plot'] = np.where(  ((non_gpl_df_roi['LOAN_PURPOSE'] == 'Plot purchase only' )| (non_gpl_df_roi['LOAN_PURPOSE'] == 'Plot loan BT' )) , 1,0 )
    non_gpl_df_roi['NON_GPL_plot'].value_counts()
    non_gpl_df_roi['NON_GPL_demi'] = np.where((non_gpl_df_roi['STEPFINANCE'] == 1) , 1,0 )
    non_gpl_df_roi['NON_GPL_demi'].value_counts()
    non_gpl_df_roi['NON_GPL_cre'] = np.where( (non_gpl_df_roi['CRE']=='CRE')|(non_gpl_df_roi['CRE']=='CRE-H') , 1, 0)
    non_gpl_df_roi['NON_GPL_cre'].value_counts()
    non_gpl_df_roi['NON_GPL_ahl'] = np.where( (non_gpl_df_roi['FINTYPE']=='AHL'),1,0)
    non_gpl_df_roi['NON_GPL_regular'] = np.where( (non_gpl_df_roi['NON_GPL_plot']==0) &  (non_gpl_df_roi['NON_GPL_cre'] ==0 )& 
                                      (non_gpl_df_roi['NON_GPL_demi'] ==0 ), 1,0 )
    non_gpl_df_roi['NON_GPL_regular'].value_counts()
    non_gpl_df_roi = non_gpl_df_roi[['BOOK_YEAR_MONTH', 'FINBRANCH', 'WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT', 'NON_GPL_regular','NON_GPL_demi','NON_GPL_cre','NON_GPL_plot','NON_GPL_ahl','NON_GPL_total','PRINCIPAL_OUTSTANDING','Net insurance income']]
    non_gpl_roi_group = non_gpl_df_roi.groupby(['BOOK_YEAR_MONTH', 'FINBRANCH' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    ngroi_group=non_gpl_df_roi.groupby(['FINBRANCH','NON_GPL_regular','NON_GPL_demi','NON_GPL_cre','NON_GPL_plot','NON_GPL_ahl','NON_GPL_total' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    non_gpl_roi_group['NON_GPL_ROI'] = non_gpl_roi_group.WROI / non_gpl_roi_group.BOOKED_AMOUNT
    non_gpl_roi_group['NON_GPL_PF'] = non_gpl_roi_group.PROCESSING_FEE / non_gpl_roi_group.BOOKED_AMOUNT
    non_gpl_roi_group['NON_GPL_NetP'] = non_gpl_roi_group.NET_PREMIUM / non_gpl_roi_group.BOOKED_AMOUNT
    
    #####################
    
    non_gpl_roi_group_total = non_gpl_df_roi.groupby([ 'BOOK_YEAR_MONTH','NON_GPL_regular','NON_GPL_demi','NON_GPL_cre','NON_GPL_plot','NON_GPL_ahl','NON_GPL_total' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
    non_gpl_roi_group_total['NON_GPL_ROI'] = non_gpl_roi_group_total.WROI / non_gpl_roi_group_total.BOOKED_AMOUNT
    non_gpl_roi_group_total['NON_GPL_PF'] = non_gpl_roi_group_total.PROCESSING_FEE / non_gpl_roi_group_total.BOOKED_AMOUNT
    non_gpl_roi_group_total['NON_GPL_NetP'] = non_gpl_roi_group_total.NET_PREMIUM / non_gpl_roi_group_total.BOOKED_AMOUNT
    non_gpl_roi_group.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    non_gpl_roi_group_total.drop(columns = ['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ],inplace = True)
    #non_gpl_roi_group_total.to_csv("non_gpl_roi_group_total.csv")
    ######################## NON GPL Completed #################################
    
    
    
    ############################## HL ##########################################
    
    HL_df=df[['FINTYPE','BOOK_YEAR_MONTH','STATUS','AHL_FLAG','FINBRANCH','WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income']] 
    
    HL_df_roi=HL_df.copy()
    HL_df_roi=HL_df_roi[HL_df_roi.FINTYPE=='AHL']
    HL_df_roi=HL_df_roi[HL_df_roi.STATUS=='Booked']
    HL_df_roi=HL_df_roi[['BOOK_YEAR_MONTH', 'FINBRANCH', 'WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','Net insurance income']]
    HL_df_roi_group=HL_df_roi.groupby(['BOOK_YEAR_MONTH', 'FINBRANCH' ],as_index=False).sum(['WROI','PROCESSING_FEE','NET_PREMIUM', 'BOOKED_AMOUNT' ])
  
    ######################## HL Completed #################################
    
    
    
    ########################## UE MTD #######################
    
    '''
    HL
    '''
    todays_date = date.today()
    MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
    mAHL_branch,mAHL=UE_Branch_HL(HL_df_roi,MTD)
    cal_c(mAHL_branch)
    drop_col_hl(mAHL_branch)
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_HL(BRANCHES)MTD.xlsx") as writer:
        mAHL_branch.to_excel(writer,sheet_name='AHL',)
    # HLF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS'],
    #      'AHL':[(mAHL.loc['Total','WROI']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','PROCESSING_FEE']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','NET_PREMIUM']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','Net insurance income']/mAHL.loc['Total','BOOKED_AMOUNT']),(sum(mAHL_branch.BOOKED_AMOUNT)),(sum(mAHL_branch.PRINCIPAL_OUTSTANDING))],
    #      }
    # MHLFORM=pd.DataFrame(data=HLF)
    '''
    GPL
    '''
    todays_date = date.today()
    MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
    
    #df_roi can be [gpl_df_roi,non_gpl_df_roi,lap_df_roi]
    #prod_name can be ['GPL_regular','GPL_demi','GPL_cre','GPL_plot','GPL_total'] for GPL

        
    
    prod_reg='GPL_regular'
    prod_demi='GPL_demi'
    prod_cre='GPL_cre'
    prod_plot='GPL_plot'
    prod_total='GPL_total'
        
    mgreg_branch,mgreg=UE_Branch_GPL(gpl_df_roi,MTD,prod_reg)   
    mgplt_branch,mgplt=UE_Branch_GPL(gpl_df_roi,MTD,prod_plot)
    mgdmi_branch,mgdmi=UE_Branch_GPL(gpl_df_roi,MTD,prod_demi)    
    mgcr_branch,mgcr=UE_Branch_GPL(gpl_df_roi,MTD,prod_cre)
    mgtot_branch,mgtot=UE_Branch_GPL(gpl_df_roi,MTD,prod_total) 
    
    
    
    mltgpl= [mgreg_branch,mgplt_branch,mgdmi_branch,mgcr_branch,mgtot_branch]   
    cal_c(mgreg_branch)    
    cal_c(mgplt_branch)    
    cal_c(mgdmi_branch)    
    cal_c(mgcr_branch) 
    cal_c(mgtot_branch)   
    
    
    for i in mltgpl:
        drop_col_gpl(i)
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_GPL(BRANCHES)(MTD).xlsx") as writer: 
        mgreg_branch.to_excel(writer,sheet_name='Regular',)
        mgplt_branch.to_excel(writer,sheet_name='Plot')
        mgdmi_branch.to_excel(writer,sheet_name='Demi')
        mgcr_branch.to_excel(writer,sheet_name='CRE')
        mgtot_branch.to_excel(writer,sheet_name='Total')
    
    GF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(mgreg.loc['Total','WROI']/mgreg.loc['Total','BOOKED_AMOUNT']),(100*mgreg.loc['Total','PROCESSING_FEE']/mgreg.loc['Total','BOOKED_AMOUNT']),(100*mgreg.loc['Total','NET_PREMIUM']/mgreg.loc['Total','BOOKED_AMOUNT']),(100*mgreg.loc['Total','Net insurance income']/mgreg.loc['Total','BOOKED_AMOUNT']),(sum(mgreg_branch.BOOKED_AMOUNT)),(sum(mgreg_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'DEMI':[(mgdmi.loc['Total','WROI']/mgdmi.loc['Total','BOOKED_AMOUNT']),(100*mgdmi.loc['Total','PROCESSING_FEE']/mgdmi.loc['Total','BOOKED_AMOUNT']),(100*mgdmi.loc['Total','NET_PREMIUM']/mgdmi.loc['Total','BOOKED_AMOUNT']),(100*mgdmi.loc['Total','Net insurance income']/mgdmi.loc['Total','BOOKED_AMOUNT']),(sum(mgdmi_branch.BOOKED_AMOUNT)),(sum(mgdmi_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'CRE':[(mgcr.loc['Total','WROI']/mgcr.loc['Total','BOOKED_AMOUNT']),(100*mgcr.loc['Total','PROCESSING_FEE']/mgcr.loc['Total','BOOKED_AMOUNT']),(100*mgcr.loc['Total','NET_PREMIUM']/mgcr.loc['Total','BOOKED_AMOUNT']),(100*mgcr.loc['Total','Net insurance income']/mgcr.loc['Total','BOOKED_AMOUNT']),(sum(mgcr_branch.BOOKED_AMOUNT)),(sum(mgcr_branch.PRINCIPAL_OUTSTANDING)),1],
         'PLOT':[(mgplt.loc['Total','WROI']/mgplt.loc['Total','BOOKED_AMOUNT']),(100*mgplt.loc['Total','PROCESSING_FEE']/mgplt.loc['Total','BOOKED_AMOUNT']),(100*mgplt.loc['Total','NET_PREMIUM']/mgplt.loc['Total','BOOKED_AMOUNT']),(100*mgplt.loc['Total','Net insurance income']/mgplt.loc['Total','BOOKED_AMOUNT']),(sum(mgplt_branch.BOOKED_AMOUNT)),(sum(mgplt_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOTAL':[(mgtot.loc['Total','WROI']/mgtot.loc['Total','BOOKED_AMOUNT']),(100*mgtot.loc['Total','PROCESSING_FEE']/mgtot.loc['Total','BOOKED_AMOUNT']),(100*mgtot.loc['Total','NET_PREMIUM']/mgtot.loc['Total','BOOKED_AMOUNT']),(100*mgtot.loc['Total','Net insurance income']/mgtot.loc['Total','BOOKED_AMOUNT']),(sum(mgtot_branch.BOOKED_AMOUNT)),(sum(mgtot_branch.PRINCIPAL_OUTSTANDING)),((0.25*(sum(mgreg_branch.BOOKED_AMOUNT))+0.25*(sum(mgdmi_branch.BOOKED_AMOUNT))+(sum(mgcr_branch.BOOKED_AMOUNT))+(sum(mgplt_branch.BOOKED_AMOUNT)))/((sum(mgreg_branch.BOOKED_AMOUNT))+(sum(mgdmi_branch.BOOKED_AMOUNT))+(sum(mgcr_branch.BOOKED_AMOUNT))+(sum(mgplt_branch.BOOKED_AMOUNT))))]}
    MGPLFORM=pd.DataFrame(data=GF)
    
    ##############################################################################################
    '''
    NON GPL
    
    '''
    
    
    
    prod_nreg='NON_GPL_regular'
    prod_ndemi='NON_GPL_demi'
    prod_ncre='NON_GPL_cre'
    prod_nplot='NON_GPL_plot'
    prod_nahl='NON_GPL_ahl'
    prod_ntotal='NON_GPL_total'
        
    mreg_branch,mreg=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_nreg)   
    mplot_branch,mplt=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_nplot)
    mdemi_branch,mdmi=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_ndemi)    
    mcre_branch,mcr=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_ncre)
    mahl_branch,mahl=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_nahl)
    mtotal_branch,mtot=UE_Branch_NGPL(non_gpl_df_roi,MTD,prod_ntotal)
    
    
    
    mlt= [mreg_branch,mplot_branch,mdemi_branch,mcre_branch,mahl_branch,mtotal_branch]   
    cal_c(mreg_branch)    
    cal_c(mplot_branch)    
    cal_c(mdemi_branch) 
    cal_c(mcre_branch) 
    cal_c(mahl_branch)
    cal_c(mtotal_branch)   
    
   
    
    for i in mlt:
        drop_col_ngpl(i)
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_NONGPL(BRANCHES)(MTD).xlsx") as writer:
        
        mreg_branch.to_excel(writer,sheet_name='Regular',)
        mplot_branch.to_excel(writer,sheet_name='Plot')
        mdemi_branch.to_excel(writer,sheet_name='Demi')
        mcre_branch.to_excel(writer,sheet_name='CRE')
        mtotal_branch.to_excel(writer,sheet_name='Total')
    
    NGF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(mreg.loc['Total','WROI']/mreg.loc['Total','BOOKED_AMOUNT']),(100*mreg.loc['Total','PROCESSING_FEE']/mreg.loc['Total','BOOKED_AMOUNT']),(100*mreg.loc['Total','NET_PREMIUM']/mreg.loc['Total','BOOKED_AMOUNT']),(100*mreg.loc['Total','Net insurance income']/mreg.loc['Total','BOOKED_AMOUNT']),(sum(mreg_branch.BOOKED_AMOUNT)),(sum(mreg_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'DEMI':[(mdmi.loc['Total','WROI']/mdmi.loc['Total','BOOKED_AMOUNT']),(100*mdmi.loc['Total','PROCESSING_FEE']/mdmi.loc['Total','BOOKED_AMOUNT']),(100*mdmi.loc['Total','NET_PREMIUM']/mdmi.loc['Total','BOOKED_AMOUNT']),(100*mdmi.loc['Total','Net insurance income']/mdmi.loc['Total','BOOKED_AMOUNT']),(sum(mdemi_branch.BOOKED_AMOUNT)),(sum(mdemi_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'CRE':[(mcr.loc['Total','WROI']/mcr.loc['Total','BOOKED_AMOUNT']),(100*mcr.loc['Total','PROCESSING_FEE']/mcr.loc['Total','BOOKED_AMOUNT']),(100*mcr.loc['Total','NET_PREMIUM']/mcr.loc['Total','BOOKED_AMOUNT']),(100*mcr.loc['Total','Net insurance income']/mcr.loc['Total','BOOKED_AMOUNT']),(sum(mcre_branch.BOOKED_AMOUNT)),(sum(mcre_branch.PRINCIPAL_OUTSTANDING)),1],
         'PLOT':[(mplt.loc['Total','WROI']/mplt.loc['Total','BOOKED_AMOUNT']),(100*mplt.loc['Total','PROCESSING_FEE']/mplt.loc['Total','BOOKED_AMOUNT']),(100*mplt.loc['Total','NET_PREMIUM']/mplt.loc['Total','BOOKED_AMOUNT']),(100*mplt.loc['Total','Net insurance income']/mplt.loc['Total','BOOKED_AMOUNT']),(sum(mplot_branch.BOOKED_AMOUNT)),(sum(mplot_branch.PRINCIPAL_OUTSTANDING)),1],
         'AHL':[(mAHL.loc['Total','WROI']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','PROCESSING_FEE']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','NET_PREMIUM']/mAHL.loc['Total','BOOKED_AMOUNT']),(100*mAHL.loc['Total','Net insurance income']/mAHL.loc['Total','BOOKED_AMOUNT']),(sum(mAHL_branch.BOOKED_AMOUNT)),(sum(mAHL_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'TOTAL':[(mtot.loc['Total','WROI']/mtot.loc['Total','BOOKED_AMOUNT']),(100*mtot.loc['Total','PROCESSING_FEE']/mtot.loc['Total','BOOKED_AMOUNT']),(100*mtot.loc['Total','NET_PREMIUM']/mtot.loc['Total','BOOKED_AMOUNT']),(100*mtot.loc['Total','Net insurance income']/mtot.loc['Total','BOOKED_AMOUNT']),(sum(mtotal_branch.BOOKED_AMOUNT)),(sum(mtotal_branch.PRINCIPAL_OUTSTANDING)),(0.25*(sum(mreg_branch.BOOKED_AMOUNT))+0.25*(sum(mdemi_branch.BOOKED_AMOUNT))+(sum(mcre_branch.BOOKED_AMOUNT)) +(sum(mplot_branch.BOOKED_AMOUNT)) +0.25*(sum(mAHL_branch.BOOKED_AMOUNT)))/((sum(mreg_branch.BOOKED_AMOUNT))+(sum(mdemi_branch.BOOKED_AMOUNT))+(sum(mcre_branch.BOOKED_AMOUNT))+(sum(mplot_branch.BOOKED_AMOUNT))+(sum(mAHL_branch.BOOKED_AMOUNT)))]}
    MNONGPLFORM=pd.DataFrame(data=NGF)
    
    ###################################################################################
    '''
    LAP 
    
    '''
    
    prod_lreg='Lap_regular'
    prod_lind='Lap_industrial'
    prod_lboost='Lap_booster'
    prod_ltopup='Lap_topup'
    prod_ltotal='Lap_total'
    prod_llrd='Lap_lrd'
        
    mlapreg_branch,ml_reg=UE_Branch_LAP(lap_df_roi,MTD,prod_lreg)   
    mlapind_branch,ml_ind=UE_Branch_LAP(lap_df_roi,MTD,prod_lind)
    mlapboost_branch,ml_boost=UE_Branch_LAP(lap_df_roi,MTD,prod_lboost)    
    mlaptopup_branch,ml_tup=UE_Branch_LAP(lap_df_roi,MTD,prod_ltopup)
    mlaptotal_branch,ml_tot=UE_Branch_LAP(lap_df_roi,MTD,prod_ltotal)
    mlaplrd_branch,ml_lrd=UE_Branch_LAP(lap_df_roi,MTD,prod_llrd)
    
    
    
    mlaplt= [mlapreg_branch,mlapind_branch,mlapboost_branch,mlaptopup_branch,mlaplrd_branch,mlaptotal_branch]   
    cal_c(mlapreg_branch)    
    cal_c(mlapind_branch)    
    cal_c(mlapboost_branch)    
    cal_c(mlaptopup_branch) 
    cal_c(mlaptotal_branch)   
    cal_c(mlaplrd_branch)
    
    for i in mlaplt:
        drop_col_lap(i)
        
        
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_LAP(BRANCHES)(MTD).xlsx") as writer:
        
        mlapreg_branch.to_excel(writer,sheet_name='Regular',)
        mlapind_branch.to_excel(writer,sheet_name='Industrial')
        mlapboost_branch.to_excel(writer,sheet_name='Booster')
        mlaptopup_branch.to_excel(writer,sheet_name='Topup')
        mlaplrd_branch.to_excel(writer,sheet_name='LRD')
        mlaptotal_branch.to_excel(writer,sheet_name='Total')
        
    
    LAPF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(ml_reg.loc['Total','WROI']/ml_reg.loc['Total','BOOKED_AMOUNT']),(100*ml_reg.loc['Total','PROCESSING_FEE']/ml_reg.loc['Total','BOOKED_AMOUNT']),(100*ml_reg.loc['Total','NET_PREMIUM']/ml_reg.loc['Total','BOOKED_AMOUNT']),(100*(ml_reg.loc['Total','Net insurance income']/ml_reg.loc['Total','BOOKED_AMOUNT'])),(sum(mlapreg_branch.BOOKED_AMOUNT)),(sum(mlapreg_branch.PRINCIPAL_OUTSTANDING)),1],
         'INDUSTRIAL':[(ml_ind.loc['Total','WROI']/ml_ind.loc['Total','BOOKED_AMOUNT']),(100*ml_ind.loc['Total','PROCESSING_FEE']/ml_ind.loc['Total','BOOKED_AMOUNT']),(100*ml_ind.loc['Total','NET_PREMIUM']/ml_ind.loc['Total','BOOKED_AMOUNT']),(100*ml_ind.loc['Total','Net insurance income']/ml_ind.loc['Total','BOOKED_AMOUNT']),(sum(mlapind_branch.BOOKED_AMOUNT)),(sum(mlapind_branch.PRINCIPAL_OUTSTANDING)),1],
         'BOOSTER':[(ml_boost.loc['Total','WROI']/ml_boost.loc['Total','BOOKED_AMOUNT']),(100*ml_boost.loc['Total','PROCESSING_FEE']/ml_boost.loc['Total','BOOKED_AMOUNT']),(100*ml_boost.loc['Total','NET_PREMIUM']/ml_boost.loc['Total','BOOKED_AMOUNT']),(100*ml_boost.loc['Total','Net insurance income']/ml_boost.loc['Total','BOOKED_AMOUNT']),(sum(mlapboost_branch.BOOKED_AMOUNT)),(sum(mlapboost_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOPUP':[(ml_tup.loc['Total','WROI']/ml_tup.loc['Total','BOOKED_AMOUNT']),(100*ml_tup.loc['Total','PROCESSING_FEE']/ml_tup.loc['Total','BOOKED_AMOUNT']),(100*ml_tup.loc['Total','NET_PREMIUM']/ml_tup.loc['Total','BOOKED_AMOUNT']),(100*ml_tup.loc['Total','Net insurance income']/ml_tup.loc['Total','BOOKED_AMOUNT']),(sum(mlaptopup_branch.BOOKED_AMOUNT)),(sum(mlaptopup_branch.PRINCIPAL_OUTSTANDING)),1],
         'LRD':[(ml_lrd.loc['Total','WROI']/ml_lrd.loc['Total','BOOKED_AMOUNT']),(100*ml_lrd.loc['Total','PROCESSING_FEE']/ml_lrd.loc['Total','BOOKED_AMOUNT']),(100*ml_lrd.loc['Total','NET_PREMIUM']/ml_lrd.loc['Total','BOOKED_AMOUNT']),(100*ml_lrd.loc['Total','Net insurance income']/ml_lrd.loc['Total','BOOKED_AMOUNT']),(sum(mlaptopup_branch.BOOKED_AMOUNT)),(sum(mlaptopup_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOTAL':[(ml_tot.loc['Total','WROI']/ml_tot.loc['Total','BOOKED_AMOUNT']),(100*ml_tot.loc['Total','PROCESSING_FEE']/ml_tot.loc['Total','BOOKED_AMOUNT']),(100*ml_tot.loc['Total','NET_PREMIUM']/ml_tot.loc['Total','BOOKED_AMOUNT']),(100*ml_tot.loc['Total','Net insurance income']/ml_tot.loc['Total','BOOKED_AMOUNT']),(sum(mlaptotal_branch.BOOKED_AMOUNT)),(sum(mlaptotal_branch.PRINCIPAL_OUTSTANDING)),1]}
    MLAPFORM=pd.DataFrame(data=LAPF)
    ####################################################################################
    
    
    ######CONSOLIDATING##################
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_ALL(MTD).xlsx") as writer:
        MLAPFORM.to_excel(writer, sheet_name="LAP",index=False)
        MGPLFORM.to_excel(writer, sheet_name="GPL",index=False)
        MNONGPLFORM.to_excel(writer, sheet_name="NGPL",index=False)
        # MHLFORM.to_excel(writer, sheet_name="HL",index=False)
    
        # a={}
        # for i in df1.columns:
    # mux = pd.MultiIndex.from_product([['GPL'], ['ROI','NET_INSURANCE_INCOME','PF']])      
    gpl_bdf={"MTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'GPL REGULAR':['','','','',''],
          'GPL DEMI':['','','','',''],
          'GPL CRE':['','','','',''],
          'GPL PLOT':['','','','',''],
          'GPL TOTAL':['','','','','']}      
    gpl_b=pd.DataFrame(data=gpl_bdf)
    
    
    Ngpl_bdf={"MTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'N-GPL REGULAR':['','','','',''],
          'N-GPL DEMI':['','','','',''],
          'N-GPL CRE':['','','','',''],
          'N-GPL PLOT':['','','','',''],
          'AHL':['','','','',''],
          'N-GPL TOTAL':['','','','','']}      
    Ngpl_b=pd.DataFrame(data=Ngpl_bdf)  
    
    lap_bdf={"MTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'LAP REGULAR':['','','','',''],
          'LAP INDUSTRIAL':['','','','',''],
          'LAP BOOSTER':['','','','',''],
          'LAP TOPUP':['','','','',''],
          'LAP LRD':['','','','',''],
          'LAP TOTAL':['','','','','']} 
    lap_b=pd.DataFrame(data=lap_bdf) 
    
    
    # hl_bdf={"MTD":['ROI','NET INSURANCE INCOME','PF'],
    #       'AHL':['','','']}            
    # hl_b=pd.DataFrame(data=hl_bdf)
    
    # TEMP1=pd.merge(hl_b, lap_b,left_on=['MTD'],right_on=['MTD'], how='outer')
    TEMP=pd.merge(gpl_b, Ngpl_b,left_on=['MTD'],right_on=['MTD'], how='outer')    
    final=pd.merge(TEMP, lap_b,left_on=['MTD'],right_on=['MTD'], how='outer')
    
    
    
    
        #     a[i]=final
            
        # return a
                   
    mdf1,mdf2,mdf3,mdf4,mdf5,mdf6,mdf7,mdf8,mdf9,mdf10,mdf11,mdf12,mdf13,mdf14,mdf15,mdf16,mdf17=branch_func(mgreg_branch, mgdmi_branch, mgcr_branch ,mgplt_branch, mgtot_branch,mreg_branch,mdemi_branch,mcre_branch,mplot_branch,mAHL_branch,mtotal_branch,mlapreg_branch,mlapind_branch,mlapboost_branch,mlaptopup_branch,mlaplrd_branch,mlaptotal_branch)    
    
    mdf1=mdf1.transpose()
    mdf1['RBI Provisions']=0.25
    mdf1=mdf1.transpose()
    mdf2=mdf2.transpose()
    mdf2['RBI Provisions']=0.25
    mdf2=mdf2.transpose()
    mdf3=mdf3.transpose()
    mdf3['RBI Provisions']=1
    mdf3=mdf3.transpose()
    mdf4=mdf4.transpose()
    mdf4['RBI Provisions']=1
    mdf4=mdf4.transpose()
    mdf5=mdf5.transpose()
    mdf5['RBI Provisions']=0
    mdf5=mdf5.transpose()
    mdf6=mdf6.transpose()
    mdf6['RBI Provisions']=0.25
    mdf6=mdf6.transpose()
    mdf7=mdf7.transpose()
    mdf7['RBI Provisions']=0.25
    mdf7=mdf7.transpose()
    mdf8=mdf8.transpose()
    mdf8['RBI Provisions']=1
    mdf8=mdf8.transpose()
    mdf9=mdf9.transpose()
    mdf9['RBI Provisions']=1
    mdf9=mdf9.transpose()
    mdf10=mdf10.transpose()
    mdf10['RBI Provisions']=0.25
    mdf10=mdf10.transpose()
    mdf11=mdf11.transpose()
    mdf11['RBI Provisions']=0
    mdf11=mdf11.transpose()
    mdf12=mdf12.transpose()
    mdf12['RBI Provisions']=1
    mdf12=mdf12.transpose()
    mdf13=mdf13.transpose()
    mdf13['RBI Provisions']=1
    mdf13=mdf13.transpose()
    mdf14=mdf14.transpose()
    mdf14['RBI Provisions']=1
    mdf14=mdf14.transpose()
    mdf15=mdf15.transpose()
    mdf15['RBI Provisions']=1
    mdf15=mdf15.transpose()
    mdf16=mdf16.transpose()
    mdf16['RBI Provisions']=1
    mdf16=mdf16.transpose()
    mdf17=mdf17.transpose()
    mdf17['RBI Provisions']=1
    mdf17=mdf17.transpose()
    mdf_br=[mdf1,mdf2,mdf3,mdf4,mdf5,mdf6,mdf7,mdf8,mdf9,mdf10,mdf11,mdf12,mdf13,mdf14,mdf15,mdf16,mdf17]
    branches=df['FINBRANCH'].unique()
    brnc_dc={}
    for i in branches:
        brnc_dc[i]=final.copy()
        brnc_dc[i].set_index('MTD',inplace=True)
        for j in brnc_dc[i].index:
            for z in range(len(brnc_dc[i].columns)):
                if i in (mdf_br[z].columns):
                    brnc_dc[i].loc[j,(brnc_dc[i].columns[z])]=(mdf_br[z]).loc[j,i]
                else:
                    brnc_dc[i].loc[j,brnc_dc[i].columns[z]]=0
        brnc_dc[i].replace(to_replace=[np.nan,'',' '],value=0,inplace=True)   
            
    m1=pd.merge(MGPLFORM,MNONGPLFORM,left_on=[''],right_on=[''], how='outer')
    # m2=pd.merge(MHLFORM,MLAPFORM,left_on=[''],right_on=[''], how='outer')
    MTOTAL=pd.merge(m1,MLAPFORM,left_on=[''],right_on=[''], how='outer')
    
      #####################################    YTD BM     ########################################
    
    
    todays_date = date.today()
    cur_month=todays_date.month-1
    YTD=[]
    while cur_month>=4:
        YTD.append(cur_month)
        cur_month=cur_month-1
    YTD
    ytd=[]
    for i in YTD:
        ytd.append(str(todays_date.year)+"-"+str(i))
    ytd
    type(MTD)
    '''
    HL YTD
    '''
    
    AHL_branch,AHL=UE_Branch_HL(HL_df_roi,ytd)
    cal_c(AHL_branch)
    drop_col_hl(AHL_branch)
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_HL(BRANCHES)(YTD).xlsx") as writer:
        AHL_branch.to_excel(writer,sheet_name='AHL',)
    # HLF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS'],
    #      'AHL':[(AHL.loc['Total','WROI']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','PROCESSING_FEE']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','NET_PREMIUM']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','Net insurance income']/AHL.loc['Total','BOOKED_AMOUNT']),(sum(AHL_branch.BOOKED_AMOUNT)),(sum(AHL_branch.PRINCIPAL_OUTSTANDING))],
    #      }
    # HLFORM=pd.DataFrame(data=HLF)
            
    ''' 
     
     GPL YTD
     
     
    '''
    
    
    greg_branch,greg=UE_Branch_GPL(gpl_df_roi,ytd,prod_reg)   
    gplt_branch,gplt=UE_Branch_GPL(gpl_df_roi,ytd,prod_plot)
    gdmi_branch,gdmi=UE_Branch_GPL(gpl_df_roi,ytd,prod_demi)    
    gcr_branch,gcr=UE_Branch_GPL(gpl_df_roi,ytd,prod_cre)
    gtot_branch,gtot=UE_Branch_GPL(gpl_df_roi,ytd,prod_total) 
    
    ltgpl= [greg_branch,gplt_branch,gdmi_branch,gcr_branch,gtot_branch]   
    cal_c(greg_branch)    
    cal_c(gplt_branch)    
    cal_c(gdmi_branch)    
    cal_c(gcr_branch) 
    cal_c(gtot_branch)   
    
    
    for i in ltgpl:
        drop_col_gpl(i)
    greg_branch.to_csv('gpl_regular.csv')
    gplt_branch.to_csv('gpl_plot.csv')
    gdmi_branch.to_csv('gpl_demi.csv')
    gcr_branch.to_csv('gpl_cre.csv')
    gtot_branch.to_csv('gpl_total.csv')
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_GPL(BRANCHES)(YTD).xlsx") as writer:
        greg_branch.to_excel(writer,sheet_name='Regular',)
        gplt_branch.to_excel(writer,sheet_name='Plot')
        gdmi_branch.to_excel(writer,sheet_name='Demi')
        gcr_branch.to_excel(writer,sheet_name='CRE')
        gtot_branch.to_excel(writer,sheet_name='Total')
    
        
    
    
    GF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(greg.loc['Total','WROI']/greg.loc['Total','BOOKED_AMOUNT']),(100*greg.loc['Total','PROCESSING_FEE']/greg.loc['Total','BOOKED_AMOUNT']),(100*greg.loc['Total','NET_PREMIUM']/greg.loc['Total','BOOKED_AMOUNT']),(100*greg.loc['Total','Net insurance income']/greg.loc['Total','BOOKED_AMOUNT']),(sum(greg_branch.BOOKED_AMOUNT)),(sum(greg_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'DEMI':[(gdmi.loc['Total','WROI']/gdmi.loc['Total','BOOKED_AMOUNT']),(100*gdmi.loc['Total','PROCESSING_FEE']/gdmi.loc['Total','BOOKED_AMOUNT']),(100*gdmi.loc['Total','NET_PREMIUM']/gdmi.loc['Total','BOOKED_AMOUNT']),(100*gdmi.loc['Total','Net insurance income']/gdmi.loc['Total','BOOKED_AMOUNT']),(sum(gdmi_branch.BOOKED_AMOUNT)),(sum(gdmi_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'CRE':[(gcr.loc['Total','WROI']/gcr.loc['Total','BOOKED_AMOUNT']),(100*gcr.loc['Total','PROCESSING_FEE']/gcr.loc['Total','BOOKED_AMOUNT']),(100*gcr.loc['Total','NET_PREMIUM']/gcr.loc['Total','BOOKED_AMOUNT']),(100*gcr.loc['Total','Net insurance income']/gcr.loc['Total','BOOKED_AMOUNT']),(sum(gcr_branch.BOOKED_AMOUNT)),(sum(gcr_branch.PRINCIPAL_OUTSTANDING)),1],
         'PLOT':[(gplt.loc['Total','WROI']/gplt.loc['Total','BOOKED_AMOUNT']),(100*gplt.loc['Total','PROCESSING_FEE']/gplt.loc['Total','BOOKED_AMOUNT']),(100*gplt.loc['Total','NET_PREMIUM']/gplt.loc['Total','BOOKED_AMOUNT']),(100*gplt.loc['Total','Net insurance income']/gplt.loc['Total','BOOKED_AMOUNT']),(sum(gplt_branch.BOOKED_AMOUNT)),(sum(gplt_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOTAL':[(gtot.loc['Total','WROI']/gtot.loc['Total','BOOKED_AMOUNT']),(100*gtot.loc['Total','PROCESSING_FEE']/gtot.loc['Total','BOOKED_AMOUNT']),(100*gtot.loc['Total','NET_PREMIUM']/gtot.loc['Total','BOOKED_AMOUNT']),(100*gtot.loc['Total','Net insurance income']/gtot.loc['Total','BOOKED_AMOUNT']),(sum(gtot_branch.BOOKED_AMOUNT)),(sum(gtot_branch.PRINCIPAL_OUTSTANDING)),((0.25*(sum(greg_branch.BOOKED_AMOUNT))+0.25*(sum(gdmi_branch.BOOKED_AMOUNT))+(sum(gcr_branch.BOOKED_AMOUNT))+(sum(gplt_branch.BOOKED_AMOUNT)))/((sum(greg_branch.BOOKED_AMOUNT))+(sum(gdmi_branch.BOOKED_AMOUNT))+(sum(gcr_branch.BOOKED_AMOUNT))+(sum(gplt_branch.BOOKED_AMOUNT))))]}
    GPLFORM=pd.DataFrame(data=GF)
    
    '''
    NON GPL YTD
    
    '''
    reg_branch,reg=UE_Branch_NGPL(non_gpl_df_roi,ytd,prod_nreg)   
    plot_branch,plt=UE_Branch_NGPL(non_gpl_df_roi,ytd,prod_nplot)
    demi_branch,dmi=UE_Branch_NGPL(non_gpl_df_roi,ytd,prod_ndemi)    
    cre_branch,cr=UE_Branch_NGPL(non_gpl_df_roi,ytd,prod_ncre)
    total_branch,tot=UE_Branch_NGPL(non_gpl_df_roi,ytd,prod_ntotal)
    
    
    lt= [reg_branch,plot_branch,demi_branch,cre_branch,total_branch]   
    cal_c(reg_branch)    
    cal_c(plot_branch)    
    cal_c(demi_branch)    
    cal_c(cre_branch) 
    cal_c(total_branch)   
    
    
    for i in lt:
        drop_col_ngpl(i)
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_NONGPL(BRANCHES)(YTD).xlsx") as writer:
        
        reg_branch.to_excel(writer,sheet_name='Regular',)
        plot_branch.to_excel(writer,sheet_name='Plot')
        demi_branch.to_excel(writer,sheet_name='Demi')
        cre_branch.to_excel(writer,sheet_name='CRE')
        total_branch.to_excel(writer,sheet_name='Total')
    
    NGF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(reg.loc['Total','WROI']/reg.loc['Total','BOOKED_AMOUNT']),(100*reg.loc['Total','PROCESSING_FEE']/reg.loc['Total','BOOKED_AMOUNT']),(100*reg.loc['Total','NET_PREMIUM']/reg.loc['Total','BOOKED_AMOUNT']),(100*reg.loc['Total','Net insurance income']/reg.loc['Total','BOOKED_AMOUNT']),(sum(reg_branch.BOOKED_AMOUNT)),(sum(reg_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'DEMI':[(dmi.loc['Total','WROI']/dmi.loc['Total','BOOKED_AMOUNT']),(100*dmi.loc['Total','PROCESSING_FEE']/dmi.loc['Total','BOOKED_AMOUNT']),(100*dmi.loc['Total','NET_PREMIUM']/dmi.loc['Total','BOOKED_AMOUNT']),(100*dmi.loc['Total','Net insurance income']/dmi.loc['Total','BOOKED_AMOUNT']),(sum(demi_branch.BOOKED_AMOUNT)),(sum(demi_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'CRE':[(cr.loc['Total','WROI']/cr.loc['Total','BOOKED_AMOUNT']),(100*cr.loc['Total','PROCESSING_FEE']/cr.loc['Total','BOOKED_AMOUNT']),(100*cr.loc['Total','NET_PREMIUM']/cr.loc['Total','BOOKED_AMOUNT']),(100*cr.loc['Total','Net insurance income']/cr.loc['Total','BOOKED_AMOUNT']),(sum(cre_branch.BOOKED_AMOUNT)),(sum(cre_branch.PRINCIPAL_OUTSTANDING)),1],
         'PLOT':[(plt.loc['Total','WROI']/plt.loc['Total','BOOKED_AMOUNT']),(100*plt.loc['Total','PROCESSING_FEE']/plt.loc['Total','BOOKED_AMOUNT']),(100*plt.loc['Total','NET_PREMIUM']/plt.loc['Total','BOOKED_AMOUNT']),(100*plt.loc['Total','Net insurance income']/plt.loc['Total','BOOKED_AMOUNT']),(sum(plot_branch.BOOKED_AMOUNT)),(sum(plot_branch.PRINCIPAL_OUTSTANDING)),1],
         'AHL':[(AHL.loc['Total','WROI']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','PROCESSING_FEE']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','NET_PREMIUM']/AHL.loc['Total','BOOKED_AMOUNT']),(100*AHL.loc['Total','Net insurance income']/AHL.loc['Total','BOOKED_AMOUNT']),(sum(AHL_branch.BOOKED_AMOUNT)),(sum(AHL_branch.PRINCIPAL_OUTSTANDING)),0.25],
         'TOTAL':[(tot.loc['Total','WROI']/tot.loc['Total','BOOKED_AMOUNT']),(100*tot.loc['Total','PROCESSING_FEE']/tot.loc['Total','BOOKED_AMOUNT']),(100*tot.loc['Total','NET_PREMIUM']/tot.loc['Total','BOOKED_AMOUNT']),(100*tot.loc['Total','Net insurance income']/tot.loc['Total','BOOKED_AMOUNT']),(sum(total_branch.BOOKED_AMOUNT)),(sum(total_branch.PRINCIPAL_OUTSTANDING)),(0.25*(sum(reg_branch.BOOKED_AMOUNT))+0.25*(sum(demi_branch.BOOKED_AMOUNT))+(sum(cre_branch.BOOKED_AMOUNT)) +(sum(plot_branch.BOOKED_AMOUNT)) +0.25*(sum(AHL_branch.BOOKED_AMOUNT)))/((sum(reg_branch.BOOKED_AMOUNT))+(sum(demi_branch.BOOKED_AMOUNT))+(sum(cre_branch.BOOKED_AMOUNT))+(sum(plot_branch.BOOKED_AMOUNT))+(sum(AHL_branch.BOOKED_AMOUNT)))]}
    NONGPLFORM=pd.DataFrame(data=NGF)
    
    
    '''
    LAP YTD
    
    '''
    
    
    lapreg_branch,l_reg=UE_Branch_LAP(lap_df_roi,ytd,prod_lreg)   
    lapind_branch,l_ind=UE_Branch_LAP(lap_df_roi,ytd,prod_lind)
    lapboost_branch,l_boost=UE_Branch_LAP(lap_df_roi,ytd,prod_lboost)    
    laptopup_branch,l_tup=UE_Branch_LAP(lap_df_roi,ytd,prod_ltopup)
    laplrd_branch,l_lrd=UE_Branch_LAP(lap_df_roi,ytd,prod_llrd)
    laptotal_branch,l_tot=UE_Branch_LAP(lap_df_roi,ytd,prod_ltotal)
    
    
    laplt= [lapreg_branch,lapind_branch,lapboost_branch,laptopup_branch,laplrd_branch,laptotal_branch]   
    cal_c(lapreg_branch)    
    cal_c(lapind_branch)    
    cal_c(lapboost_branch)    
    cal_c(laptopup_branch)
    cal_c(laplrd_branch)
    cal_c(laptotal_branch)   
    
    
    for i in laplt:
        drop_col_lap(i)
    
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_LAP(BRANCHES)(YTD).xlsx") as writer:
        
        lapreg_branch.to_excel(writer,sheet_name='Regular',)
        lapind_branch.to_excel(writer,sheet_name='Industrial')
        lapboost_branch.to_excel(writer,sheet_name='Booster')
        laptopup_branch.to_excel(writer,sheet_name='Topup')
        laplrd_branch.to_excel(writer,sheet_name='LRD')
        laptotal_branch.to_excel(writer,sheet_name='Total')
    # tot.loc['Total','WROI']
    # tot.loc['Total','BOOKED_AMOUNT']
    # tot.loc['Total','NET_PREMIUM']
    # tot.loc['Total','PROCESSING_FEE']
    LAPF={"":['ROI','PF','GROSS','NET_INSURANCE_INCOME','BOOKED_AMOUNT','POS','RBI Provisions'],
         'REGULAR':[(l_reg.loc['Total','WROI']/l_reg.loc['Total','BOOKED_AMOUNT']),(100*l_reg.loc['Total','PROCESSING_FEE']/l_reg.loc['Total','BOOKED_AMOUNT']),(100*l_reg.loc['Total','NET_PREMIUM']/l_reg.loc['Total','BOOKED_AMOUNT']),(100*(l_reg.loc['Total','Net insurance income']/l_reg.loc['Total','BOOKED_AMOUNT'])),(sum(lapreg_branch.BOOKED_AMOUNT)),(sum(lapreg_branch.PRINCIPAL_OUTSTANDING)),1],
         'INDUSTRIAL':[(l_ind.loc['Total','WROI']/l_ind.loc['Total','BOOKED_AMOUNT']),(100*l_ind.loc['Total','PROCESSING_FEE']/l_ind.loc['Total','BOOKED_AMOUNT']),(100*l_ind.loc['Total','NET_PREMIUM']/l_ind.loc['Total','BOOKED_AMOUNT']),(100*l_ind.loc['Total','Net insurance income']/l_ind.loc['Total','BOOKED_AMOUNT']),(sum(lapind_branch.BOOKED_AMOUNT)),(sum(lapind_branch.PRINCIPAL_OUTSTANDING)),1],
         'BOOSTER':[(l_boost.loc['Total','WROI']/l_boost.loc['Total','BOOKED_AMOUNT']),(100*l_boost.loc['Total','PROCESSING_FEE']/l_boost.loc['Total','BOOKED_AMOUNT']),(100*l_boost.loc['Total','NET_PREMIUM']/l_boost.loc['Total','BOOKED_AMOUNT']),(100*l_boost.loc['Total','Net insurance income']/l_boost.loc['Total','BOOKED_AMOUNT']),(sum(lapboost_branch.BOOKED_AMOUNT)),(sum(lapboost_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOPUP':[(l_tup.loc['Total','WROI']/l_tup.loc['Total','BOOKED_AMOUNT']),(100*l_tup.loc['Total','PROCESSING_FEE']/l_tup.loc['Total','BOOKED_AMOUNT']),(100*l_tup.loc['Total','NET_PREMIUM']/l_tup.loc['Total','BOOKED_AMOUNT']),(100*l_tup.loc['Total','Net insurance income']/l_tup.loc['Total','BOOKED_AMOUNT']),(sum(laptopup_branch.BOOKED_AMOUNT)),(sum(laptopup_branch.PRINCIPAL_OUTSTANDING)),1],
         'LRD':[(l_lrd.loc['Total','WROI']/l_lrd.loc['Total','BOOKED_AMOUNT']),(100*l_lrd.loc['Total','PROCESSING_FEE']/l_lrd.loc['Total','BOOKED_AMOUNT']),(100*l_lrd.loc['Total','NET_PREMIUM']/l_lrd.loc['Total','BOOKED_AMOUNT']),(100*l_lrd.loc['Total','Net insurance income']/l_lrd.loc['Total','BOOKED_AMOUNT']),(sum(laptopup_branch.BOOKED_AMOUNT)),(sum(laptopup_branch.PRINCIPAL_OUTSTANDING)),1],
         'TOTAL':[(l_tot.loc['Total','WROI']/l_tot.loc['Total','BOOKED_AMOUNT']),(100*l_tot.loc['Total','PROCESSING_FEE']/l_tot.loc['Total','BOOKED_AMOUNT']),(100*l_tot.loc['Total','NET_PREMIUM']/l_tot.loc['Total','BOOKED_AMOUNT']),(100*l_tot.loc['Total','Net insurance income']/l_tot.loc['Total','BOOKED_AMOUNT']),(sum(laptotal_branch.BOOKED_AMOUNT)),(sum(laptotal_branch.PRINCIPAL_OUTSTANDING)),1]}
    LAPFORM=pd.DataFrame(data=LAPF)
    # LAPFORM['LAP_PF']=100*LAPFORM['LAP_PF']
    # LAPFORM['LAP_GROSS']=100*LAPFORM['LAP_GROSS']
    # GPLFORM['LAP_PF']=100*GPLFORM['LAP_PF']
    # GPLFORM['LAP_GROSS']=100*GPLFORM['LAP_GROSS']
    # NONGPLFORM['LAP_PF']=100*NONGPLFORM['LAP_PF']
    # NONGPLFORM['LAP_GROSS']=100*NONGPLFORM['LAP_GROSS']
    #LAPFORM.to_excel('LAP_UNIT_ECONOMICS.csv',index=False)
    # from IPython.display import display_html 
    
    # df1 = GPLFORM
    # df2= NONGPLFORM
    # df3 = LAPFORM
    
    # df1_styler = df1.style.set_table_attributes("style='display:inline'").set_caption('GPL')
    # df2_styler = df2.style.set_table_attributes("style='display:inline'").set_caption('NON GPL')
    # df3_styler = df3.style.set_table_attributes("style='display:inline'").set_caption('LAP')
    
    # DF_FINAL=display_html(df1_styler._repr_html_()+df2_styler._repr_html_(), df2_styler._repr_html_(),raw=True)
    # DF_FINAL.to_excel(r'C:\Users\VAIBHAV.SRIVASTAV01\OPS_review_output\unit_economics.xlsx')
    with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\UNIT_ECONIMICS_ALL(YTD).xlsx") as writer:
        LAPFORM.to_excel(writer, sheet_name="LAP",index=False)
        GPLFORM.to_excel(writer, sheet_name="GPL",index=False)
        NONGPLFORM.to_excel(writer, sheet_name="NGPL",index=False)
        # HLFORM.to_excel(writer, sheet_name="HL",index=False)
    # def transpose(df):
    #     df=df[['ROI','NET INSURANCE INCOME','PF']]
    #     df=df.transpose()
    #     return df
           
    ym1=pd.merge(GPLFORM,NONGPLFORM,left_on=[''],right_on=[''], how='outer')
    # ym2=pd.merge(HLFORM,LAPFORM,left_on=[''],right_on=[''], how='outer')
    TOTAL=pd.merge(ym1,LAPFORM,left_on=[''],right_on=[''], how='outer')
    # def branch_func(df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15):
    #     df1=transpose(df1)
    #     df2=transpose(df2)
    #     df3=transpose(df3)
    #     df4=transpose(df4)
    #     df5=transpose(df5)
    #     df6=transpose(df6)
    #     df7=transpose(df7)
    #     df8=transpose(df8)
    #     df9=transpose(df9)
    #     df10=transpose(df10)
    #     df11=transpose(df11)
    #     df12=transpose(df12)
    #     df13=transpose(df13)
    #     df14=transpose(df14)
    #     df15=transpose(df15)
    #     return df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15
        
        # a={}
        # for i in df1.columns:
    # mux = pd.MultiIndex.from_product([['GPL'], ['ROI','NET_INSURANCE_INCOME','PF']])       
    gpl_bdf={"YTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'GPL REGULAR':['','','','',''],
          'GPL DEMI':['','','','',''],
          'GPL CRE':['','','','',''],
          'GPL PLOT':['','','','',''],
          'GPL TOTAL':['','','','','']}      
    gpl_b=pd.DataFrame(data=gpl_bdf)
    
    
    Ngpl_bdf={"YTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'N-GPL REGULAR':['','','','',''],
          'N-GPL DEMI':['','','','',''],
          'N-GPL CRE':['','','','',''],
          'N-GPL PLOT':['','','','',''],
          'AHL':['','','','',''],
          'N-GPL TOTAL':['','','','','']}      
    Ngpl_b=pd.DataFrame(data=Ngpl_bdf)  
    
    lap_bdf={"YTD":['ROI','NET INSURANCE INCOME','PF','RBI Provisions','BOOKED_AMOUNT'],
          'LAP REGULAR':['','','','',''],
          'LAP INDUSTRIAL':['','','','',''],
          'LAP BOOSTER':['','','','',''],
          'LAP TOPUP':['','','','',''],
          'LAP LRD':['','','','',''],
          'LAP TOTAL':['','','','','']} 
    lap_b=pd.DataFrame(data=lap_bdf)  
    # hl_bdf={"YTD":['ROI','NET INSURANCE INCOME','PF'],
    #       'AHL':['','','']}            
    # hl_b=pd.DataFrame(data=hl_bdf)
    
    # TEMP1=pd.merge(hl_b, lap_b,left_on=['YTD'],right_on=['YTD'], how='outer')
    TEMP=pd.merge(gpl_b, Ngpl_b,left_on=['YTD'],right_on=['YTD'], how='outer')    
    ytd_final=pd.merge(TEMP,lap_b ,left_on=['YTD'],right_on=['YTD'], how='outer')
    
    
    
        #     a[i]=final
            
        # return a
                   
    df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15,df16,df17=branch_func(greg_branch, gdmi_branch, gcr_branch ,gplt_branch, gtot_branch,reg_branch,demi_branch,cre_branch,plot_branch,AHL_branch,total_branch,lapreg_branch,lapind_branch,lapboost_branch,laptopup_branch,laplrd_branch,laptotal_branch)    
    
    df1=df1.transpose()
    df1['RBI Provisions']=0.25
    df1=df1.transpose()
    df2=df2.transpose()
    df2['RBI Provisions']=0.25
    df2=df2.transpose()
    df3=df3.transpose()
    df3['RBI Provisions']=1
    df3=df3.transpose()
    df4=df4.transpose()
    df4['RBI Provisions']=1
    df4=df4.transpose()
    df5=df5.transpose()
    df5['RBI Provisions']=0
    df5=df5.transpose()
    df6=df6.transpose()
    df6['RBI Provisions']=0.25
    df6=df6.transpose()
    df7=df7.transpose()
    df7['RBI Provisions']=0.25
    df7=df7.transpose()
    df8=df8.transpose()
    df8['RBI Provisions']=1
    df8=df8.transpose()
    df9=df9.transpose()
    df9['RBI Provisions']=1
    df9=df9.transpose()
    df10=df10.transpose()
    df10['RBI Provisions']=0.25
    df10=df10.transpose()
    df11=df11.transpose()
    df11['RBI Provisions']=0
    df11=df11.transpose()
    df12=df12.transpose()
    df12['RBI Provisions']=1
    df12=df12.transpose()
    df13=df13.transpose()
    df13['RBI Provisions']=1
    df13=df13.transpose()
    df14=df14.transpose()
    df14['RBI Provisions']=1
    df14=df14.transpose()
    df15=df15.transpose()
    df15['RBI Provisions']=1
    df15=df15.transpose()
    df16=df16.transpose()
    df16['RBI Provisions']=1
    df16=df16.transpose()
    df17=df17.transpose()
    df17['RBI Provisions']=1
    df17=df17.transpose()
    df_br=[df1,df2,df3,df4,df5,df6,df7,df8,df9,df10,df11,df12,df13,df14,df15,df16,df17]
    branches=df['FINBRANCH'].unique()
    ytd_brnc_dc={}
    for i in branches:
        ytd_brnc_dc[i]=ytd_final.copy()
        ytd_brnc_dc[i].set_index('YTD',inplace=True)
        for j in ytd_brnc_dc[i].index:
            for z in range(len(ytd_brnc_dc[i].columns)):
                if i in (df_br[z].columns):
                    ytd_brnc_dc[i].loc[j,(brnc_dc[i].columns[z])]=(df_br[z]).loc[j,i]
                else:
                    ytd_brnc_dc[i].loc[j,brnc_dc[i].columns[z]]=0
        ytd_brnc_dc[i].replace(to_replace=[np.nan,'',' '],value=0,inplace=True)   
                
    writer = pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\BRANCH_UE.xlsx",engine='xlsxwriter')
    workbook=writer.book
    for x in branches:
        worksheet=workbook.add_worksheet(x)
        writer.sheets[x] = worksheet
        worksheet.write_string(0, 0, "MTD")
        
        brnc_dc[x].to_excel(writer,sheet_name=x,startrow=1 , startcol=0)
        worksheet.write_string(brnc_dc[x].shape[0] + 4, 0, 'YTD')
        ytd_brnc_dc[x].to_excel(writer,sheet_name=x,startrow=brnc_dc[x].shape[0] + 5, startcol=0)
    writer.save()
    return brnc_dc,ytd_brnc_dc,MTOTAL,TOTAL
    
mtd_dc,ytd_dc,MTOTAL,TOTAL=Unit_Economics(mydata)
for i in mtd_dc.keys():
    mtd_dc[i].loc['RBI Provisions','GPL TOTAL']=(0.25*mtd_dc[i].loc['BOOKED_AMOUNT','GPL REGULAR']+0.25*mtd_dc[i].loc['BOOKED_AMOUNT','GPL DEMI']+mtd_dc[i].loc['BOOKED_AMOUNT','GPL CRE']+mtd_dc[i].loc['BOOKED_AMOUNT','GPL PLOT'])/(mtd_dc[i].loc['BOOKED_AMOUNT','GPL REGULAR']+mtd_dc[i].loc['BOOKED_AMOUNT','GPL DEMI']+mtd_dc[i].loc['BOOKED_AMOUNT','GPL CRE']+mtd_dc[i].loc['BOOKED_AMOUNT','GPL PLOT'])
    mtd_dc[i].loc['RBI Provisions','N-GPL TOTAL']=(0.25*mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL REGULAR']+0.25*mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL DEMI']+mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL CRE']+mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL PLOT']+0.25*mtd_dc[i].loc['BOOKED_AMOUNT','AHL'])/(mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL REGULAR']+mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL DEMI']+mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL CRE']+mtd_dc[i].loc['BOOKED_AMOUNT','N-GPL PLOT']+mtd_dc[i].loc['BOOKED_AMOUNT','AHL'])
    mtd_dc[i].loc['RBI Provisions','LAP TOTAL']=1
    ytd_dc[i].loc['RBI Provisions','GPL TOTAL']=(0.25*ytd_dc[i].loc['BOOKED_AMOUNT','GPL REGULAR']+0.25*ytd_dc[i].loc['BOOKED_AMOUNT','GPL DEMI']+ytd_dc[i].loc['BOOKED_AMOUNT','GPL CRE']+ytd_dc[i].loc['BOOKED_AMOUNT','GPL PLOT'])/(ytd_dc[i].loc['BOOKED_AMOUNT','GPL REGULAR']+ytd_dc[i].loc['BOOKED_AMOUNT','GPL DEMI']+ytd_dc[i].loc['BOOKED_AMOUNT','GPL CRE']+ytd_dc[i].loc['BOOKED_AMOUNT','GPL PLOT'])
    ytd_dc[i].loc['RBI Provisions','N-GPL TOTAL']=(0.25*ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL REGULAR']+0.25*ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL DEMI']+ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL CRE']+ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL PLOT']+0.25*ytd_dc[i].loc['BOOKED_AMOUNT','AHL'])/(ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL REGULAR']+ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL DEMI']+ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL CRE']+ytd_dc[i].loc['BOOKED_AMOUNT','N-GPL PLOT']+ytd_dc[i].loc['BOOKED_AMOUNT','AHL'])
    ytd_dc[i].loc['RBI Provisions','LAP TOTAL']=1
    ytd_dc[i]=ytd_dc[i].transpose()
    ytd_dc[i].drop(columns='BOOKED_AMOUNT', inplace=True)    
    ytd_dc[i]=ytd_dc[i].transpose()
    mtd_dc[i]=mtd_dc[i].transpose()
    mtd_dc[i].drop(columns='BOOKED_AMOUNT', inplace=True)    
    mtd_dc[i]=mtd_dc[i].transpose()




MTOTAL.set_index('',inplace=True)
MTOTAL=MTOTAL.transpose()
MTOTAL.rename(columns = {'NET_INSURANCE_INCOME':'NET INSURANCE INCOME'}, inplace = True)

MTOTAL=MTOTAL[['ROI','NET INSURANCE INCOME','PF','RBI Provisions']]
MTOTAL=MTOTAL.transpose()
TOTAL.set_index('',inplace=True)
TOTAL=TOTAL.transpose()
TOTAL.rename(columns = {'NET_INSURANCE_INCOME':'NET INSURANCE INCOME'}, inplace = True)
TOTAL=TOTAL[['ROI','NET INSURANCE INCOME','PF','RBI Provisions']]
TOTAL=TOTAL.transpose()
mtd_dc['Total']=MTOTAL
ytd_dc['Total']=TOTAL
for i in mtd_dc.keys():
    mtd_dc[i]=mtd_dc[i].round(decimals=2)
    ytd_dc[i]=ytd_dc[i].round(decimals=2)
colum=['Regular',
       'Demi',
       'Cre',
 'Plot',
 'Total',
 'Regular',
 'Demi',
 'Cre',
 'Plot',
 'Ahl',
 'Total',
 'Regular',
 'Industrial',
 'Booster',
 'Topup',
 'Lrd',
 'Total']

for i in mtd_dc.keys():
    
    mtd_dc[i]=mtd_dc[i].transpose()
    mtd_dc[i]['Loss on Morat']=" "
    mtd_dc[i]['COF']=" "
    mtd_dc[i]['Spread']=" "
    mtd_dc[i]['Gross NIM']=" "
    mtd_dc[i]['COA']=" "
    mtd_dc[i]['Net NIM']=" "
    mtd_dc[i]['MTD']=' '
    mtd_dc[i].reset_index(level=0, inplace=True)
    mtd_dc[i].rename(columns = {'index':' '}, inplace = True)
    mtd_dc[i]=mtd_dc[i][['MTD',' ','ROI','Loss on Morat','COF','Spread','NET INSURANCE INCOME','PF','Gross NIM','COA','RBI Provisions','Net NIM']]
    mtd_dc[i]['NET INSURANCE INCOME']=' '
    mtd_dc[i]['MTD']= mtd_dc[i]['MTD'].str.upper().str.title()
    mtd_dc[i]=mtd_dc[i].transpose()
    for j in range(len(mtd_dc[i].columns)):
        (mtd_dc[i].loc[' ',j])=(colum[j])
    mtd_dc[i].reset_index(level=0, inplace=True)
    mtd_dc[i].insert(0,calendar.month_name[todays_date.month][:3],' ')
    mtd_dc[i].iloc[0,0]='MTD '+str(calendar.month_name[todays_date.month-1][:3])
    

for i in ytd_dc.keys():
    
    ytd_dc[i]=ytd_dc[i].transpose()
    ytd_dc[i]['Loss on Morat']=" "
    ytd_dc[i]['COF']=" "
    ytd_dc[i]['Spread']=" "
    ytd_dc[i]['Gross NIM']=" "
    ytd_dc[i]['COA']=" "
    ytd_dc[i]['Net NIM']=" "
    ytd_dc[i]['YTD']=' '
    ytd_dc[i].reset_index(level=0, inplace=True)
    ytd_dc[i].rename(columns = {'index':' '}, inplace = True)
    ytd_dc[i]=ytd_dc[i][['YTD',' ','ROI','Loss on Morat','COF','Spread','NET INSURANCE INCOME','PF','Gross NIM','COA','RBI Provisions','Net NIM']]
    ytd_dc[i]['NET INSURANCE INCOME']=' '
    ytd_dc[i]['YTD']= ytd_dc[i]['YTD'].str.upper().str.title()
    ytd_dc[i]=ytd_dc[i].transpose()
    for j in range(len(ytd_dc[i].columns)):
        (ytd_dc[i].loc[' ',j])=(colum[j])
    ytd_dc[i].reset_index(level=0, inplace=True)
    ytd_dc[i].insert(0,'YTD 22',' ')
    ytd_dc[i].iloc[0,0]='YTD 22'
MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
prs = Presentation(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\OPERATING REVIEW QUARTER_"+str(get_quarter(int(MTD[5]))[1])+".pptx")

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

todays_date = date.today()
def UE_PPT(title,mtd,ytd,prs,c):
    prs.slides[c].shapes[0].text=title
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    # slide.shapes.title.text = title
    # slide.shapes.title.top = Inches(0)
    # slide.shapes.title.left = Inches(0)
    # slide.shapes.title.width = Inches(10)
    # slide.shapes.title.height = Inches(0.3)
    # slide.shapes.title.text_frame.paragraphs[0].font.size=Pt(12)
    # slide.shapes.title.text_frame.paragraphs[0].font.underline = True
    
# ---add table to slide---
    # x, y, cx, cy = Inches(0.5), Inches(0.25), Inches(11), Inches(2)
    # shape1 = slide.shapes.add_table(len(mtd),len(mtd.columns),x, y, cx,cy)
    
    # tbl1 =  shape1._element.graphic.graphicData.tbl
    # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    # tbl1[0][-1].text = style_id    
    # table1 = shape1.table

    # x, y, cx, cy = Inches(0.5), Inches(3.75), Inches(11), Inches(2)
    # shape2 = slide.shapes.add_table(len(ytd),len(ytd.columns),x, y, cx,cy)
    
    # tbl2 =  shape2._element.graphic.graphicData.tbl
    # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    # tbl2[0][-1].text = style_id    
    # table2 = shape2.table
    
    # cols1=table1.columns
    # cols1[1].width=Inches(1.75)
    
    # cols2=table2.columns
    # cols2[1].width=Inches(1.75)
    # for i in range(len(cols1)):
    #     if i==0:
    #         continue
    #     # else:
    #     #     cols1[i].width=Inches(0.58)
    #     #     cols2[i].width=Inches(0.58)
    for i in range(len(mtd)):
        for j in range(len(mtd.columns)):
            prs.slides[c].shapes[3].table.cell(i,j).text=str(mtd.iloc[i,j])
            # prs.slides[c].shapes[5].table.cell(i,j).text=str(0)
            # if (i==0 & ('N-GPL ' in str(mtd.iloc[i,j]))& j!=0):
            #     print(mtd.iloc[i,j])
            #     prs.slides[c].shapes[3].table.cell(i,j).text=str(mtd.iloc[i,j]).split('N-GPL ')[1]
            # elif (i==0 & ('GPL ' in str(mtd.iloc[i,j]))& j!=0):
            #     table1.cell(i,j).text=str(mtd.iloc[i,j]).split('GPL ')[1]
            # elif (i==0 & ('LAP ' in str(mtd.iloc[i,j]))& j!=0) :
            #     table1.cell(i,j).text=str(mtd.iloc[i,j]).split('LAP ')[1]
            # else:
            #     table1.cell(i,j).text=str(mtd.iloc[i,j])
                    
                # table.cell(i,j).text_frame.margin_bottom = Inches(0.08)
                # table.cell(i,j).text_frame.margin_left = 0
                # table.cell(i,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
                # table.cell(i,j).text_frame.word_wrap = False
                # table.cell(i,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                # table.cell(i,j).text_frame.fit_text()
                # if i==0:
                #     cell=table1.cell(i, j)
                #     fill = cell.fill
                #     fill.solid()
                #     fill.fore_color.rgb = RGBColor(173, 216, 230)
    for i in range(len(ytd)):
        for j in range(len(ytd.columns)):
            prs.slides[c].shapes[5].table.cell(i,j).text=str(ytd.iloc[i,j])
            # prs.slides[c].shapes[5].table.cell(i,j).text=str(0)
            # if (i==0 & ('N-GPL ' in str(ytd.iloc[i,j]))& j!=0):
            #     table2.cell(i,j).text=str(ytd.iloc[i,j]).text.split('N-GPL ')[1]
            # elif (i==0 & ('GPL ' in str(ytd.iloc[i,j]))& j!=0):
            #     table2.cell(i,j).text=str(ytd.iloc[i,j]).text.split('GPL ')[1]
            # elif (i==0 & ('LAP ' in str(ytd.iloc[i,j]))& j!=0):
            #     table2.cell(i,j).text=str(ytd.iloc[i,j]).text.split('LAP ')[1]
            # else:
            #     table2.cell(i,j).text=str(ytd.iloc[i,j])
                    
                # table.cell(i,j).text_frame.margin_bottom = Inches(0.08)
                # table.cell(i,j).text_frame.margin_left = 0
                # table.cell(i,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
                # table.cell(i,j).text_frame.word_wrap = False
                # table.cell(i,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                # table.cell(i,j).text_frame.fit_text()
                # if i==0:
                #     cell=table2.cell(i, j)
                #     fill = cell.fill
                #     fill.solid()
                #     fill.fore_color.rgb = RGBColor(173, 216, 230)
    
    prs.slides[c].shapes[3].table.cell(0,2).text='GPL'
    
    prs.slides[c].shapes[3].table.cell(0,7).text='NGPL'
    
    prs.slides[c].shapes[3].table.cell(0,12).text='HL'
    prs.slides[c].shapes[3].table.cell(0,13).text='LAP'
    
    prs.slides[c].shapes[5].table.cell(0,2).text='GPL'
    
    prs.slides[c].shapes[5].table.cell(0,7).text='NGPL'
    
    prs.slides[c].shapes[5].table.cell(0,12).text='HL'
    prs.slides[c].shapes[5].table.cell(0,13).text='LAP'
    
    
    prs.slides[c].shapes[3].table.cell(4,0).text='Less'
    prs.slides[c].shapes[5].table.cell(4,0).text='Less'
    prs.slides[c].shapes[3].table.cell(6,0).text='Add'
    prs.slides[c].shapes[5].table.cell(6,0).text='Add'
    prs.slides[c].shapes[3].table.cell(7,0).text='Add'
    prs.slides[c].shapes[5].table.cell(7,0).text='Add'
    prs.slides[c].shapes[3].table.cell(9,0).text='Less'
    prs.slides[c].shapes[5].table.cell(9,0).text='Less'
    prs.slides[c].shapes[3].table.cell(10,0).text='Less'
    prs.slides[c].shapes[5].table.cell(10,0).text='Less'
    cell = prs.slides[c].shapes[5].table.rows[0].cells[0]
    
    for cell in iter_cells(prs.slides[c].shapes[3].table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0,0,0)
                run.font.style= 'Calibri'
                run.font.bold=False
    for cell in iter_cells(prs.slides[c].shapes[5].table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0,0,0)
                run.font.style= 'Calibri'
                run.font.bold=False
    # for i in range(len(ytd)):
    #     for j in range(len(ytd.columns)):
    #         if i<2 & j<2:
    #             cell = prs.slides[c].shapes[5].table.cell(i,j)
    #             paragraphs = cell.text_frame.paragraphs
    #             for paragraph in paragraphs:
    #                 paragraph.font.size = Pt(13)
    #                 paragraph.font.bold=True
    #             cell = prs.slides[c].shapes[3].table.cell(i,j)
    #             paragraphs = cell.text_frame.paragraphs
    #             for paragraph in paragraphs:
    #                 paragraph.font.size = Pt(13)
    #                 paragraph.font.bold=True 
    return(prs)
count=1
for i in mtd_dc.keys():
    UE_PPT(("Unit Economics:"+str(i)),mtd_dc[i],ytd_dc[i],prs,count)
    count=count+1
    
prs.save(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\Q2\OPERATING REVIEW QUARTER_"+str(get_quarter(int(MTD[5]))[1])+".pptx")
#mydata.to_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\DTCRON_"+str(todays_date)+".xlsx")

