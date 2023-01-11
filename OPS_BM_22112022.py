# -*- coding: utf-8 -*-
"""
Created on Tue Nov 22 11:19:31 2022

@author: VAIBHAV.SRIVASTAV01
"""

####### OPS TRANSITION BM




import pandas as pd
import os
from datetime import date
import calendar
import numpy as np
from pptx import Presentation
import pyodbc
pyodbc.drivers()
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import sys
import collections


con_dm = pyodbc.connect('DSN=GHF_BI_CONN_DM;UID=GHF_BI_CONN;PWD=Godrej@123')
con = pyodbc.connect('DSN=GHF_BI_CONN;UID=GHF_BI_CONN;PWD=Godrej@123')
con_dm_nbfc = pyodbc.connect('DSN=GHF_BI_CONN_DM_NBFC;UID=GHF_BI_CONN;PWD=Godrej@123')

# dtcron=pd.read_sql("SELECT * from dtcron;", con)
ghf_loan_view=pd.read_sql("SELECT * from prod_da_db.serve.loan_view WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_loan_view['NBFC_FLAG'] = 'N'
ghf_base_view = pd.read_sql("SELECT * from prod_da_db.serve.BASE_VIEW WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_base_view['NBFC_FLAG'] = 'N'
ghf_customer_base = pd.read_sql("SELECT * from prod_da_db.serve.customer_base WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_customer_base['NBFC_FLAG'] = 'N'
ghf_insurance_view = pd.read_sql("SELECT * from prod_da_db.serve.INSURANCE_VIEW WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_insurance_view ['NBFC_FLAG'] = 'N'
ghf_lanvas=pd.read_sql("SELECT * from prod_da_db.serve.x_ref_lan_to_vas WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_lanvas['NBFC_FLAG'] = 'N'
ghf_lancif=pd.read_sql("SELECT * from prod_da_db.serve.x_ref_lan_to_cif WHERE DH_RECORD_ACTIVE_FLAG = 'Y' and APPLICANT_TYPE='APPLICANT'", con_dm)
ghf_lancif['NBFC_FLAG'] = 'N'
ghf_cifvas=pd.read_sql("SELECT * from prod_da_db.serve.x_ref_cif_to_vas WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_cifvas['NBFC_FLAG'] = 'N'
ghf_lancollat=pd.read_sql("SELECT * from prod_da_db.serve.x_ref_lan_to_collat WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm)
ghf_lancollat['NBFC_FLAG'] = 'N'

gfl_loan_view=pd.read_sql("SELECT * from prod_gfl_da_db.serve.loan_view WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_loan_view['NBFC_FLAG'] = 'Y'
gfl_base_view = pd.read_sql("SELECT * from prod_gfl_da_db.serve.BASE_VIEW WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_base_view['NBFC_FLAG'] = 'Y'
gfl_customer_base = pd.read_sql("SELECT * from prod_gfl_da_db.serve.customer_base WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_customer_base['NBFC_FLAG'] = 'Y'
gfl_insurance_view = pd.read_sql("SELECT * from prod_gfl_da_db.serve.INSURANCE_VIEW WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_insurance_view['NBFC_FLAG'] = 'Y'
gfl_lanvas=pd.read_sql("SELECT * from prod_gfl_da_db.serve.x_ref_lan_to_vas WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_lanvas['NBFC_FLAG'] = 'Y'
gfl_lancif=pd.read_sql("SELECT * from prod_gfl_da_db.serve.x_ref_lan_to_cif WHERE DH_RECORD_ACTIVE_FLAG = 'Y' and APPLICANT_TYPE='APPLICANT'", con_dm_nbfc)
gfl_lancif['NBFC_FLAG'] = 'Y'
gfl_cifvas=pd.read_sql("SELECT * from prod_gfl_da_db.serve.x_ref_cif_to_vas WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_cifvas['NBFC_FLAG'] = 'Y'
gfl_lancollat=pd.read_sql("SELECT * from prod_gfl_da_db.serve.x_ref_lan_to_collat WHERE DH_RECORD_ACTIVE_FLAG = 'Y'", con_dm_nbfc)
gfl_lancollat['NBFC_FLAG'] = 'Y'

loan_view=pd.concat([ghf_loan_view,gfl_loan_view]) #REFERENCE
base_view=pd.concat([ghf_base_view,gfl_base_view]) # LAN_ID
ghf_lancif=ghf_lancif[[ 'CUSTOMER_CIF','LAN_ID','APPLICANT_TYPE']]
ghf_customer_base=ghf_lancif.merge(ghf_customer_base,on = 'CUSTOMER_CIF',how='left')
gfl_lancif=gfl_lancif[[ 'CUSTOMER_CIF','LAN_ID','APPLICANT_TYPE']]
gfl_customer_base=gfl_lancif.merge(gfl_customer_base,on = 'CUSTOMER_CIF',how='left')

customer_view=pd.concat([ghf_customer_base,gfl_customer_base]) # CIF
insurance_view=pd.concat([ghf_insurance_view,gfl_insurance_view]) # LAN_ID
lanvas=pd.concat([ghf_lanvas,gfl_lanvas])
lancif=pd.concat([ghf_lancif,gfl_lancif])
cifvas=pd.concat([ghf_cifvas,gfl_cifvas])
lancollat=pd.concat([ghf_lancollat,gfl_lancollat])
loan_view.rename(columns={'REFERENCE':'LAN_ID'},inplace=True)

def error_lans():
    a=['GHF1001FL0000328',
 'GHF1001FL0002120',
 'GHF1002FL0002436',
 'GHF1001FL0002119',
 'GHF1002FL0000299',
 'GHF1002FL0000640',
 'GHF1002HL0002943',
 'GHF1002HL0000836',
 'GHF1002FT0001610',
 'GHF1002FL0001993',
 'GHF1002HL0002773',
 'GHF1001FT0000552',
 'GHF1101LP0003918',
 'GHF1401HL0000123',
 'GHF1002HL0003082',
 'GHF1002FL0000605',
 'GHF1002HL0002977',
 'GHF1001FL0000327',
 'GHF1002HL0003004',
 'GHF1002FL0005334',
 'GHF1001FL0000912',
 'GHF1002FL0000257',
 'GHF1002HL0003367',
 'GHF1002FL0000744',
 'GHF1002FT0001009',
 'GHF1001FL0000329',
 'GHF1401HL0000122',
 'GHF1001HL0000342',
 'GHF1002HL0001014',
 'GHF1002HL0003484',
 'GHF1401HL0001902',
 'GHF1001HL0000900',
 'GHF1002FL0000063',
 'GHF1001FL0002118',
 'GHF1401HL0000124',
 'GHF1002FL0002252',
 'GHF1401HL0000125',
 'GHF1002FL0003178',
 'GHF1002FL0000659',
 'GHF1001HL0000967',
 'GHF1001HL0000966']
    return a
def remove_error_lans(df):
    try:
        df=df[df['REFERENCE'].isin(error_lans())==False]
    except:
        pass
    try:
        df=df[df['LAN_ID'].isin(error_lans())==False]
    except:
        pass
    try:
        df=df[df['FINREFERENCE'].isin(error_lans())==False]
    except:
        pass
    try:
        df=df[df['FINANCE_REFERENCE'].isin(error_lans())==False]
    except:
        pass
    return df
def find_columns(base_view,dtcron):
    
    base_view=remove_error_lans(base_view)
    dtcron_dict={}
    dt=dtcron.columns.to_list()
    dt.sort()
    for i in dt:
        dtcron_dict[i]=list(dtcron[i].unique())
        
    bv_dict={}
    bv=base_view.columns.to_list()
    bv.sort()
    for i in bv:
        bv_dict[i]=list(base_view[i].unique())
    a=[]
    for j in bv: 
        for i in dt:        
            if len(bv_dict[j])>2 and (collections.Counter(bv_dict[j]) == collections.Counter(dtcron_dict[i]) or set(bv_dict[j]).issubset(set(dtcron_dict[i])) or j==i) :
                a.append(j)
                print(j,i)
                break
    return a         
          
# bv_cols=find_columns(base_view,dtcron)  
# b=['BOOKING_AMOUNT',
#  'EOMCLTRL',
#  'EOMLOGN',
#  'EOMRJCT',
#  'EOMSNCTN',
#  'REQUESTED_AMOUNT',
#  'LOGIN_STATUS']      
# c=list(set(bv_cols)-set(b))
bv_cols=['FIRST_REJECT_DATE',
 'DETAILED_STATUS',
 'SANCTION_AMOUNT',
 'REPORTING_BRANCH',
 'LAN_ID',
 'QUEUE',
 'LOAN_STATUS',
 'STATUS',
 'STATUS_SEG',
 'FIRST_ENTRY_TO_COLLATERAL',
 'FINANCE_TYPE',
 'LOAN_TYPE',
 'BOOKING_DATE',
 'LOGIN_DATE',
 'BOOKING_AMOUNT',
 'EOMCLTRL',
 'EOMLOGN',
 'EOMRJCT',
 'EOMSNCTN',
 'REQUESTED_AMOUNT','LOGIN_STATUS','NBFC_FLAG']

base_view=base_view[bv_cols]      
# lv_cols=find_columns(loan_view,dtcron) 
lv_cols=['BORROWER_TYPE',
 'BT_LOAN_LAN',
 'BT_LOAN_START_DATE',
 'BT_OUTSTANDING',
 'DELAY_REASON',
 'DST_CODE',
 'END_USAGE_FUNDS',
 'END_USE_FOR_TOPUP',
 'FINAL_FOIR',
 'FINAL_LOAN_AMOUNT',
 'FINAL_LTV',
 'FTR',
 'GHFAM',
 'GHFAM_BD',
 'GHFAM_BD_NAME',
 'GHFAM_NAME',
 'GHFSM',
 'GHFSM_NAME',
 'GHF_AM',
 'GRACE_TERMS',
 'GPL_FLAG',
 'INCOME_PROGRAM_TYPE',
 'INDIVIDUAL_DEVIATION_FLAG',
 'LAN_ID',
 'LEADID',
 'LOAN_PURPOSE',
 'LOAN_TYPE',
 'NO_FINANCE_DEVIATIONS',
 'PEP',
 'PRINCIPAL_OUTSTANDING',
 'PSL',
 'RISK_CATEGORIZATION',
 'ROI',
 'SELF_EMPLOYED_RISK',
 'TOTAL_TENOR',
 'TYPES_OF_REJECT']
loan_view=loan_view[lv_cols]
# cv_cols=find_columns(customer_view, dtcron)
cv_cols=['CUSTOMER_CITY_NAME', 'CUSTOMER_COUNTRY_DESC', 'CUSTOMER_INDUSTRY_DESC', 'CUSTOMER_QUALIFICATION_DESCRIPTION', 'CUSTOMER_RESIDENTIAL_STATUS', 'CUSTOMER_SECTOR_DESC', 'CUSTOMER_SUB_SECTOR_CODE', 'CUSTOMER_SUB_SECTOR_DESC', 'EMPLOYERNAME', 'LAN_ID', 'NUMBER_OF_DEPENDENTS', 'OCCUPATION_CATEGORY', 'QUALIFICATION', 'SCORE', 'SUB_CATEGORY']
cv_cols.append('CUSTOMER_CIF')
customer_view=customer_view[cv_cols]
# iv_cols=find_columns(insurance_view,dtcron)
iv_cols= ['CHANNEL_CODE', 'LAN_ID', 'NET_PREMIUM']
insurance_view=insurance_view[iv_cols]
insurance_view.set_index('LAN_ID', inplace=True)
insurance_view.reset_index(drop=False, inplace=True)
loan_view.set_index('LAN_ID', inplace=True)
loan_view.reset_index(drop=False, inplace=True)
base_view.set_index('LAN_ID', inplace=True)
base_view.reset_index(drop=False, inplace=True)
customer_view.set_index('LAN_ID', inplace=True)
customer_view.reset_index(drop=False, inplace=True)




insurance_view=remove_error_lans(insurance_view)
loan_view=remove_error_lans(loan_view)
base_view=remove_error_lans(base_view)
customer_view=remove_error_lans(customer_view)
base_view=base_view.merge(customer_view,on='LAN_ID',how='left')
base_view=base_view.merge(loan_view, left_on='LAN_ID',right_on='LAN_ID',how='left')
insurance_view=insurance_view.groupby('LAN_ID',as_index=False).sum('NET_PREMIUM')
base_view=base_view.merge(insurance_view,left_on='LAN_ID',right_on='LAN_ID',how='left')
base_view['WROI']=base_view.BOOKING_AMOUNT*base_view.ROI
base_view['PROCESSING_FEE']=0
base_view['BOOKING_DATE'].replace(to_replace=[np.nan,'',' '],value=date(1900,1,1),inplace=True)
base_view['BOOKING_YEAR']=pd.to_datetime(base_view['BOOKING_DATE']).dt.year
# base_view['BOOKING_YEAR']=base_view['BOOKING_YEAR'].astype(int)
base_view['BOOKING_MONTH']=pd.to_datetime(base_view['BOOKING_DATE']).dt.month
# base_view['BOOKING_MONTH']=base_view['BOOKING_MONTH'].astype(int)
base_view['BOOK_YEAR_MONTH']=base_view['BOOKING_YEAR'].astype(str)+'-'+ base_view['BOOKING_MONTH'].astype(str)
base_view['FINTYPE']=np.where(base_view['LOAN_PURPOSE'].isin(['LAP Balance Transfer plus Top-up','Loan against Property', 'Industrial LAP Balance Transfer','Industrial LAP Balance Transfer plus Top-up','LAP Balance Transfer','Loan against industrial property','LAP Top Up']),'LP',base_view['FINANCE_TYPE'])
base_view['GPLFLAG_SANCTIONS']=np.where((base_view['FINTYPE'].isin(['LP','NP'])==False) & (base_view['GPL_FLAG']=='YES') & (base_view['NBFC_FLAG']=='N') , 'GPL',np.where((base_view['FINTYPE'].isin(['LP','NP'])==False) &(base_view['GPL_FLAG']=='NO') & (base_view['NBFC_FLAG']=='N'),'NON GPL','NIL'))


ghf_disb=pd.read_sql("select FINANCE_REFERENCE as LAN_ID, DISBURSEMENT_AMOUNT as DISBAMOUNT,DISBURSEMENT_SEQUENCE as DISBSEQ from  prod_da_db.serve.f_finance_disbursement_details WHERE DH_RECORD_ACTIVE_FLAG = 'Y';", con_dm)
ghf_disb['NBFC_FLAG'] = 'N'
gfl_disb=pd.read_sql("select FINANCE_REFERENCE as LAN_ID, DISBURSEMENT_AMOUNT as DISBAMOUNT,DISBURSEMENT_SEQUENCE as DISBSEQ from  prod_gfl_da_db.serve.f_finance_disbursement_details WHERE DH_RECORD_ACTIVE_FLAG = 'Y';", con_dm_nbfc)
gfl_disb['NBFC_FLAG'] = 'Y'
disb=pd.concat([ghf_disb,gfl_disb])


fd=disb.copy()


'''
#################################################################################
'''








# query = "SELECT * FROM DTCRON_MASTER_05112022;"
# q2="select * from V_FINDISBURSEMENTDETAILS;"
# mydata = pd.read_sql(query, conx)
mydata=base_view.copy()
# fd=pd.read_sql(q2,conx)
df=mydata.copy()
target=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\October targets.xlsx")
target_df=target.copy()
target_df=target_df[['Unnamed: 20','Unnamed: 21','Unnamed: 22','Unnamed: 23']]
length=int(len(target_df)/2)

def change(df):
    if ('AMD' in df.index)|('NCR' in df.index)|('Total ' in df.index)|('Mumbai ' in df.index)|('Bangalore ' in df.index):
        a=df.index.get_indexer(['AMD','NCR','Total','Mumbai','Bangalore'])
        b=df.index.to_list()
        b[a[0]]='Ahmedabad'
        b[a[1]]='Delhi'
        b[a[2]]='Total'
        b[a[3]]='Mumbai'
        b[a[4]]='Bangalore'
        df.index=b
    return df
mtd_target_df=target_df.iloc[4:length]
mtd_target_df.reset_index(drop=True,inplace=True)
ytd_target_df=target_df.iloc[length+3:]
ytd_target_df.reset_index(drop=True,inplace=True)
ytd_target_df.loc[0,'Unnamed: 23']='AUM'
mtd_target_df.loc[0,'Unnamed: 23']='AUM'
ytd_target_df.loc[0,'Unnamed: 20']='BRANCH'
mtd_target_df.loc[0,'Unnamed: 20']='BRANCH'
ytd_target_df.columns = ytd_target_df.iloc[0]
mtd_target_df.columns = mtd_target_df.iloc[0]
ytd_target_df.drop(0,inplace=True)
mtd_target_df.drop(0,inplace=True)

mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 3','Unnamed: 4','Unnamed: 5','Unnamed: 6']]
length=int(len(target_df)/2)
GPL_mtd_target_df=target_df.iloc[4:length]
GPL_mtd_target_df.reset_index(drop=True,inplace=True)
GPL_ytd_target_df=target_df.iloc[length+3:]
GPL_ytd_target_df.reset_index(drop=True,inplace=True)
GPL_ytd_target_df.loc[0,'Unnamed: 6']='AUM'
GPL_mtd_target_df.loc[0,'Unnamed: 6']='AUM'
GPL_ytd_target_df.loc[0,'Unnamed: 3']='BRANCH'
GPL_mtd_target_df.loc[0,'Unnamed: 3']='BRANCH'
GPL_ytd_target_df.columns = GPL_ytd_target_df.iloc[0]
GPL_mtd_target_df.columns = GPL_mtd_target_df.iloc[0]
GPL_ytd_target_df.drop(0,inplace=True)
GPL_mtd_target_df.drop(0,inplace=True)
GPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
GPL_mtd_target_df.loc[6,'BRANCH']=0
GPL_mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 9','Unnamed: 10','Unnamed: 11','Unnamed: 12']]
length=int(len(target_df)/2)
NGPL_mtd_target_df=target_df.iloc[4:length]
NGPL_mtd_target_df.reset_index(drop=True,inplace=True)
NGPL_ytd_target_df=target_df.iloc[length+3:]
NGPL_ytd_target_df.reset_index(drop=True,inplace=True)
NGPL_ytd_target_df.loc[0,'Unnamed: 12']='AUM'
NGPL_mtd_target_df.loc[0,'Unnamed: 12']='AUM'
NGPL_ytd_target_df.loc[0,'Unnamed: 9']='BRANCH'
NGPL_mtd_target_df.loc[0,'Unnamed: 9']='BRANCH'
NGPL_ytd_target_df.columns = NGPL_ytd_target_df.iloc[0]
NGPL_mtd_target_df.columns = NGPL_mtd_target_df.iloc[0]
NGPL_ytd_target_df.drop(0,inplace=True)
NGPL_mtd_target_df.drop(0,inplace=True)
NGPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
NGPL_mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 15','Unnamed: 16','Unnamed: 17','Unnamed: 18']]
length=int(len(target_df)/2)
LAP_mtd_target_df=target_df.iloc[4:length]
LAP_mtd_target_df.reset_index(drop=True,inplace=True)
LAP_ytd_target_df=target_df.iloc[length+3:]
LAP_ytd_target_df.reset_index(drop=True,inplace=True)
LAP_ytd_target_df.loc[0,'Unnamed: 18']='AUM'
LAP_mtd_target_df.loc[0,'Unnamed: 18']='AUM'
LAP_ytd_target_df.loc[0,'Unnamed: 15']='BRANCH'
LAP_mtd_target_df.loc[0,'Unnamed: 15']='BRANCH'
LAP_ytd_target_df.columns = LAP_ytd_target_df.iloc[0]
LAP_mtd_target_df.columns = LAP_mtd_target_df.iloc[0]
LAP_ytd_target_df.drop(0,inplace=True)
LAP_mtd_target_df.drop(0,inplace=True)
LAP_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
LAP_mtd_target_df.set_index('BRANCH',inplace=True)


target_df=target.copy()
target_df=target_df[['Unnamed: 25','Unnamed: 26','Unnamed: 27','Unnamed: 28']]
length=int(len(target_df)/2)
HL_mtd_target_df=target_df.iloc[4:length]
HL_mtd_target_df.reset_index(drop=True,inplace=True)
HL_ytd_target_df=target_df.iloc[length+3:]
HL_ytd_target_df.reset_index(drop=True,inplace=True)
HL_ytd_target_df.loc[0,'Unnamed: 28']='AUM'
HL_mtd_target_df.loc[0,'Unnamed: 28']='AUM'
HL_ytd_target_df.loc[0,'Unnamed: 25']='BRANCH'
HL_mtd_target_df.loc[0,'Unnamed: 25']='BRANCH'
HL_ytd_target_df.columns = HL_ytd_target_df.iloc[0]
HL_mtd_target_df.columns = HL_mtd_target_df.iloc[0]
HL_ytd_target_df.drop(0,inplace=True)
HL_mtd_target_df.drop(0,inplace=True)
HL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
HL_mtd_target_df.set_index('BRANCH',inplace=True)

HL_mtd_target_df=change(HL_mtd_target_df)
LAP_mtd_target_df=change(LAP_mtd_target_df)
NGPL_mtd_target_df=change(NGPL_mtd_target_df)
GPL_mtd_target_df=change(GPL_mtd_target_df)
mtd_target_df=change(mtd_target_df)
HL_mtd_target_df=HL_mtd_target_df.round(decimals=0)
LAP_mtd_target_df=LAP_mtd_target_df.round(decimals=0)
mtd_target_df=mtd_target_df.round(decimals=0)
GPL_mtd_target_df=GPL_mtd_target_df.round(decimals=0)
NGPL_mtd_target_df=NGPL_mtd_target_df.round(decimals=0)
'''
BOOKING_DATE BOOKING_DATE
BORROWER_TYPE BORROWER_TYPE
BT_LOAN_LAN BT_LOAN_LAN
BT_LOAN_START_DATE BT_LOAN_START_DATE
BT_OUTSTANDING BT_OUTSTANDING
CUSTINDUSTRY CUSTOMER_INDUSTRY_DESC
DETAILED_STATUS DETAILED_STATUS
DST_CODE DST_CODE
EMPLOYERNAME EMPLOYERNAME
EMPLOYMENT_TYPE SUB_CATEGORY
END_USAGE_FUNDS END_USAGE_FUNDS
REPORTING_BRANCH REPORTING_BRANCH
LAN_ID LAN_ID
FINTYPE FINTYPE
FINTYPE_TECH FINTYPE
FIRST_ENTRY_TO_COLLATERAL FIRST_ENTRY_TO_COLLATERAL
FIRST_REJECT_DATE EOMRJCT
FTR FTR
GHFAM GHFAM
GHFAM_BD GHFAM_BD
GHFAM_BD_NAME GHFAM_BD_NAME
GHFAM_NAME GHFAM_NAME
GHFSM GHFSM
GHFSM_NAME GHFSM_NAME
INCOME_PROGRAM_TYPE INCOME_PROGRAM_TYPE
INDIVIDUAL_DEVIATION_FLAG INDIVIDUAL_DEVIATION_FLAG
LEADID LEADID
LOAN_PURPOSE LOAN_PURPOSE
LOAN_STATUS LOAN_STATUS
LOGIN_DATE LOGIN_DATE
LOGIN_MONTH NO_FINANCE_DEVIATIONS
OCCUPATION OCCUPATION_CATEGORY
PEP PEP
PRINCIPAL_OUTSTANDING PRINCIPAL_OUTSTANDING
PSL PSL
QUALIFICATION CUSTOMER_QUALIFICATION_DESCRIPTION
QUEUE QUEUE
RESIDENTIAL_STATUS CUSTOMER_RESIDENTIAL_STATUS
ROI ROI
SANCTION_AMOUNT SANCTION_AMOUNT
SECTOR CUSTOMER_SECTOR_DESC
STATUS STATUS
STATUS_SEG STATUS_SEG
TOTAL_TENOR TOTAL_TENOR
TYPES_OF_REJECT TYPES_OF_REJECT
'''

'''
#--------------LAP_BUSINESS METRICS-------------------

#------------BOOK-----------

LP_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')

#------------ROI------------
LP_ROI<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
               & BM$BOOK_YEAR_MONTH == MTD)
#------------PF-------------
LP_PF<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
                     & BM$BOOK_YEAR_MONTH == MTD)

#------------LTD (BOOKING TO LOGIN)----------
LTD_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')
LTD_LOGIN<-subset (BM,BM$FINTYPE %in% c('LP','NP') & BM$LOGIN_STATUS == 'A) Login')

LTD_BOOK1<-ddply(LTD_BOOK,c('REPORTING_BRANCH'),summarise,
              Book_value=sum(BOOKING_AMOUNT/10000000), 
              Book_vol=length(LAN_ID))
LTD_LOGIN1<-ddply(LTD_LOGIN,c('REPORTING_BRANCH'),summarise,
                  Log_value=sum(REQUESTED_AMOUNT/1000000000),
                  Log_Vol=length(LAN_ID))

LTD_BOOK_LOG_mrg<-merge(LTD_LOGIN1,LTD_BOOK1,by.x = c('REPORTING_BRANCH'),
                        by.y = c('REPORTING_BRANCH'),all.x = T)
'''


todays_date = date.today()

def zonwise(newb,zone):
    newb=newb.transpose()
    newb.reset_index(inplace=True)
    newb.loc[11]=newb.sum()
    newb.loc[11,'REPORTING_BRANCH']=zone
    newb['BOOK_VALUE'] = newb.BOOKING_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['REPORTING_BRANCH','one_count','BOOKING_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]
    newb['BOOK_VOL'] = newb.one_count_MTD
    newb['DISB AMNT Tranch 1']=newb.DISBAMOUNT/1000000000
    newb['DISB_Act (%)']=(newb.DISBAMOUNT)/newb.BOOKING_AMOUNT
    newb['AUM']=newb.PRINCIPAL_OUTSTANDING/10000000
    newb['ROI MTD'] = newb.WROI / newb.BOOKING_AMOUNT
    newb['ROI M-1'] = newb.WROI_M1 /newb.BOOKING_AMOUNT_M1
    newb['GROSS%']=100*newb.NET_PREMIUM/newb.BOOKING_AMOUNT
    newb['PF%']=100*newb.PROCESSING_FEE / newb.BOOKING_AMOUNT
    newb['LTD_volume']=100*newb.one_count_tot_book/newb.one_count
    newb['LTD_Value']=100*(newb.TOTAL_BOOKING_AMOUNTytd*100)/newb.REQUESTED_AMOUNT
    newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    newb.loc['Total']= newb.sum()
    newb.loc['Total','DISB_Act (%)']=np.float64(sum(newb.DISBAMOUNT))/np.float64(sum(newb.BOOKING_AMOUNT))
    newb.loc['Total','ROI MTD']=np.float64(sum(newb.WROI )) /np.float64(sum( newb.BOOKING_AMOUNT))
    newb.loc['Total','ROI M-1']=np.float64(sum(newb.WROI_M1))/np.float64(sum(newb.BOOKING_AMOUNT_M1))
    newb.loc['Total','GROSS%']=100*np.float64(sum(newb.NET_PREMIUM))/np.float64(sum(newb.BOOKING_AMOUNT))
    newb.loc['Total','PF%']=100*np.float64(sum(newb.PROCESSING_FEE )) / np.float64(sum(newb.BOOKING_AMOUNT))
    newb.loc['Total','LTD_volume']=100*np.float64(sum(newb.one_count_tot_book))/np.float64(sum(newb.one_count))
    newb.loc['Total','LTD_Value']=100*np.float64(sum(newb.TOTAL_BOOKING_AMOUNTytd))*(100)/np.float64(sum(newb.REQUESTED_AMOUNT))
    newb.loc['Total','REPORTING_BRANCH']='Total'
    newb.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKING_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKING_AMOUNT','BOOKING_AMOUNT_M1','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    newb["new"] = range(1,len(newb)+1)
    newb.loc['Total','new'] = 0
    newb=newb.sort_values(by="new").drop('new', axis=1)
    newb=newb.round(decimals=2)
    newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    newb.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    newb.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    newb.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    newb.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    newb.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    newb.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    newb.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    newb.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    newb['Booking Tar(Cr)']=' '
    newb['Booking Tar(#)']=' '
    newb['AUM Tar(Cr)']=' '
    newb['Net']=' '
    newb['Empaneled('+calendar.month_name[todays_date.month-1][:3]+')']=' '
    newb['Active']=' '
    newb['Channel']=' '
    newb['Employee']=' '
    newb['Cumulative']=' '
    newb.columns
    newb=newb.round(decimals=2)
    newb['Value%']=round(newb['Value%'])
    newb['Volume%']=round(newb['Volume%'])
    newb=newb[['REPORTING_BRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Booking Tar(#)', 'Booking Act(#)','DISB_Act (%)', 'DISB Tr1','AUM Tar(Cr)','AUM Act(Cr)', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled('+calendar.month_name[todays_date.month-1][:3]+')', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #newb.set_index('REPORTING_BRANCH',inplace=True)
    finaln=newb.transpose()
    finaln.insert(0,'Target',' ')
    finaln.loc['REPORTING_BRANCH','Target']='Target'
    return finaln


def buss_metrics(LOAN,df,fd,todays_date):

    # MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
    MTD=str(2022)+"-"+str(12)
    YTD=[]
    cur_month=todays_date.month
    while cur_month>=4:
        YTD.append(cur_month)
        cur_month=cur_month-1
    ytd=[]
    for i in YTD:
        ytd.append(str(todays_date.year)+"-"+str(i))
    
    ytd=['2022-12','2022-11','2022-10','2022-9','2022-8','2022-7','2022-6','2022-5','2022-4']
    BM_df = df[[ 'LAN_ID','NET_PREMIUM','BOOK_YEAR_MONTH','REPORTING_BRANCH','FINTYPE','STATUS', 'WROI','BOOKING_AMOUNT' ,'PROCESSING_FEE','REQUESTED_AMOUNT','LOGIN_STATUS','GPLFLAG_SANCTIONS','PRINCIPAL_OUTSTANDING']]
    findisb=fd[['LAN_ID','DISBSEQ','DISBAMOUNT']]
    findisb=findisb[(findisb['DISBSEQ']==1)]
    findisb=findisb.sort_values(by='LAN_ID')
    BM_df=BM_df.sort_values(by='LAN_ID')
    fd_bm= pd.merge(BM_df, findisb,left_on=['LAN_ID'],right_on=['LAN_ID'], how='inner')
    

    if LOAN== 'LAP':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGIN_STATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'LP') | (fd_bm['FINTYPE'] =='NP'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'LP') |  (bm_df['FINTYPE'] =='NP'))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'HL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGIN_STATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'HL') | (fd_bm['FINTYPE'] =='HT') | (fd_bm['FINTYPE'] =='FL') | (fd_bm['FINTYPE'] =='FT') | (fd_bm['FINTYPE'] =='LT')| (fd_bm['FINTYPE'] =='AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'HL') | (bm_df['FINTYPE'] =='HT') | (bm_df['FINTYPE'] =='FL') | (bm_df['FINTYPE'] =='FT') | (bm_df['FINTYPE'] =='LT')| (bm_df['FINTYPE'] =='AHL'))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'GPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGIN_STATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'] == 'GPL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'] == 'GPL')&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'NGPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGIN_STATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL'])))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL']))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'AHL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGIN_STATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login'))&(bm_df['FINTYPE'] == 'AHL')&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    
    # bm_df_hl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login') &((bm_df['FINTYPE'] == 'HL') |  (bm_df['FINTYPE'] =='HT'))]
    # bm_df_gpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='GPL')]
    # bm_df_ngpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGIN_STATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='NON GPL')]
    bm_df['one_count']=np.ones(len(bm_df), dtype = int)
    
    bm_df_group=bm_df.groupby(['REPORTING_BRANCH'],as_index=False).sum(['DISBSEQ','LAN_ID','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKING_AMOUNT','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ])
    # def buss_metrics(bm_df):
        #bm_df_group = bm_df.groupby(['REPORTING_BRANCH'],as_index=False).sum(['WROI','PROCESSING_FEE','BOOKING_AMOUNT','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING' ])
        #bm_df_group['BOOK_TARGET'] =
        
    fd_bm_group=fd_bm.groupby(['REPORTING_BRANCH'],as_index=False).sum(['WROI', 'BOOKING_AMOUNT', 'PROCESSING_FEE','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING',  'DISBSEQ', 'DISBAMOUNT'])
    fd_bm_group=fd_bm_group[['REPORTING_BRANCH','DISBAMOUNT']]
    bm_df_book=bm_df[(bm_df.STATUS== 'Booked')&(bm_df.BOOK_YEAR_MONTH==MTD)]
    # bm_df_book=pd.merge(bm_df_book, findisb,left_on=['LAN_ID'],right_on=['LAN_ID'], how='inner')
    bm_df_book['one_count']=np.ones(len(bm_df_book), dtype = int)
    bm_df_BOOK_group=bm_df_book.groupby(['REPORTING_BRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKING_AMOUNT','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ])
    
    
    b=bm_df[(bm_df.STATUS== 'Booked')]
    c=bm_df[((bm_df.STATUS== 'Booked') & (bm_df.BOOK_YEAR_MONTH.isin(ytd)))]
    c_group=c.groupby(['REPORTING_BRANCH'],as_index=False).sum(['BOOKING_AMOUNT','WROI'])
    c_group=c_group[['REPORTING_BRANCH','BOOKING_AMOUNT','WROI']]
    c_group.rename(columns = {'BOOKING_AMOUNT':'BOOKING_AMOUNT_M1'}, inplace = True)
    c_group.rename(columns = {'WROI':'WROI_M1'}, inplace = True)
    b_group=b.groupby(['REPORTING_BRANCH'],as_index=False).sum(['BOOKING_AMOUNT','PRINCIPAL_OUTSTANDING','one_count' ])
    b_group=b_group[['REPORTING_BRANCH','PRINCIPAL_OUTSTANDING','one_count','BOOKING_AMOUNT' ]]
    b_group.rename(columns = {'one_count':'one_count_tot_book'}, inplace = True)
    b_group.rename(columns = {'BOOKING_AMOUNT':'TOTAL_BOOKING_AMOUNTytd'}, inplace = True)
    bm_df_group=bm_df.groupby(['REPORTING_BRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKING_AMOUNT','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ])
    bm_df_BOOK_group=bm_df_BOOK_group[['REPORTING_BRANCH','BOOKING_AMOUNT','WROI','NET_PREMIUM','one_count','PROCESSING_FEE']]
    bm_df_BOOK_group.rename(columns = {'one_count':'one_count_MTD'}, inplace = True)
    bm_df_group=bm_df_group[['REPORTING_BRANCH','REQUESTED_AMOUNT','one_count']]
    
    merged=pd.merge(bm_df_group, bm_df_BOOK_group,left_on=['REPORTING_BRANCH'],right_on=['REPORTING_BRANCH'], how='outer')
    merged2=pd.merge(merged, c_group,left_on=['REPORTING_BRANCH'],right_on=['REPORTING_BRANCH'], how='outer')
    merged3=pd.merge(merged2, b_group,left_on=['REPORTING_BRANCH'],right_on=['REPORTING_BRANCH'], how='outer')
    merged4=pd.merge(merged3, fd_bm_group,left_on=['REPORTING_BRANCH'],right_on=['REPORTING_BRANCH'], how='outer')
    '''
    CHANGES AFTER TRANSITION
    '''
    if LOAN=='LAP':
        newb=merged4.copy()
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.set_index('REPORTING_BRANCH',inplace=True)
        newb=newb.transpose()
        new=newb.copy()
        finaln=zonwise(new[['Hyderabad','Indore', 'Jaipur', 'Chandigarh', 'Chennai', 'Surat']],'New Branches')
        new=newb.copy()
        north=zonwise(new[['Delhi','Chandigarh','Jaipur']],'North')
        new=newb.copy()
        south=zonwise(new[['Bangalore','Chennai','Hyderabad']],'South')
        new=newb.copy()
        west1=zonwise(new[['Mumbai', 'Ahmedabad','Surat']],'West 1')
        new=newb.copy()
        west2=zonwise(new[['Indore','Pune']],'West 2')
        
    
    merged4['BOOK_VALUE'] = merged4.BOOKING_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['REPORTING_BRANCH','one_count','BOOKING_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]

    merged4['BOOK_VOL'] = merged4.one_count_MTD
    merged4['DISB AMNT Tranch 1']=merged4.DISBAMOUNT/1000000000
    merged4['DISB_Act (%)']=(merged4.DISBAMOUNT)/merged4.BOOKING_AMOUNT
    merged4['AUM']=merged4.PRINCIPAL_OUTSTANDING/10000000
    merged4['ROI MTD'] = merged4.WROI / merged4.BOOKING_AMOUNT
    merged4['ROI M-1'] = merged4.WROI_M1 /merged4.BOOKING_AMOUNT_M1
    merged4['GROSS%']=100*merged4.NET_PREMIUM/merged4.BOOKING_AMOUNT
    merged4['PF%']=100*merged4.PROCESSING_FEE / merged4.BOOKING_AMOUNT
    merged4['LTD_volume']=100*merged4.one_count_tot_book/merged4.one_count
    merged4['LTD_Value']=100*(merged4.TOTAL_BOOKING_AMOUNTytd*100)/merged4.REQUESTED_AMOUNT
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    merged4.loc['Total']= merged4.sum()
    merged4.loc['Total','DISB_Act (%)']=np.float64(sum(merged4.DISBAMOUNT))/np.float64(sum(merged4.BOOKING_AMOUNT))
    merged4.loc['Total','ROI MTD']=np.float64(sum(merged4.WROI )) /np.float64(sum( merged4.BOOKING_AMOUNT))
    merged4.loc['Total','ROI M-1']=np.float64(sum(merged4.WROI_M1))/np.float64(sum(merged4.BOOKING_AMOUNT_M1))
    merged4.loc['Total','GROSS%']=100*np.float64(sum(merged4.NET_PREMIUM))/np.float64(sum(merged4.BOOKING_AMOUNT))
    merged4.loc['Total','PF%']=100*np.float64(sum(merged4.PROCESSING_FEE )) / np.float64(sum(merged4.BOOKING_AMOUNT))
    merged4.loc['Total','LTD_volume']=100*np.float64(sum(merged4.one_count_tot_book))/np.float64(sum(merged4.one_count))
    merged4.loc['Total','LTD_Value']=100*np.float64(sum(merged4.TOTAL_BOOKING_AMOUNTytd))*(100)/np.float64(sum(merged4.REQUESTED_AMOUNT))
    merged4.loc['Total','REPORTING_BRANCH']='Total'
    merged4.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKING_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKING_AMOUNT','BOOKING_AMOUNT_M1','REQUESTED_AMOUNT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    merged4["new"] = range(1,len(merged4)+1)
    merged4.loc['Total','new'] = 0
    merged4=merged4.sort_values(by="new").drop('new', axis=1)
    merged4=merged4.round(decimals=2)
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    merged4.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    merged4.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    merged4.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    merged4.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    merged4.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    merged4.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    merged4['Booking Tar(Cr)']=' '
    merged4['Booking Tar(#)']=' '
    merged4['AUM Tar(Cr)']=' '
    merged4['Net']=' '
    merged4['Empaneled('+calendar.month_name[todays_date.month-1][:3]+')']=' '
    merged4['Active']=' '
    merged4['Channel']=' '
    merged4['Employee']=' '
    merged4['Cumulative']=' '
    merged4.columns
    merged4=merged4.round(decimals=2)
    merged4['Value%']=round(merged4['Value%'])
    merged4['Volume%']=round(merged4['Volume%'])
    merged4=merged4[['REPORTING_BRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Booking Tar(#)', 'Booking Act(#)','DISB_Act (%)', 'DISB Tr1','AUM Tar(Cr)','AUM Act(Cr)', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled('+calendar.month_name[todays_date.month-1][:3]+')', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #merged4.set_index('REPORTING_BRANCH',inplace=True)
    

    final=merged4.transpose()
    
    final.insert(0,'Target',' ')
    final.loc['REPORTING_BRANCH','Target']='Target'
    if LOAN=='LAP':
        final.reset_index(inplace=True)
        finaln.reset_index(inplace=True)
        finaln=finaln[['index',11]]
        final=pd.merge(final,finaln,left_on=('index'),right_on=('index'),how='outer')
        final.set_index('index',inplace=True)
        return final,north,west1,west2,south
    return(final)
# return(final)
lap_metrics,n,w1,w2,s=buss_metrics('LAP',df,fd,todays_date)
n.loc[:,'Total']=lap_metrics.loc[:,'Total']
w1.loc[:,'Total']=lap_metrics.loc[:,'Total']
w2.loc[:,'Total']=lap_metrics.loc[:,'Total']
s.loc[:,'Total']=lap_metrics.loc[:,'Total']
lap_metrics=lap_metrics[['Target', 'Total', 1, 4, 8, 9, 0, 2, 3, 5, 6, 7, 10,11]]
lap_metrics.columns=lap_metrics.loc['REPORTING_BRANCH']
w2.columns=w2.loc['REPORTING_BRANCH']
n.columns=n.loc['REPORTING_BRANCH']
s.columns=s.loc['REPORTING_BRANCH']
w1.columns=w1.loc['REPORTING_BRANCH']
def set_target(lap_metrics,LAP_mtd_target_df):
    
    for i in LAP_mtd_target_df.index:
        if  (i != 0) & (i in lap_metrics.columns.to_list() ): 
            lap_metrics.loc['Booking Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'Value (Cr)']
            lap_metrics.loc['Booking Act(Cr)',i]=str(round(lap_metrics.loc['Booking Act(Cr)',i]))+'('+str(round(100*lap_metrics.loc['Booking Act(Cr)',i]/LAP_mtd_target_df.loc[i,'Value (Cr)']))+'%)'
            lap_metrics.loc['Booking Tar(#)',i]=LAP_mtd_target_df.loc[i,'#']
            lap_metrics.loc['Booking Act(#)',i]=str(round(lap_metrics.loc['Booking Act(#)',i]))+'('+str(round(100*lap_metrics.loc['Booking Act(#)',i]/LAP_mtd_target_df.loc[i,'#']))+'%)'
            lap_metrics.loc['AUM Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'AUM']
            lap_metrics.loc['AUM Act(Cr)',i]=str(round(lap_metrics.loc['AUM Act(Cr)',i]))+'('+str(round(100*lap_metrics.loc['AUM Act(Cr)',i]/LAP_mtd_target_df.loc[i,'AUM']))+'%)'
    return lap_metrics
hl_metrics=buss_metrics('HL',df,fd,todays_date)
hl_metrics.columns=hl_metrics.loc['REPORTING_BRANCH']
hl_metrics=hl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]


gpl_metrics=buss_metrics('GPL',df,fd,todays_date )
gpl_metrics.columns=gpl_metrics.loc['REPORTING_BRANCH']
gpl_metrics=gpl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]
ngpl_metrics=buss_metrics('NGPL',df,fd,todays_date )

ngpl_metrics.columns=ngpl_metrics.loc['REPORTING_BRANCH']
ngpl_metrics=ngpl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]

# ahl_metrics=buss_metrics('AHL',df,fd,todays_date )
# ahl_metrics.columns=ahl_metrics.loc['REPORTING_BRANCH']


hl_metrics=set_target(hl_metrics,HL_mtd_target_df)
lap_metrics=set_target(lap_metrics,LAP_mtd_target_df)
n=set_target(n,LAP_mtd_target_df)
s=set_target(s,LAP_mtd_target_df)
w1=set_target(w1,LAP_mtd_target_df)
w2=set_target(w2,LAP_mtd_target_df)
gpl_metrics=set_target(gpl_metrics,GPL_mtd_target_df)
ngpl_metrics=set_target(ngpl_metrics,NGPL_mtd_target_df)


lap_metrics.loc['DISB_Act (%)','Target']='90%'
lap_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
lap_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
lap_metrics.loc['GROSS%','Target']='1.5%'
lap_metrics.loc['PF%','Target']='0.9%'
lap_metrics.loc['Cumulative','Target']='1.35%'
lap_metrics.loc['Value%','Target']='33%'
lap_metrics.loc['Volume%','Target']='33%'

n.loc['DISB_Act (%)','Target']='90%'
n.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
n.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
n.loc['GROSS%','Target']='1.5%'
n.loc['PF%','Target']='0.9%'
n.loc['Cumulative','Target']='1.35%'
n.loc['Value%','Target']='33%'
n.loc['Volume%','Target']='33%'

s.loc['DISB_Act (%)','Target']='90%'
s.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
s.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
s.loc['GROSS%','Target']='1.5%'
s.loc['PF%','Target']='0.9%'
s.loc['Cumulative','Target']='1.35%'
s.loc['Value%','Target']='33%'
s.loc['Volume%','Target']='33%'

w1.loc['DISB_Act (%)','Target']='90%'
w1.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
w1.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
w1.loc['GROSS%','Target']='1.5%'
w1.loc['PF%','Target']='0.9%'
w1.loc['Cumulative','Target']='1.35%'
w1.loc['Value%','Target']='33%'
w1.loc['Volume%','Target']='33%'

w2.loc['DISB_Act (%)','Target']='90%'
w2.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
w2.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
w2.loc['GROSS%','Target']='1.5%'
w2.loc['PF%','Target']='0.9%'
w2.loc['Cumulative','Target']='1.35%'
w2.loc['Value%','Target']='33%'
w2.loc['Volume%','Target']='33%'

# hl_metrics.loc['DISB_Act (%)','Target']='90%'
hl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
hl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
hl_metrics.loc['GROSS%','Target']='2%'
hl_metrics.loc['PF%','Target']='0.02%'
hl_metrics.loc['Cumulative','Target']='0.45%'
hl_metrics.loc['Value%','Target']='50%'
hl_metrics.loc['Volume%','Target']='50%'

gpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
gpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
gpl_metrics.loc['GROSS%','Target']='2%'
gpl_metrics.loc['PF%','Target']='0.02%'
gpl_metrics.loc['Cumulative','Target']='0.45%'
gpl_metrics.loc['Value%','Target']='50%'
gpl_metrics.loc['Volume%','Target']='50%'


ngpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
ngpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
ngpl_metrics.loc['GROSS%','Target']='2%'
ngpl_metrics.loc['PF%','Target']='0.02%'
ngpl_metrics.loc['Cumulative','Target']='0.45%'
ngpl_metrics.loc['Value%','Target']='50%'
ngpl_metrics.loc['Volume%','Target']='50%'


MTD=str(2022)+"-"+str(12)
with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\1"+MTD+".xlsx") as writer:
    lap_metrics.to_excel(writer, sheet_name="LAP BM")
    s.to_excel(writer, sheet_name="SOUTH BM")
    n.to_excel(writer, sheet_name="NORTH BM")
    w1.to_excel(writer, sheet_name="WEST1 BM")
    w2.to_excel(writer, sheet_name="WEST2 BM")
    hl_metrics.to_excel(writer, sheet_name="HL BM")
    gpl_metrics.to_excel(writer, sheet_name="GPL BM")
    ngpl_metrics.to_excel(writer, sheet_name="NGPL BM")
    # ahl_metrics.to_excel(writer, sheet_name="AHL BM")

from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# ---create presentation with 1 slide---
prs = Presentation(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\OPSREVIEWDECK.pptx")
# prs.slide_width = 11887200
# prs.slide_height = 6686550
def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
def create_slide(prs,lap_metrics,title,c,target):
    prs.slides[c].shapes[0].text=title
    lap_metrics.reset_index(level=0, inplace=True)
    lap_metrics.insert(0,' ',' ')
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    # slide.shapes.title.text = title
    # slide.shapes.title.top = Inches(0.5)
    # slide.shapes.title.left = Inches(0.5)
    # slide.shapes.title.width = Inches(10)

    
    # # ---add table to slide---
    # x, y, cx, cy = Inches(0.5), Inches(1), Inches(11), Inches(3)
    # shape = slide.shapes.add_table(len(lap_metrics),len(lap_metrics.columns),x, y, cx,cy)
    
    # tbl =  shape._element.graphic.graphicData.tbl
    # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    # tbl[0][-1].text = style_id    
    # table = shape.table
    
    cols=prs.slides[c].shapes[2].table.columns
    cols[1].width=Inches(1.4)
    prs.slides[c].shapes[2].table.width=Inches(12)
    prs.slides[c].shapes[2].table.height=Inches(5)
    # cols[0].width=Inches(1.4)
    # rows=table.rows
    for i in range(len(lap_metrics)):
        for j in range(len(lap_metrics.columns)):
            if ((i in (5,11,13,19,20))& (j>=2)):
                prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])+'%'
            else:
                prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])
                    
                
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_bottom = Inches(0.08)
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_left = 0
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.word_wrap = False
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                #prs.slides[c].shapes[2].table.cell(i,j).text_frame.fit_text()
                # if i==0:
                #     cell=table.cell(i, j)
                #     fill = cell.fill
                #     fill.solid()
                #     fill.fore_color.rgb = RGBColor(173, 216, 230)
    
                
    
    
    prs.slides[c].shapes[2].table.cell(1,0).text="Business Number"
    prs.slides[c].shapes[2].table.cell(9,0).text="ROI"
    prs.slides[c].shapes[2].table.cell(11,0).text="%Insurance"
    prs.slides[c].shapes[2].table.cell(13,0).text="PF%"
    prs.slides[c].shapes[2].table.cell(14,0).text="Number of DSAs"
    prs.slides[c].shapes[2].table.cell(16,0).text="Incentive"
    prs.slides[c].shapes[2].table.cell(19,0).text="LTD Booking to Login"
    
    for cell in iter_cells(prs.slides[c].shapes[2].table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if title=="LAP Business Metrics":
                    run.font.size = Pt(11)
                else:
                    run.font.size = Pt(14)
                run.font.style='Calibri'
                run.font.color.rgb = RGBColor(0,0,0)
                run.font.bold = False
    
    return(prs)

MTD=str(2022)+"-"+str(12)
create_slide(prs,lap_metrics,"LAP Business Metrics",13,target)            
create_slide(prs,hl_metrics,"HL Business Metrics",14,target)       
create_slide(prs,gpl_metrics,"GPL Business Metrics",15,target)
create_slide(prs,ngpl_metrics,"N-GPL Business Metrics",16,target)
create_slide(prs,n,"LAP Business Metrics North",17,target)   
create_slide(prs,s,"LAP Business Metrics South",18,target)   
create_slide(prs,w1,"LAP Business Metrics West 1",19,target)   
create_slide(prs,w2,"LAP Business Metrics West 2",20,target)   
# create_slide(prs,ahl_metrics,"AHL Business Metrics",17,target)
prs.save(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\NOV_OPS_REVIEW\OPERATING REVIEW"+MTD +"NEW.pptx")



