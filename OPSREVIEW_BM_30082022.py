# -*- coding: utf-8 -*-
"""
Created on Tue Aug 30 14:17:13 2022

@author: VAIBHAV.SRIVASTAV01
"""





import pandas as pd
import os
from datetime import date
import calendar
import numpy as np
from pptx import Presentation

import pyodbc
pyodbc.drivers()
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
config = {}
with open(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\config.txt","r") as file:
    for line in file:
        key, value = line.strip().split(",")
        config[key] = value

# del key,line,value,file
# mydata=pd.read_csv(r"C:\Users\VAIBHAV.SRIVASTAV01\.spyder-py3\DTCRON_MASTER_03072022.csv")
# mydata=pd.read_excel(r"C:\Users\vaibhav.srivastav01\Downloads\DTCRON (18 AUG).xlsx")
# #fd=pd.read_csv(r"C:\Users\VAIBHAV.SRIVASTAV01\.spyder-py3\FINDISBURSEMENTDETAILS (30th June).csv")
# fd=pd.read_csv(r"C:\Users\vaibhav.srivastav01\Desktop\findisb1808.csv")


del key,line,value,file

conx = pyodbc.connect('DRIVER={SnowflakeDSIIDriver};SERVER='+str(config['server'])+';Warehouse = COMPUTE_WH;DATABASE='+str(config['db1'])+';UID='+ str(config['user1']) +';PWD='+ str(config['password']),autocommit=True)
cur = conx.cursor()


cur.execute("USE DATABASE SAMPLE_DB")

query = "SELECT * FROM DTCRON_MASTER_05012023;"
q2="select * from V_FINDISBURSEMENTDETAILS;"
mydata = pd.read_sql(query, conx)
fd=pd.read_sql(q2,conx)
df=mydata.copy()
target=pd.read_excel(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\OPS\December Targets.xlsx")
target_df=target.copy()
target_df=target_df[['Unnamed: 20','Unnamed: 21','Unnamed: 22','Unnamed: 23']]
length=int(len(target_df)/2)

def change(df):
    if ('AMD' in df.index)|('NCR' in df.index)|('Total ' in df.index)|('Mumbai ' in df.index)|('Bangalore ' in df.index):
        a=df.index.get_indexer(['AMD','NCR','Total','Mumbai','Bangalore'])
        b=df.index.to_list()
        b[a[0]]='Ahmedabad'
        b[a[1]]='Delhi'
        b[a[2]]='Total'
        b[a[3]]='Mumbai'
        b[a[4]]='Bangalore'
        df.index=b
    return df
mtd_target_df=target_df.iloc[4:length]
mtd_target_df.reset_index(drop=True,inplace=True)
ytd_target_df=target_df.iloc[length+3:]
ytd_target_df.reset_index(drop=True,inplace=True)
ytd_target_df.loc[0,'Unnamed: 23']='AUM'
mtd_target_df.loc[0,'Unnamed: 23']='AUM'
ytd_target_df.loc[0,'Unnamed: 20']='BRANCH'
mtd_target_df.loc[0,'Unnamed: 20']='BRANCH'
ytd_target_df.columns = ytd_target_df.iloc[0]
mtd_target_df.columns = mtd_target_df.iloc[0]
ytd_target_df.drop(0,inplace=True)
mtd_target_df.drop(0,inplace=True)

mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 3','Unnamed: 4','Unnamed: 5','Unnamed: 6']]
length=int(len(target_df)/2)
GPL_mtd_target_df=target_df.iloc[4:length]
GPL_mtd_target_df.reset_index(drop=True,inplace=True)
GPL_ytd_target_df=target_df.iloc[length+3:]
GPL_ytd_target_df.reset_index(drop=True,inplace=True)
GPL_ytd_target_df.loc[0,'Unnamed: 6']='AUM'
GPL_mtd_target_df.loc[0,'Unnamed: 6']='AUM'
GPL_ytd_target_df.loc[0,'Unnamed: 3']='BRANCH'
GPL_mtd_target_df.loc[0,'Unnamed: 3']='BRANCH'
GPL_ytd_target_df.columns = GPL_ytd_target_df.iloc[0]
GPL_mtd_target_df.columns = GPL_mtd_target_df.iloc[0]
GPL_ytd_target_df.drop(0,inplace=True)
GPL_mtd_target_df.drop(0,inplace=True)
GPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
GPL_mtd_target_df.loc[6,'BRANCH']=0
GPL_mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 9','Unnamed: 10','Unnamed: 11','Unnamed: 12']]
length=int(len(target_df)/2)
NGPL_mtd_target_df=target_df.iloc[4:length]
NGPL_mtd_target_df.reset_index(drop=True,inplace=True)
NGPL_ytd_target_df=target_df.iloc[length+3:]
NGPL_ytd_target_df.reset_index(drop=True,inplace=True)
NGPL_ytd_target_df.loc[0,'Unnamed: 12']='AUM'
NGPL_mtd_target_df.loc[0,'Unnamed: 12']='AUM'
NGPL_ytd_target_df.loc[0,'Unnamed: 9']='BRANCH'
NGPL_mtd_target_df.loc[0,'Unnamed: 9']='BRANCH'
NGPL_ytd_target_df.columns = NGPL_ytd_target_df.iloc[0]
NGPL_mtd_target_df.columns = NGPL_mtd_target_df.iloc[0]
NGPL_ytd_target_df.drop(0,inplace=True)
NGPL_mtd_target_df.drop(0,inplace=True)
NGPL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
NGPL_mtd_target_df.set_index('BRANCH',inplace=True)

target_df=target.copy()
target_df=target_df[['Unnamed: 15','Unnamed: 16','Unnamed: 17','Unnamed: 18']]
length=int(len(target_df)/2)
LAP_mtd_target_df=target_df.iloc[4:length]
LAP_mtd_target_df.reset_index(drop=True,inplace=True)
LAP_ytd_target_df=target_df.iloc[length+3:]
LAP_ytd_target_df.reset_index(drop=True,inplace=True)
LAP_ytd_target_df.loc[0,'Unnamed: 18']='AUM'
LAP_mtd_target_df.loc[0,'Unnamed: 18']='AUM'
LAP_ytd_target_df.loc[0,'Unnamed: 15']='BRANCH'
LAP_mtd_target_df.loc[0,'Unnamed: 15']='BRANCH'
LAP_ytd_target_df.columns = LAP_ytd_target_df.iloc[0]
LAP_mtd_target_df.columns = LAP_mtd_target_df.iloc[0]
LAP_ytd_target_df.drop(0,inplace=True)
LAP_mtd_target_df.drop(0,inplace=True)
LAP_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
LAP_mtd_target_df.set_index('BRANCH',inplace=True)


target_df=target.copy()
target_df=target_df[['Unnamed: 25','Unnamed: 26','Unnamed: 27','Unnamed: 28']]
length=int(len(target_df)/2)
HL_mtd_target_df=target_df.iloc[4:length]
HL_mtd_target_df.reset_index(drop=True,inplace=True)
HL_ytd_target_df=target_df.iloc[length+3:]
HL_ytd_target_df.reset_index(drop=True,inplace=True)
HL_ytd_target_df.loc[0,'Unnamed: 28']='AUM'
HL_mtd_target_df.loc[0,'Unnamed: 28']='AUM'
HL_ytd_target_df.loc[0,'Unnamed: 25']='BRANCH'
HL_mtd_target_df.loc[0,'Unnamed: 25']='BRANCH'
HL_ytd_target_df.columns = HL_ytd_target_df.iloc[0]
HL_mtd_target_df.columns = HL_mtd_target_df.iloc[0]
HL_ytd_target_df.drop(0,inplace=True)
HL_mtd_target_df.drop(0,inplace=True)
HL_mtd_target_df.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
HL_mtd_target_df.set_index('BRANCH',inplace=True)

HL_mtd_target_df=change(HL_mtd_target_df)
LAP_mtd_target_df=change(LAP_mtd_target_df)
NGPL_mtd_target_df=change(NGPL_mtd_target_df)
GPL_mtd_target_df=change(GPL_mtd_target_df)
mtd_target_df=change(mtd_target_df)
HL_mtd_target_df=HL_mtd_target_df.round(decimals=0)
LAP_mtd_target_df=LAP_mtd_target_df.round(decimals=0)
mtd_target_df=mtd_target_df.round(decimals=0)
GPL_mtd_target_df=GPL_mtd_target_df.round(decimals=0)
NGPL_mtd_target_df=NGPL_mtd_target_df.round(decimals=0)


'''
#--------------LAP_BUSINESS METRICS-------------------

#------------BOOK-----------

LP_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')

#------------ROI------------
LP_ROI<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
               & BM$BOOK_YEAR_MONTH == MTD)
#------------PF-------------
LP_PF<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked' 
                     & BM$BOOK_YEAR_MONTH == MTD)

#------------LTD (BOOKING TO LOGIN)----------
LTD_BOOK<-subset(BM,BM$FINTYPE %in% c('LP','NP') & BM$STATUS == 'Booked')
LTD_LOGIN<-subset (BM,BM$FINTYPE %in% c('LP','NP') & BM$LOGINSTATUS == 'A) Login')

LTD_BOOK1<-ddply(LTD_BOOK,c('FINBRANCH'),summarise,
              Book_value=sum(BOOKED_AMOUNT/10000000), 
              Book_vol=length(FINREFERENCE))
LTD_LOGIN1<-ddply(LTD_LOGIN,c('FINBRANCH'),summarise,
                  Log_value=sum(REQLOANAMT/1000000000),
                  Log_Vol=length(FINREFERENCE))

LTD_BOOK_LOG_mrg<-merge(LTD_LOGIN1,LTD_BOOK1,by.x = c('FINBRANCH'),
                        by.y = c('FINBRANCH'),all.x = T)
'''


todays_date = date.today()

def zonwise(newb,zone):
    newb=newb.transpose()
    newb.reset_index(inplace=True)
    newb.loc[11]=newb.sum()
    newb.loc[11,'FINBRANCH']=zone
    newb['BOOK_VALUE'] = newb.BOOKED_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]
    newb['BOOK_VOL'] = newb.one_count_MTD
    newb['DISB AMNT Tranch 1']=newb.DISBAMOUNT/1000000000
    newb['DISB_Act (%)']=(newb.DISBAMOUNT)/newb.BOOKED_AMOUNT
    newb['AUM']=newb.PRINCIPAL_OUTSTANDING/10000000
    newb['ROI MTD'] = newb.WROI / newb.BOOKED_AMOUNT
    newb['ROI M-1'] = newb.WROI_M1 /newb.BOOKED_AMOUNT_M1
    newb['GROSS%']=100*newb.NET_PREMIUM/newb.BOOKED_AMOUNT
    newb['PF%']=100*newb.PROCESSING_FEE / newb.BOOKED_AMOUNT
    newb['LTD_volume']=100*newb.one_count_tot_book/newb.one_count
    newb['LTD_Value']=100*(newb.TOTAL_BOOKED_AMOUNTytd*100)/newb.REQLOANAMT
    newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    newb.loc['Total']= newb.sum()
    newb.loc['Total','DISB_Act (%)']=np.float64(sum(newb.DISBAMOUNT))/np.float64(sum(newb.BOOKED_AMOUNT))
    newb.loc['Total','ROI MTD']=np.float64(sum(newb.WROI )) /np.float64(sum( newb.BOOKED_AMOUNT))
    newb.loc['Total','ROI M-1']=np.float64(sum(newb.WROI_M1))/np.float64(sum(newb.BOOKED_AMOUNT_M1))
    newb.loc['Total','GROSS%']=100*np.float64(sum(newb.NET_PREMIUM))/np.float64(sum(newb.BOOKED_AMOUNT))
    newb.loc['Total','PF%']=100*np.float64(sum(newb.PROCESSING_FEE )) / np.float64(sum(newb.BOOKED_AMOUNT))
    newb.loc['Total','LTD_volume']=100*np.float64(sum(newb.one_count_tot_book))/np.float64(sum(newb.one_count))
    newb.loc['Total','LTD_Value']=100*np.float64(sum(newb.TOTAL_BOOKED_AMOUNTytd))*(100)/np.float64(sum(newb.REQLOANAMT))
    newb.loc['Total','FINBRANCH']='Total'
    newb.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    newb["new"] = range(1,len(newb)+1)
    newb.loc['Total','new'] = 0
    newb=newb.sort_values(by="new").drop('new', axis=1)
    newb=newb.round(decimals=2)
    newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    newb.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    newb.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    newb.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    newb.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    newb.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    newb.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    newb.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    newb.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    newb['Booking Tar(Cr)']=' '
    newb['Booking Tar(#)']=' '
    newb['AUM Tar(Cr)']=' '
    newb['Net']=' '
    newb['Empaneled('+calendar.month_name[todays_date.month-1][:3]+')']=' '
    newb['Active']=' '
    newb['Channel']=' '
    newb['Employee']=' '
    newb['Cumulative']=' '
    newb.columns
    newb=newb.round(decimals=2)
    newb['Value%']=round(newb['Value%'])
    newb['Volume%']=round(newb['Volume%'])
    newb=newb[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Booking Tar(#)', 'Booking Act(#)','DISB_Act (%)', 'DISB Tr1','AUM Tar(Cr)','AUM Act(Cr)', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled('+calendar.month_name[todays_date.month-1][:3]+')', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #newb.set_index('FINBRANCH',inplace=True)
    finaln=newb.transpose()
    finaln.insert(0,'Target',' ')
    finaln.loc['FINBRANCH','Target']='Target'
    return finaln


def buss_metrics(LOAN,df,fd,todays_date):

    # MTD=str(todays_date.year)+"-"+str(todays_date.month-1)
    MTD=str(2022)+"-"+str(12)
    YTD=[]
    cur_month=todays_date.month
    while cur_month>=4:
        YTD.append(cur_month)
        cur_month=cur_month-1
    ytd=[]
    for i in YTD:
        ytd.append(str(todays_date.year)+"-"+str(i))
    
    ytd=['2022-12','2022-11','2022-10','2022-9','2022-8','2022-7','2022-6','2022-5','2022-4']
    BM_df = df[[ 'FINREFERENCE','NET_PREMIUM','BOOK_YEAR_MONTH','FINBRANCH','FINTYPE','STATUS', 'WROI','BOOKED_AMOUNT' ,'PROCESSING_FEE','REQLOANAMT','LOGINSTATUS','GPLFLAG_SANCTIONS','PRINCIPAL_OUTSTANDING','DISB_AMOUNT']]
    findisb=fd[['FINREFERENCE','DISBSEQ','DISBAMOUNT']]
    findisb=findisb[(findisb['DISBSEQ']==1)]
    findisb=findisb.sort_values(by='FINREFERENCE')
    BM_df=BM_df.sort_values(by='FINREFERENCE')
    fd_bm= pd.merge(BM_df, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')
    

    if LOAN== 'LAP':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'LP') | (fd_bm['FINTYPE'] =='NP'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'LP') |  (bm_df['FINTYPE'] =='NP'))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'HL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'HL') | (fd_bm['FINTYPE'] =='HT') | (fd_bm['FINTYPE'] =='FL') | (fd_bm['FINTYPE'] =='FT') | (fd_bm['FINTYPE'] =='LT')| (fd_bm['FINTYPE'] =='AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&((bm_df['FINTYPE'] == 'HL') | (bm_df['FINTYPE'] =='HT') | (bm_df['FINTYPE'] =='FL') | (bm_df['FINTYPE'] =='FT') | (bm_df['FINTYPE'] =='LT')| (bm_df['FINTYPE'] =='AHL'))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'GPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'] == 'GPL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'] == 'GPL')&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'NGPL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL'])))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['GPLFLAG_SANCTIONS'].isin(['NON GPL','AHL']))&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    elif LOAN== 'AHL':
        fd_bm= fd_bm[ (fd_bm.STATUS== 'Booked')& (fd_bm.LOGINSTATUS == 'A) Login')&((fd_bm['FINTYPE'] == 'AHL'))&(fd_bm.BOOK_YEAR_MONTH==MTD)]
        bm_df= BM_df
        bm_df= bm_df[ ((bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login'))&(bm_df['FINTYPE'] == 'AHL')&(bm_df.BOOK_YEAR_MONTH !='2023-01')]
    
    # bm_df_hl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &((bm_df['FINTYPE'] == 'HL') |  (bm_df['FINTYPE'] =='HT'))]
    # bm_df_gpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='GPL')]
    # bm_df_ngpl= bm_df[(bm_df.STATUS== 'Booked')| (bm_df.LOGINSTATUS == 'A) Login') &(bm_df['GPLFLAG_SANCTIONS'] =='NON GPL')]
    bm_df['one_count']=np.ones(len(bm_df), dtype = int)
    
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','FINREFERENCE','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    # def buss_metrics(bm_df):
        #bm_df_group = bm_df.groupby(['FINBRANCH'],as_index=False).sum(['WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING' ])
        #bm_df_group['BOOK_TARGET'] =
        
    fd_bm_group=fd_bm.groupby(['FINBRANCH'],as_index=False).sum(['WROI', 'BOOKED_AMOUNT', 'PROCESSING_FEE','REQLOANAMT','PRINCIPAL_OUTSTANDING', 'DISB_AMOUNT', 'DISBSEQ', 'DISBAMOUNT'])
    fd_bm_group=fd_bm_group[['FINBRANCH','DISBAMOUNT']]
    bm_df_book=bm_df[(bm_df.STATUS== 'Booked')&(bm_df.BOOK_YEAR_MONTH==MTD)]
    # bm_df_book=pd.merge(bm_df_book, findisb,left_on=['FINREFERENCE'],right_on=['FINREFERENCE'], how='inner')
    bm_df_book['one_count']=np.ones(len(bm_df_book), dtype = int)
    bm_df_BOOK_group=bm_df_book.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    
    
    b=bm_df[(bm_df.STATUS== 'Booked')]
    c=bm_df[((bm_df.STATUS== 'Booked') & (bm_df.BOOK_YEAR_MONTH.isin(ytd)))]
    c_group=c.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','WROI'])
    c_group=c_group[['FINBRANCH','BOOKED_AMOUNT','WROI']]
    c_group.rename(columns = {'BOOKED_AMOUNT':'BOOKED_AMOUNT_M1'}, inplace = True)
    c_group.rename(columns = {'WROI':'WROI_M1'}, inplace = True)
    b_group=b.groupby(['FINBRANCH'],as_index=False).sum(['BOOKED_AMOUNT','PRINCIPAL_OUTSTANDING','one_count' ])
    b_group=b_group[['FINBRANCH','PRINCIPAL_OUTSTANDING','one_count','BOOKED_AMOUNT' ]]
    b_group.rename(columns = {'one_count':'one_count_tot_book'}, inplace = True)
    b_group.rename(columns = {'BOOKED_AMOUNT':'TOTAL_BOOKED_AMOUNTytd'}, inplace = True)
    bm_df_group=bm_df.groupby(['FINBRANCH'],as_index=False).sum(['DISBSEQ','DISBAMOUNT','NET_PREMIUM','WROI','PROCESSING_FEE','BOOKED_AMOUNT','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISB_AMOUNT','one_count' ])
    bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','BOOKED_AMOUNT','WROI','NET_PREMIUM','one_count','PROCESSING_FEE']]
    bm_df_BOOK_group.rename(columns = {'one_count':'one_count_MTD'}, inplace = True)
    bm_df_group=bm_df_group[['FINBRANCH','REQLOANAMT','one_count']]
    
    merged=pd.merge(bm_df_group, bm_df_BOOK_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged2=pd.merge(merged, c_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged3=pd.merge(merged2, b_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    merged4=pd.merge(merged3, fd_bm_group,left_on=['FINBRANCH'],right_on=['FINBRANCH'], how='outer')
    '''
    CHANGES AFTER TRANSITION
    '''
    if LOAN=='LAP':
        newb=merged4.copy()
        newb.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
        newb.set_index('FINBRANCH',inplace=True)
        newb=newb.transpose()
        new=newb.copy()
        finaln=zonwise(new[['Hyderabad','Indore', 'Jaipur', 'Chandigarh', 'Chennai', 'Surat']],'New Branches')
        new=newb.copy()
        north=zonwise(new[['Delhi','Chandigarh','Jaipur']],'North')
        new=newb.copy()
        south=zonwise(new[['Bangalore','Chennai','Hyderabad']],'South')
        new=newb.copy()
        west1=zonwise(new[['Mumbai', 'Ahmedabad','Surat']],'West 1')
        new=newb.copy()
        west2=zonwise(new[['Indore','Pune']],'West 2')
        
    
    merged4['BOOK_VALUE'] = merged4.BOOKED_AMOUNT/10000000
    #bm_df_BOOK_group=bm_df_BOOK_group[['FINBRANCH','one_count','BOOKED_AMOUNT','WROI','NET_PREMIUM','PROCESSING_FEE']]

    merged4['BOOK_VOL'] = merged4.one_count_MTD
    merged4['DISB AMNT Tranch 1']=merged4.DISBAMOUNT/1000000000
    merged4['DISB_Act (%)']=(merged4.DISBAMOUNT)/merged4.BOOKED_AMOUNT
    merged4['AUM']=merged4.PRINCIPAL_OUTSTANDING/10000000
    merged4['ROI MTD'] = merged4.WROI / merged4.BOOKED_AMOUNT
    merged4['ROI M-1'] = merged4.WROI_M1 /merged4.BOOKED_AMOUNT_M1
    merged4['GROSS%']=100*merged4.NET_PREMIUM/merged4.BOOKED_AMOUNT
    merged4['PF%']=100*merged4.PROCESSING_FEE / merged4.BOOKED_AMOUNT
    merged4['LTD_volume']=100*merged4.one_count_tot_book/merged4.one_count
    merged4['LTD_Value']=100*(merged4.TOTAL_BOOKED_AMOUNTytd*100)/merged4.REQLOANAMT
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)

    merged4.loc['Total']= merged4.sum()
    merged4.loc['Total','DISB_Act (%)']=np.float64(sum(merged4.DISBAMOUNT))/np.float64(sum(merged4.BOOKED_AMOUNT))
    merged4.loc['Total','ROI MTD']=np.float64(sum(merged4.WROI )) /np.float64(sum( merged4.BOOKED_AMOUNT))
    merged4.loc['Total','ROI M-1']=np.float64(sum(merged4.WROI_M1))/np.float64(sum(merged4.BOOKED_AMOUNT_M1))
    merged4.loc['Total','GROSS%']=100*np.float64(sum(merged4.NET_PREMIUM))/np.float64(sum(merged4.BOOKED_AMOUNT))
    merged4.loc['Total','PF%']=100*np.float64(sum(merged4.PROCESSING_FEE )) / np.float64(sum(merged4.BOOKED_AMOUNT))
    merged4.loc['Total','LTD_volume']=100*np.float64(sum(merged4.one_count_tot_book))/np.float64(sum(merged4.one_count))
    merged4.loc['Total','LTD_Value']=100*np.float64(sum(merged4.TOTAL_BOOKED_AMOUNTytd))*(100)/np.float64(sum(merged4.REQLOANAMT))
    merged4.loc['Total','FINBRANCH']='Total'
    merged4.drop(columns = ['one_count_MTD','one_count_tot_book','TOTAL_BOOKED_AMOUNTytd','NET_PREMIUM','WROI','WROI_M1','PROCESSING_FEE','BOOKED_AMOUNT','BOOKED_AMOUNT_M1','REQLOANAMT','PRINCIPAL_OUTSTANDING','DISBAMOUNT','one_count' ],inplace = True)
    merged4["new"] = range(1,len(merged4)+1)
    merged4.loc['Total','new'] = 0
    merged4=merged4.sort_values(by="new").drop('new', axis=1)
    merged4=merged4.round(decimals=2)
    merged4.replace(to_replace=[np.nan,'',' '],value=0,inplace=True)
    merged4.rename(columns = {'BOOK_VALUE':'Booking Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'BOOK_VOL':'Booking Act(#)'}, inplace = True)
    merged4.rename(columns = {'AUM':'AUM Act(Cr)'}, inplace = True)
    merged4.rename(columns = {'ROI MTD':'ROI ('+calendar.month_name[todays_date.month-1][:3]+')'}, inplace = True)
    merged4.rename(columns = {'ROI M-1':'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')'}, inplace = True)
    merged4.rename(columns = {'LTD_volume':'Volume%'}, inplace = True)
    merged4.rename(columns = {'LTD_Value':'Value%'}, inplace = True)
    merged4.rename(columns = {'DISB AMNT Tranch 1':'DISB Tr1'}, inplace = True)
    merged4['Booking Tar(Cr)']=' '
    merged4['Booking Tar(#)']=' '
    merged4['AUM Tar(Cr)']=' '
    merged4['Net']=' '
    merged4['Empaneled('+calendar.month_name[todays_date.month-1][:3]+')']=' '
    merged4['Active']=' '
    merged4['Channel']=' '
    merged4['Employee']=' '
    merged4['Cumulative']=' '
    merged4.columns
    merged4=merged4.round(decimals=2)
    merged4['Value%']=round(merged4['Value%'])
    merged4['Volume%']=round(merged4['Volume%'])
    merged4=merged4[['FINBRANCH', 'Booking Tar(Cr)', 'Booking Act(Cr)','Booking Tar(#)', 'Booking Act(#)','DISB_Act (%)', 'DISB Tr1','AUM Tar(Cr)','AUM Act(Cr)', 'ROI ('+calendar.month_name[todays_date.month-1][:3]+')', 'ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')', 'GROSS%',  'Net', 'PF%', 'Empaneled('+calendar.month_name[todays_date.month-1][:3]+')', 'Active', 'Channel', 'Employee','Cumulative', 'Value%','Volume%']]
    #merged4.set_index('FINBRANCH',inplace=True)
    

    final=merged4.transpose()
    
    final.insert(0,'Target',' ')
    final.loc['FINBRANCH','Target']='Target'
    if LOAN=='LAP':
        final.reset_index(inplace=True)
        finaln.reset_index(inplace=True)
        finaln=finaln[['index',11]]
        final=pd.merge(final,finaln,left_on=('index'),right_on=('index'),how='outer')
        final.set_index('index',inplace=True)
        return final,north,west1,west2,south
    return(final)
# return(final)
lap_metrics,n,w1,w2,s=buss_metrics('LAP',df,fd,todays_date)
n.loc[:,'Total']=lap_metrics.loc[:,'Total']
w1.loc[:,'Total']=lap_metrics.loc[:,'Total']
w2.loc[:,'Total']=lap_metrics.loc[:,'Total']
s.loc[:,'Total']=lap_metrics.loc[:,'Total']
lap_metrics=lap_metrics[['Target', 'Total', 1, 4, 8, 9, 0, 2, 3, 5, 6, 7, 10,11]]
lap_metrics.columns=lap_metrics.loc['FINBRANCH']
w2.columns=w2.loc['FINBRANCH']
n.columns=n.loc['FINBRANCH']
s.columns=s.loc['FINBRANCH']
w1.columns=w1.loc['FINBRANCH']
def set_target(lap_metrics,LAP_mtd_target_df):
    
    for i in LAP_mtd_target_df.index:
        if  (i != 0) & (i in lap_metrics.columns.to_list() ): 
            lap_metrics.loc['Booking Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'Value (Cr)']
            lap_metrics.loc['Booking Act(Cr)',i]=str(round(lap_metrics.loc['Booking Act(Cr)',i]))+'('+str(round(100*lap_metrics.loc['Booking Act(Cr)',i]/LAP_mtd_target_df.loc[i,'Value (Cr)']))+'%)'
            lap_metrics.loc['Booking Tar(#)',i]=LAP_mtd_target_df.loc[i,'#']
            lap_metrics.loc['Booking Act(#)',i]=str(round(lap_metrics.loc['Booking Act(#)',i]))+'('+str(round(100*lap_metrics.loc['Booking Act(#)',i]/LAP_mtd_target_df.loc[i,'#']))+'%)'
            lap_metrics.loc['AUM Tar(Cr)',i]=LAP_mtd_target_df.loc[i,'AUM']
            lap_metrics.loc['AUM Act(Cr)',i]=str(round(lap_metrics.loc['AUM Act(Cr)',i]))+'('+str(round(100*lap_metrics.loc['AUM Act(Cr)',i]/LAP_mtd_target_df.loc[i,'AUM']))+'%)'
    return lap_metrics
hl_metrics=buss_metrics('HL',df,fd,todays_date)
hl_metrics.columns=hl_metrics.loc['FINBRANCH']
hl_metrics=hl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]


gpl_metrics=buss_metrics('GPL',df,fd,todays_date )
gpl_metrics.columns=gpl_metrics.loc['FINBRANCH']
gpl_metrics=gpl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]
ngpl_metrics=buss_metrics('NGPL',df,fd,todays_date )

ngpl_metrics.columns=ngpl_metrics.loc['FINBRANCH']
ngpl_metrics=ngpl_metrics[['Target', 'Total', 'Ahmedabad', 'Bangalore', 'Delhi', 'Mumbai', 'Pune']]

# ahl_metrics=buss_metrics('AHL',df,fd,todays_date )
# ahl_metrics.columns=ahl_metrics.loc['FINBRANCH']


hl_metrics=set_target(hl_metrics,HL_mtd_target_df)
lap_metrics=set_target(lap_metrics,LAP_mtd_target_df)
n=set_target(n,LAP_mtd_target_df)
s=set_target(s,LAP_mtd_target_df)
w1=set_target(w1,LAP_mtd_target_df)
w2=set_target(w2,LAP_mtd_target_df)
gpl_metrics=set_target(gpl_metrics,GPL_mtd_target_df)
ngpl_metrics=set_target(ngpl_metrics,NGPL_mtd_target_df)


lap_metrics.loc['DISB_Act (%)','Target']='90%'
lap_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
lap_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
lap_metrics.loc['GROSS%','Target']='1.5%'
lap_metrics.loc['PF%','Target']='0.9%'
lap_metrics.loc['Cumulative','Target']='1.35%'
lap_metrics.loc['Value%','Target']='33%'
lap_metrics.loc['Volume%','Target']='33%'

n.loc['DISB_Act (%)','Target']='90%'
n.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
n.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
n.loc['GROSS%','Target']='1.5%'
n.loc['PF%','Target']='0.9%'
n.loc['Cumulative','Target']='1.35%'
n.loc['Value%','Target']='33%'
n.loc['Volume%','Target']='33%'

s.loc['DISB_Act (%)','Target']='90%'
s.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
s.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
s.loc['GROSS%','Target']='1.5%'
s.loc['PF%','Target']='0.9%'
s.loc['Cumulative','Target']='1.35%'
s.loc['Value%','Target']='33%'
s.loc['Volume%','Target']='33%'

w1.loc['DISB_Act (%)','Target']='90%'
w1.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
w1.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
w1.loc['GROSS%','Target']='1.5%'
w1.loc['PF%','Target']='0.9%'
w1.loc['Cumulative','Target']='1.35%'
w1.loc['Value%','Target']='33%'
w1.loc['Volume%','Target']='33%'

w2.loc['DISB_Act (%)','Target']='90%'
w2.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='8.75%'
w2.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='8.75%'
w2.loc['GROSS%','Target']='1.5%'
w2.loc['PF%','Target']='0.9%'
w2.loc['Cumulative','Target']='1.35%'
w2.loc['Value%','Target']='33%'
w2.loc['Volume%','Target']='33%'

# hl_metrics.loc['DISB_Act (%)','Target']='90%'
hl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
hl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
hl_metrics.loc['GROSS%','Target']='2%'
hl_metrics.loc['PF%','Target']='0.02%'
hl_metrics.loc['Cumulative','Target']='0.45%'
hl_metrics.loc['Value%','Target']='50%'
hl_metrics.loc['Volume%','Target']='50%'

gpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
gpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
gpl_metrics.loc['GROSS%','Target']='2%'
gpl_metrics.loc['PF%','Target']='0.02%'
gpl_metrics.loc['Cumulative','Target']='0.45%'
gpl_metrics.loc['Value%','Target']='50%'
gpl_metrics.loc['Volume%','Target']='50%'


ngpl_metrics.loc['ROI ('+calendar.month_name[todays_date.month-1][:3]+')','Target']='7.75%'
ngpl_metrics.loc['ROI (Till'+calendar.month_name[todays_date.month-2][:3]+')','Target']='7.75%'
ngpl_metrics.loc['GROSS%','Target']='2%'
ngpl_metrics.loc['PF%','Target']='0.02%'
ngpl_metrics.loc['Cumulative','Target']='0.45%'
ngpl_metrics.loc['Value%','Target']='50%'
ngpl_metrics.loc['Volume%','Target']='50%'


MTD=str(2022)+"-"+str(12)
with pd.ExcelWriter(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\1"+MTD+".xlsx") as writer:
    lap_metrics.to_excel(writer, sheet_name="LAP BM")
    s.to_excel(writer, sheet_name="SOUTH BM")
    n.to_excel(writer, sheet_name="NORTH BM")
    w1.to_excel(writer, sheet_name="WEST1 BM")
    w2.to_excel(writer, sheet_name="WEST2 BM")
    hl_metrics.to_excel(writer, sheet_name="HL BM")
    gpl_metrics.to_excel(writer, sheet_name="GPL BM")
    ngpl_metrics.to_excel(writer, sheet_name="NGPL BM")
    # ahl_metrics.to_excel(writer, sheet_name="AHL BM")

from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# ---create presentation with 1 slide---
prs = Presentation(r"C:\Users\VAIBHAV.SRIVASTAV01\Documents\OPSREVIEWDECK.pptx")
# prs.slide_width = 11887200
# prs.slide_height = 6686550
def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
def create_slide(prs,lap_metrics,title,c,target):
    prs.slides[c].shapes[0].text=title
    lap_metrics.reset_index(level=0, inplace=True)
    lap_metrics.insert(0,' ',' ')
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    # slide.shapes.title.text = title
    # slide.shapes.title.top = Inches(0.5)
    # slide.shapes.title.left = Inches(0.5)
    # slide.shapes.title.width = Inches(10)

    
    # # ---add table to slide---
    # x, y, cx, cy = Inches(0.5), Inches(1), Inches(11), Inches(3)
    # shape = slide.shapes.add_table(len(lap_metrics),len(lap_metrics.columns),x, y, cx,cy)
    
    # tbl =  shape._element.graphic.graphicData.tbl
    # style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    # tbl[0][-1].text = style_id    
    # table = shape.table
    
    cols=prs.slides[c].shapes[2].table.columns
    cols[1].width=Inches(1.4)
    prs.slides[c].shapes[2].table.width=Inches(12)
    prs.slides[c].shapes[2].table.height=Inches(5)
    # cols[0].width=Inches(1.4)
    # rows=table.rows
    for i in range(len(lap_metrics)):
        for j in range(len(lap_metrics.columns)):
            if ((i in (5,11,13,19,20))& (j>=2)):
                prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])+'%'
            else:
                prs.slides[c].shapes[2].table.cell(i,j).text=str(lap_metrics.iloc[i,j])
                    
                
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_bottom = Inches(0.08)
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.margin_left = 0
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.vertical_anchor = MSO_ANCHOR.TOP
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.word_wrap = False
                prs.slides[c].shapes[2].table.cell(i,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                #prs.slides[c].shapes[2].table.cell(i,j).text_frame.fit_text()
                # if i==0:
                #     cell=table.cell(i, j)
                #     fill = cell.fill
                #     fill.solid()
                #     fill.fore_color.rgb = RGBColor(173, 216, 230)
    
                
    
    
    prs.slides[c].shapes[2].table.cell(1,0).text="Business Number"
    prs.slides[c].shapes[2].table.cell(9,0).text="ROI"
    prs.slides[c].shapes[2].table.cell(11,0).text="%Insurance"
    prs.slides[c].shapes[2].table.cell(13,0).text="PF%"
    prs.slides[c].shapes[2].table.cell(14,0).text="Number of DSAs"
    prs.slides[c].shapes[2].table.cell(16,0).text="Incentive"
    prs.slides[c].shapes[2].table.cell(19,0).text="LTD Booking to Login"
    
    for cell in iter_cells(prs.slides[c].shapes[2].table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if title=="LAP Business Metrics":
                    run.font.size = Pt(11)
                else:
                    run.font.size = Pt(14)
                run.font.style='Calibri'
                run.font.color.rgb = RGBColor(0,0,0)
                run.font.bold = False
    
    return(prs)

MTD=str(2022)+"-"+str(12)
create_slide(prs,lap_metrics,"LAP Business Metrics",13,target)            
create_slide(prs,hl_metrics,"HL Business Metrics",14,target)       
create_slide(prs,gpl_metrics,"GPL Business Metrics",15,target)
create_slide(prs,ngpl_metrics,"N-GPL Business Metrics",16,target)
create_slide(prs,n,"LAP Business Metrics North",17,target)   
create_slide(prs,s,"LAP Business Metrics South",18,target)   
create_slide(prs,w1,"LAP Business Metrics West 1",19,target)   
create_slide(prs,w2,"LAP Business Metrics West 2",20,target)   
# create_slide(prs,ahl_metrics,"AHL Business Metrics",17,target)
prs.save(r"C:\Users\VAIBHAV.SRIVASTAV01\Desktop\NOV_OPS_REVIEW\OPERATING REVIEW"+MTD +"NEW.pptx")



